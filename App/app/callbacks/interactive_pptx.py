# =============================================================================
# MSI Analysis Application - Interactive PPTX Export
# インタラクティブ解析 PPTX エクスポート コールバック
#
# interactive_callbacks.py から分離された PPTX 関連の
# ヘルパー関数・コールバックをまとめたモジュール。
# =============================================================================

import json
import logging
import re
from datetime import datetime
from io import BytesIO
from pathlib import Path

import numpy as np
import pandas as pd
import plotly.graph_objects as go
from dash import (Input, Output, State, callback, dcc, no_update)
from dash.exceptions import PreventUpdate

from app.services.seurat_bridge import SeuratBridge
from app.services.caveats import banner_text as _caveat_banner
from app.utils.color_utils import (
    cluster_sort_key as _cluster_sort_key,
    get_cluster_color_map as _get_cluster_color_map,
    cluster_display_name as _cluster_display_name,
)
from app.utils.display_helpers import (
    display_name as _display_name,
)
from app.utils.deg_utils import (
    is_meaningful_annotation as _is_meaningful_annotation,
    get_top_n_features_for_cluster as _get_top_n_features_for_cluster,
)
from app.utils.label_persistence import (
    compute_annotation_offsets as _compute_annotation_offsets,
    load_label_positions as _load_label_positions_util,
    load_cluster_name_map as _load_cluster_name_map,
)
from app.utils.pptx_helpers import (
    render_png as _fig_to_png_bytes,
    shutdown_shared_queue as _shutdown_render_pool,
    pptx_add_title_bar as _pptx_add_title_bar,
    pptx_add_image as _pptx_add_image,
    pptx_add_image_preserve_ratio as _pptx_add_image_preserve_ratio,
    square_tile_dims as _square_tile_dims,
    build_cluster_legend_fig as _build_cluster_legend_fig,
    pptx_add_sections as _pptx_add_sections,
)
from app.utils import raster as _raster

# 共有状態・ヘルパーを分離モジュールから参照
from app.callbacks.interactive_callbacks import (
    _interactive_data,
    _bridge,
    _load_deg_results,
    _load_label_positions,
)
from app.callbacks.interactive_umap import (
    _build_umap_integrated_fig,
    _get_merged_label_positions,
)
from app.callbacks.interactive_spatial import (
    _create_single_spatial_fig,
    _transform_coords,
    _calc_zero_gap_marker_size,
)
from app.callbacks.interactive_calibration import (
    _build_mz_to_compound_map,
    _annotate_gene_labels,
)
from app.utils.validation import coerce_count
from app.utils.display_settings import read_display_settings

logger = logging.getLogger("msi.interactive.pptx")


# ---------------------------------------------------------------------------
# PPTX (Google Slides) エクスポート
# ---------------------------------------------------------------------------

def _build_feature_plot_fig(df, feature_name, cache_dir_path, rds_path,
                            rotation_store=None, name_map=None, marker_size=2,
                            colorbar_tickformat=None, show_colorbar_title=True,
                            auto_marker_size=False, show_colorbar=True,
                            intensity_min=None, intensity_max=None):
    """単一 m/z Feature の Spatial Expression Plot figure を生成（PPTX 用）。

    intensity_min / intensity_max: 強度レンジ (%) (ver51.9)。画面側の
        `feature_intensity_min/max` と揃えるための引数。色域が変わると
        同じ図に見えて別のことを示すため、資料と画面で食い違うと危ない。
        None なら従来どおりデータ全域。

    Returns:
        go.Figure or None
    """
    from plotly.subplots import make_subplots

    if rotation_store is None:
        rotation_store = {}
    if name_map is None:
        name_map = {}

    # 発現データ取得
    expression = None
    if cache_dir_path:
        cache_dir = Path(cache_dir_path) if isinstance(cache_dir_path, str) else cache_dir_path
        expression = _bridge.get_feature_expression_fast(cache_dir, feature_name)
    if expression is None and rds_path:
        try:
            expression = _bridge.get_feature_expression(rds_path, feature_name)
        except Exception:
            return None
    if expression is None:
        return None

    # ★ ver51.9: 長さの番人。画面側 (interactive_deg.py:524) には ver51.6 で
    #   入れたのに PPTX 側だけ無かった。R フォールバックがヘッダ 1 行ぶん長い
    #   Series を返すと、代入で pandas が例外を投げる。ここは
    #   **数十分かかる背景エクスポートの終盤**なので、最後まで走ってから
    #   全部無駄になる。1 枚を諦めて資料は出す方が損失が小さい。
    _expr_arr = np.asarray(
        expression.values if hasattr(expression, "values") else expression)
    if len(_expr_arr) != len(df):
        logger.warning(
            "発現量の長さが不一致のため Feature をスキップ (feature=%s, %d != %d)",
            feature_name, len(_expr_arr), len(df))
        return None

    df_plot = df.copy()
    df_plot["_expression"] = _expr_arr

    if "SpatialX" not in df_plot.columns:
        return None

    samples = sorted(df_plot["Sample"].unique())
    n_samples = len(samples)
    subplot_titles = [_display_name(s, name_map) for s in samples]

    fig = make_subplots(
        rows=1, cols=n_samples,
        subplot_titles=subplot_titles,
        horizontal_spacing=0.03,
    )

    expr_vals = df_plot["_expression"].values
    # ★ ver51.9: 強度レンジ (%) を画面と揃える。式は画面側の
    #   `interactive_deg._feature_display_range` と同一 (別々に書くとまたずれる)。
    from app.callbacks.interactive_deg import _feature_display_range
    _rng = _feature_display_range(expr_vals, intensity_min, intensity_max)
    if _rng is None:
        global_min = float(np.nanmin(expr_vals))
        global_max = float(np.nanmax(expr_vals))
    else:
        global_min, global_max = _rng

    use_raster = _raster.raster_enabled()
    for idx, s in enumerate(samples, 1):
        df_s = df_plot[df_plot["Sample"] == s]
        transform = rotation_store.get(
            s, rotation_store.get("__all__", {"angle": 0, "flip_h": False, "flip_v": False}))
        if isinstance(transform, (int, float)):
            transform = {"angle": int(transform), "flip_h": False, "flip_v": False}

        raw_x = df_s["SpatialX"].values
        raw_y = -df_s["SpatialY"].values
        plot_x, plot_y = _transform_coords(
            raw_x, raw_y,
            transform.get("angle", 0),
            flip_h=transform.get("flip_h", False),
            flip_v=transform.get("flip_v", False),
        )

        # === ラスター経路（規則グリッドを go.Heatmap 化して kaleido を高速化）===
        # 座標つき Heatmap のため既存の軸設定（scaleanchor 等）をそのまま活かせる。
        if use_raster:
            _res = _raster.bin_to_grid(
                plot_x, plot_y, df_s["_expression"].values, agg="mean")
            if _res is not None:
                _z, _xc, _yc = _res
                _show_scale = (idx == n_samples) and show_colorbar
                _cb = None
                if _show_scale:
                    _cb = dict(len=0.8, thickness=15,
                               tickvals=[global_min, global_max],
                               ticktext=["0%", "100%"])
                    if show_colorbar_title:
                        _cb["title"] = "Intensity"
                fig.add_trace(_raster.heatmap_trace(
                    _z, _xc, _yc, "Plasma", global_min, global_max,
                    showscale=_show_scale, colorbar=_cb,
                ), row=1, col=idx)
                continue  # 散布経路をスキップ

        # サンプル別マーカーサイズ自動計算
        ms = marker_size
        if auto_marker_size and len(plot_x) > 1:
            sorted_ux = np.sort(np.unique(plot_x))
            if len(sorted_ux) > 1:
                min_sp = float(np.min(np.diff(sorted_ux)))
                xr = float(plot_x.max() - plot_x.min())
                yr = float(plot_y.max() - plot_y.min()) if len(plot_y) > 1 else 1.0
                # Feature Plot は height=280 描画、margin 約40px → 有効240px
                eff_h = 240
                if yr > 0 and xr > 0:
                    eff_w = eff_h * (xr / yr)
                    ms = max(2, round(min_sp * eff_w / xr * 1.5))

        # TIC 背景
        if "TotalCount" in df_s.columns:
            fig.add_trace(go.Scatter(
                x=plot_x, y=plot_y, mode="markers",
                marker=dict(size=ms, symbol="square", color=df_s["TotalCount"].values,
                            colorscale="Greys", opacity=0.5, showscale=False),
                hoverinfo="skip", showlegend=False,
            ), row=1, col=idx)

        # 発現量オーバーレイ
        is_last = (idx == n_samples)
        _show_scale = is_last and show_colorbar
        marker_opts = dict(
            size=ms,
            symbol="square",
            color=df_s["_expression"].values,
            colorscale="Plasma",
            cmin=global_min, cmax=global_max,
            showscale=_show_scale,
            opacity=0.8,
        )
        if _show_scale:
            cb_opts = dict(
                len=0.8, thickness=15,
                tickvals=[global_min, global_max],
                ticktext=["0%", "100%"],
            )
            if show_colorbar_title:
                cb_opts["title"] = "Intensity"
            marker_opts["colorbar"] = cb_opts
        fig.add_trace(go.Scatter(
            x=plot_x, y=plot_y, mode="markers",
            marker=marker_opts, showlegend=False,
        ), row=1, col=idx)

    # FUTURE(annot-provenance): 将来「由来表示」を足す場合、feature タイトルに
    #   app/services/annotation_sources.format_annotation_label() で由来（source）を併記する想定。
    #   取込設計が未確定のため現状は変更なし。詳細: App/docs/MVP4_IMPLEMENTATION_STATUS.md
    fig.update_layout(
        title=dict(text=feature_name, font=dict(size=20), x=0.5),
        plot_bgcolor="white", paper_bgcolor="white",
        margin=dict(l=10, r=50 if show_colorbar else 10, t=50, b=10),
    )
    # Subplot titles のフォントサイズ拡大
    fig.update_annotations(font_size=16)
    for i in range(1, n_samples + 1):
        fig.update_xaxes(visible=False, row=1, col=i)
        xanchor = f"x{i}" if i > 1 else "x"
        fig.update_yaxes(visible=False, scaleanchor=xanchor, row=1, col=i)

    return fig


# --------------------------------------------------------------------------
# 表示設定（ver51.9 / B-2）
# --------------------------------------------------------------------------
# ★ 「解析条件」スライドと図が食い違っていた。条件スライドは
#   `provenance_callbacks` が記録した利用者の閾値を出すのに、Volcano は
#   `_build_volcano_fig_for_cluster(deg_data, cl_str)` と **既定値のまま**
#   呼ばれていたので破線も有意判定も常に 0.5 / 1.3 だった。
#   同じ 1 ファイルの中で矛盾するため、外から真偽を判断できない。
#
#   直し方は「条件スライドと図が **同じ 1 つの dict** を見る」。`conditions` は
#   クリック時点で確定させているので、そこから取り出せば構造上ずれない。

# 記録が無い/壊れているときの既定値。**従来の図と同じになる値**にする
# （条件を一度も触っていないプロジェクトで見た目を変えないため）。
# ★ ver52.3 ⑥: 表示設定の取り出しは `app/utils/display_settings.py` へ移した。
#   同じ `interactive_settings.json` を出典にしている 3 系統
#   （画面 / PPTX / Lite ビュー）が、それぞれ別の読み出しコードを持っていたのが
#   「PPTX だけ直して Lite は 0.5/1.3 のまま」（ver51.9 B-2）の原因だった。
#   ここは中立モジュールへの薄い委譲だけを残す。
# ★ 旧 `_DISPLAY_DEFAULTS` / `_num` は消した。呼び出しが 0 件の別名を残すのは、
#   この版で潰している「宣言はあるが誰も使わない」(T8) を自分で作ること。
def _display_settings(conditions) -> dict:
    """「解析条件」スライドが出すのと同じ表示設定を、図の生成用に取り出す。

    PPTX は `provenance.collect_conditions()` を通すので 1 段深い。
    その 1 段を剥がすのがこの関数の役割で、中身の解釈は共通側に置く。
    """
    return read_display_settings((conditions or {}).get("interactive"))


def _apply_merge_view(df, display_settings, custom_colors):
    """マージ表示なら Cluster / UMAP 座標を統合後のものへ差し替える (ver51.9)。

    ★ 画面 (`interactive_umap.py:537-545`) と同じ変換。PPTX にだけ無かったため、
      「マージ統合」で確認した利用者が **元のクラスタリングの資料** を配っていた。
      見出しも凡例も同じ形なので、受け取った人には区別が付かない。

    Returns (df, effective_custom_colors, merged: bool)
    """
    if (display_settings or {}).get("merge_toggle") != "merged":
        return df, custom_colors, False
    if df is None or "Cluster_merged" not in df.columns:
        return df, custom_colors, False

    cols = [c for c in ("Sample", "CellID", "SpatialX", "SpatialY",
                        "Annotation") if c in df.columns]
    merged_df = df[cols].copy()
    merged_df["Cluster"] = df["Cluster_merged"].to_numpy()
    for src, dst in (("UMAP_1_merged", "UMAP_1"), ("UMAP_2_merged", "UMAP_2")):
        if src in df.columns:
            merged_df[dst] = df[src].to_numpy()
        elif dst in df.columns:
            merged_df[dst] = df[dst].to_numpy()

    from app.utils.color_utils import get_merged_cluster_color_map
    effective = get_merged_cluster_color_map(
        merged_df["Cluster"],
        mode=display_settings.get("merge_color_mode") or "shade")
    return merged_df, effective, True


def _build_volcano_fig_for_cluster(deg_data, cluster, fc_thresh=0.5, p_thresh=1.3,
                                   marker_size=8, label_top_n=5):
    """指定クラスタの Volcano Plot figure を生成（PPTX 用）。

    label_top_n: 自動ラベルを付ける up/down それぞれの件数 (ver51.9)。
        画面側 (`interactive_deg.update_volcano_plot`) は
        `volcano_label_top_n` を見るのに、ここは 5 固定だった。

    Returns:
        go.Figure or None
    """
    if not deg_data:
        return None
    df = pd.DataFrame(deg_data)
    if "p_val_adj_raw" in df.columns:
        df["p_num"] = pd.to_numeric(df["p_val_adj_raw"], errors="coerce")
    else:
        df["p_num"] = pd.to_numeric(df["p_val_adj"], errors="coerce")
    df["avg_log2FC"] = pd.to_numeric(df["avg_log2FC"], errors="coerce")
    min_nonzero_p = (
        df.loc[df["p_num"] > 0, "p_num"].min()
        if (df["p_num"] > 0).any()
        else 5e-324
    )
    df["neg_log10_p"] = -np.log10(df["p_num"].clip(lower=min_nonzero_p))

    if "annotation" in df.columns:
        df["display_text"] = df.apply(
            lambda r: f"{r['gene']}\n({r['annotation']})"
            if _is_meaningful_annotation(r.get("annotation", ""), r.get("gene", ""))
            else r["gene"],
            axis=1,
        )
    else:
        df["display_text"] = df["gene"]

    df = df[df["cluster"].astype(str) == str(cluster)]
    if df.empty:
        return None

    fig = go.Figure()
    for reg, color, label in [
        ("Up", "#FF2D2D", "Up-regulated"),
        ("Down", "#1E5BFF", "Down-regulated"),
        ("NS", "#7A7A7A", "Not significant"),
    ]:
        if reg == "Up":
            mask = (df["neg_log10_p"] >= p_thresh) & (df["avg_log2FC"] >= fc_thresh)
        elif reg == "Down":
            mask = (df["neg_log10_p"] >= p_thresh) & (df["avg_log2FC"] <= -fc_thresh)
        else:
            mask = ~(
                (df["neg_log10_p"] >= p_thresh) & (df["avg_log2FC"].abs() >= fc_thresh)
            )
        sub = df[mask]
        if len(sub) > 0:
            fig.add_trace(go.Scattergl(
                x=sub["avg_log2FC"], y=sub["neg_log10_p"],
                mode="markers",
                marker=dict(size=marker_size, color=color, opacity=0.7),
                name=label, text=sub["display_text"],
                hovertemplate=(
                    "<b>%{text}</b><br>"
                    "log2FC: %{x:.3f}<br>"
                    "-log10(p): %{y:.2f}<extra></extra>"
                ),
            ))

    fig.add_hline(y=p_thresh, line_dash="dash", line_color="gray", opacity=0.5)
    fig.add_vline(x=fc_thresh, line_dash="dash", line_color="gray", opacity=0.5)
    fig.add_vline(x=-fc_thresh, line_dash="dash", line_color="gray", opacity=0.5)

    # --- Top N アノテーション（PPTX用） ---
    sig_mask = (df["neg_log10_p"] >= p_thresh) & (df["avg_log2FC"].abs() >= fc_thresh)
    sig_df = df[sig_mask]
    # ★ ver52.3 ⑤: 画面側と**同じ関数**で解決する（式を 2 箇所に書かない）。
    _n_label = coerce_count(label_top_n, "volcano_label_top_n")
    if not sig_df.empty and _n_label > 0:
        _top_up = sig_df[sig_df["avg_log2FC"] > 0].nlargest(_n_label, "avg_log2FC")
        _top_down = sig_df[sig_df["avg_log2FC"] < 0].nsmallest(_n_label, "avg_log2FC")
        _auto_label = pd.concat([_top_up, _top_down])
        if len(_auto_label) > 0:
            _pts = list(zip(
                _auto_label["avg_log2FC"].values,
                _auto_label["neg_log10_p"].values,
                _auto_label["display_text"].values,
            ))
            _offsets = _compute_annotation_offsets(_pts)
            for (_x, _y, _txt), (_ax, _ay) in zip(_pts, _offsets):
                fig.add_annotation(
                    x=_x, y=_y, text=_txt,
                    showarrow=True, arrowhead=0, arrowwidth=1,
                    arrowcolor="#999", ax=_ax, ay=_ay,
                    font=dict(size=9, color="#333"),
                    bgcolor="rgba(255,255,255,0.8)", borderpad=2,
                )

    fig.update_layout(
        title=dict(text=f"Volcano Plot - Cluster {cluster}", font=dict(size=20), x=0.5),
        xaxis_title="avg_log2FC",
        yaxis_title="-log10(p_val_adj)",
        template="plotly_white",
        margin=dict(l=50, r=20, t=40, b=58),
    )
    # pixel-level 探索的ランキングである旨を図にも明記（PPTX/共有図に注意書きを残す）
    fig.add_annotation(
        text=_caveat_banner("en", short=True), showarrow=False,
        xref="paper", yref="paper", x=0.5, y=-0.16, xanchor="center",
        font=dict(size=10, color="#a15c00"),
    )
    return fig


def _build_heatmap_for_cluster(deg_data, cluster, df, cache_dir, top_n,
                                mrm_path=None, scale=None):
    """指定クラスタの Top N マーカーを全クラスタで比較するヒートマップを生成。

    scale: "zscore" | "raw" | None (ver51.9)。画面のスケール切替
        (`heatmap_scale`) と揃えるための引数。**Z-score か Raw かは
        「何をプロットしたか」そのもの**で、見た目だけの違いではない。
        従来は常に Z-score だったため、Raw を選んで確認した利用者が
        Z-score の図を配ることになっていた。None は従来どおり Z-score。
    """
    if not deg_data or not cache_dir:
        return None
    try:
        df_deg = pd.DataFrame(deg_data)
        df_deg["p_num"] = pd.to_numeric(df_deg["p_val_adj"], errors="coerce")
        cl_deg = df_deg[df_deg["cluster"].astype(str) == str(cluster)]
        if cl_deg.empty:
            return None
        top_markers = cl_deg.sort_values("p_num").head(top_n)
        genes = top_markers["gene"].unique().tolist()
        if not genes:
            return None

        cache_dir_p = Path(cache_dir) if isinstance(cache_dir, str) else cache_dir
        expr_path = cache_dir_p / "expression_matrix.parquet"
        if not expr_path.exists():
            return None

        # 利用可能な遺伝子を Parquet schema から一括判定（I/O 1 回で完結）
        import pyarrow.parquet as pq
        try:
            schema_names = set(pq.read_schema(expr_path).names)
        except Exception:
            schema_names = set()
        available = [g for g in genes if g in schema_names]
        if not available:
            return None

        expr_df = pd.read_parquet(expr_path, columns=["CellID"] + available)
        plot_data = df if df is not None else _interactive_data.get("plot_data")
        if plot_data is None:
            return None
        merged = expr_df.merge(
            plot_data[["CellID", "Cluster"]], on="CellID", how="inner")
        cluster_means = merged.groupby("Cluster")[available].mean()
        cluster_means = cluster_means.reindex(
            sorted(cluster_means.index, key=_cluster_sort_key))

        z_data = cluster_means.values.copy()
        is_zscore = (scale != "raw")     # None / "zscore" は従来どおり Z-score
        if is_zscore:
            col_mean = z_data.mean(axis=0)
            col_std = z_data.std(axis=0)
            col_std[col_std == 0] = 1
            z_data = (z_data - col_mean) / col_std

        # Y軸ラベル（アノテーション付き）
        y_labels = available
        gene_to_annotation = {}
        for r in deg_data:
            gene = r.get("gene", "")
            ann = r.get("annotation", "")
            if gene and _is_meaningful_annotation(ann, gene):
                gene_to_annotation[gene] = ann
        if gene_to_annotation:
            y_labels = [
                f"{g} ({gene_to_annotation[g]})" if g in gene_to_annotation else g
                for g in available
            ]
        elif mrm_path:
            mz_to_compound = _build_mz_to_compound_map(str(mrm_path), tolerance=0.1)
            y_labels = _annotate_gene_labels(available, mz_to_compound, tolerance=0.1)

        _z_label = "Z-score" if is_zscore else "Intensity"
        fig = go.Figure(go.Heatmap(
            z=z_data.T,
            x=[str(c) for c in cluster_means.index],
            y=y_labels,
            colorscale="RdBu_r",
            zmid=0 if is_zscore else None,
            ygap=2,  # Top N m/z行間の区切り線
            hovertemplate=("Cluster: %{x}<br>Gene: %{y}<br>"
                           + _z_label + ": %{z:.3f}<extra></extra>"),
        ))
        fig.update_layout(
            margin=dict(l=20, r=20, t=30, b=20),
            xaxis=dict(side="bottom"),
            yaxis=dict(autorange="reversed"),
            paper_bgcolor="white",
            plot_bgcolor="white",
        )
        return fig.to_dict()
    except Exception as e:
        logger.error(f"Per-cluster heatmap error (cluster {cluster}): {e}")
        return None


def _build_heatmap_for_pptx(heatmap_fig, deg_data, df, cache_dir, top_n,
                             mrm_path=None):
    """PPTX 用にZ-score + アノテーション付きヒートマップ figure を生成/調整する。

    アプリの heatmap_fig がすでに Z-score ならアノテーションだけ補完し、
    Raw なら expression_matrix から再計算して Z-score + アノテーション付きで返す。
    """
    # ---- deg_data / cache_dir がなければアプリ図をそのまま返す ----
    if not deg_data or not cache_dir:
        return heatmap_fig

    cache_dir_path = Path(cache_dir) if isinstance(cache_dir, str) else cache_dir
    expr_path = cache_dir_path / "expression_matrix.parquet" if cache_dir_path else None

    # ---- アノテーション辞書を構築 ----
    gene_to_annotation = {}
    for r in deg_data:
        gene = r.get("gene", "")
        ann = r.get("annotation", "")
        if gene and _is_meaningful_annotation(ann, gene):
            gene_to_annotation[gene] = ann

    # MRM fallback
    if not gene_to_annotation and mrm_path:
        try:
            mz_to_compound = _build_mz_to_compound_map(mrm_path, tolerance=0.1)
        except Exception:
            mz_to_compound = {}
    else:
        mz_to_compound = {}

    # ---- heatmap_fig が Z-score かどうかを判定 ----
    is_zscore = False
    if heatmap_fig:
        traces = heatmap_fig.get("data", [])
        for t in traces:
            if t.get("type") == "heatmap" and t.get("zmid") == 0:
                is_zscore = True
                break

    # ---- Z-score 済みの場合: アノテーション補完のみ ----
    if is_zscore and heatmap_fig:
        fig = go.Figure(heatmap_fig)
        # Y軸ラベルにアノテーションを補完
        for t in fig.data:
            if hasattr(t, "y") and t.y is not None:
                current_labels = list(t.y)
                needs_update = False
                new_labels = []
                for lbl in current_labels:
                    gene = str(lbl).split(" (")[0]  # 既存アノテーション除去
                    if gene in gene_to_annotation:
                        new_labels.append(f"{gene} ({gene_to_annotation[gene]})")
                        needs_update = True
                    elif mz_to_compound:
                        annotated = _annotate_gene_labels(
                            [gene], mz_to_compound, tolerance=0.1)
                        new_labels.append(annotated[0])
                        if annotated[0] != gene:
                            needs_update = True
                    else:
                        new_labels.append(lbl)
                if needs_update:
                    t.y = new_labels
        # 左マージンを再調整
        all_labels = []
        for t in fig.data:
            if hasattr(t, "y") and t.y is not None:
                all_labels.extend(str(l) for l in t.y)
        if all_labels:
            max_len = max(len(l) for l in all_labels)
            left_margin = min(max(max_len * 7, 120), 350)
            fig.update_layout(margin=dict(l=left_margin))
        return fig.to_dict()

    # ---- Raw / heatmap_fig がない場合: expression_matrix から新規生成 ----
    if not expr_path or not expr_path.exists():
        return heatmap_fig  # fallback: アプリ図をそのまま
    if df is None or df.empty:
        return heatmap_fig

    try:
        df_deg = pd.DataFrame(deg_data)
        df_deg["p_num"] = pd.to_numeric(df_deg.get("p_val_adj", ""), errors="coerce")
        top_markers = df_deg.sort_values("p_num").groupby("cluster").head(top_n)
        genes = top_markers["gene"].unique().tolist()
        if not genes:
            return heatmap_fig

        # 利用可能な遺伝子を Parquet schema から一括判定（I/O 1 回で完結）。
        # 旧実装は存在確認のため feature 数ぶん parquet を開いて捨てていた（無駄読み）。
        import pyarrow.parquet as pq
        try:
            schema_names = set(pq.read_schema(expr_path).names)
        except Exception:
            schema_names = set()
        available = [g for g in genes if g in schema_names]
        if not available:
            return heatmap_fig

        expr_df = pd.read_parquet(expr_path, columns=["CellID"] + available)
        merged = expr_df.merge(
            df[["CellID", "Cluster"]], on="CellID", how="inner"
        )
        cluster_means = merged.groupby("Cluster")[available].mean()
        cluster_means = cluster_means.reindex(
            sorted(cluster_means.index, key=_cluster_sort_key)
        )

        # Z-score 変換
        z_data = cluster_means.values.copy()
        col_mean = z_data.mean(axis=0)
        col_std = z_data.std(axis=0)
        col_std[col_std == 0] = 1
        z_data = (z_data - col_mean) / col_std

        # Y軸ラベル（アノテーション付き）
        y_labels = []
        for g in available:
            if g in gene_to_annotation:
                y_labels.append(f"{g} ({gene_to_annotation[g]})")
            elif mz_to_compound:
                annotated = _annotate_gene_labels(
                    [g], mz_to_compound, tolerance=0.1)
                y_labels.append(annotated[0])
            else:
                y_labels.append(g)

        fig = go.Figure(go.Heatmap(
            z=z_data.T,
            x=[str(c) for c in cluster_means.index],
            y=y_labels,
            colorscale="RdBu_r",
            zmid=0,
            hovertemplate=(
                "Cluster: %{x}<br>Gene: %{y}<br>"
                "Value: %{z:.3f}<extra></extra>"
            ),
        ))
        max_label_len = max(len(str(l)) for l in y_labels) if y_labels else 10
        left_margin = min(max(max_label_len * 7, 120), 350)
        fig.update_layout(
            title=dict(
                text=f"Top {top_n} DEG Heatmap (Z-score)",
                font=dict(size=14), x=0.5),
            xaxis_title="Cluster",
            yaxis_title="Gene / m/z",
            template="plotly_white",
            margin=dict(l=left_margin, r=20, t=40, b=40),
            yaxis=dict(autorange="reversed"),
        )
        return fig.to_dict()

    except Exception:
        return heatmap_fig  # エラー時はアプリ図にフォールバック


def _build_cluster_umap_panel_fig(df_s, cl_mask, cl_color, ux, uy,
                                  title=None, bg_gray="rgb(225,225,225)"):
    """単一サンプルの UMAP highlight パネル図（PPT 個別配置用）。

    ラスター（2D ヒストグラム→go.Heatmap）優先、失敗時は散布フォールバック。
    非ハイライトを薄グレー、ハイライトクラスタを cl_color で描く 1 パネル図。
    """
    fig = go.Figure()
    use_raster = _raster.raster_enabled()
    done = False
    if use_raster and len(df_s) > 0:
        _ucat = np.where(cl_mask.values, 1.0, 0.0)  # bg=0, hl=1
        _ures = _raster.umap_hist_grid(
            df_s["UMAP_1"].values, df_s["UMAP_2"].values, _ucat, ux, uy)
        if _ures is not None:
            _uz, _uxc, _uyc = _ures
            _ucs, _uzmin, _uzmax = _raster.build_discrete_colorscale(
                [bg_gray, cl_color])
            fig.add_trace(_raster.heatmap_trace(
                _uz, _uxc, _uyc, _ucs, _uzmin, _uzmax, showscale=False))
            done = True
    if not done:
        bg_df = df_s[~cl_mask]
        hl_df = df_s[cl_mask]
        if len(bg_df) > 0:
            fig.add_trace(go.Scattergl(
                x=bg_df["UMAP_1"], y=bg_df["UMAP_2"], mode="markers",
                marker=dict(color="lightgray", size=3, opacity=0.8),
                showlegend=False, hoverinfo="skip"))
        if len(hl_df) > 0:
            fig.add_trace(go.Scattergl(
                x=hl_df["UMAP_1"], y=hl_df["UMAP_2"], mode="markers",
                marker=dict(color=cl_color, size=3),
                showlegend=False, hoverinfo="skip"))
    fig.update_xaxes(showticklabels=False, showgrid=False, zeroline=False,
                     range=list(ux) if ux else None, visible=False)
    fig.update_yaxes(showticklabels=False, showgrid=False, zeroline=False,
                     range=list(uy) if uy else None,
                     scaleanchor="x", visible=False)
    fig.update_layout(
        showlegend=False, plot_bgcolor="white", paper_bgcolor="white",
        margin=dict(l=5, r=5, t=(44 if title else 5), b=5))
    if title:
        fig.update_layout(
            title=dict(text=str(title), x=0.5, font=dict(size=28)))
    return fig


def _build_cluster_spatial_panel_fig(df_s, cl_mask, cl_color, transform,
                                     title=None, bg_gray="rgb(225,225,225)"):
    """単一サンプルの Spatial highlight パネル図（PPT 個別配置用）。

    背景=TIC(TotalCount) のグレー濃淡、前景=highlight クラスタ単色 の 2 層。ラスター優先。
    向きは App(`_create_single_spatial_fig`)と一致させるため **y 軸は非反転**
    （`raw_y=-SpatialY` をデータ座標 heatmap で非反転軸に描く）→ App と PPT の上下反転を解消。
    cl_mask: df_s に整列した bool Series（True=highlight クラスタ）。
    """
    if "SpatialX" not in df_s.columns:
        return None
    fig = go.Figure()
    raw_x = df_s["SpatialX"].values.astype(float)
    raw_y = -df_s["SpatialY"].values.astype(float)  # Y軸反転（App と同じ）
    tx, ty = _transform_coords(
        raw_x, raw_y,
        transform.get("angle", 0),
        flip_h=transform.get("flip_h", False),
        flip_v=transform.get("flip_v", False))
    bg_mask_arr = (~cl_mask).values
    hl_mask_arr = cl_mask.values
    use_raster = _raster.raster_enabled()
    done = False
    if use_raster:
        _gi = _raster.grid_index(tx, ty)
        if _gi is not None:
            _six, _siy, _sxc, _syc = _gi
            _ny, _nx = len(_syc), len(_sxc)
            # --- 背景: TIC 濃淡（無ければ薄グレー一色）---
            if bg_mask_arr.any():
                _zbg = np.full((_ny, _nx), np.nan, dtype=float)
                _has_tic = ("TotalCount" in df_s.columns
                            and df_s["TotalCount"].notna().any())
                if _has_tic:
                    _tic = df_s["TotalCount"].values.astype(float)
                    _zbg[_siy[bg_mask_arr], _six[bg_mask_arr]] = _tic[bg_mask_arr]
                    _tf = _tic[bg_mask_arr]
                    _tf = _tf[np.isfinite(_tf)]
                    _bgcs = "Greys"
                    _bgmin = float(_tf.min()) if _tf.size else 0.0
                    _bgmax = float(_tf.max()) if _tf.size else 1.0
                    if _bgmax <= _bgmin:
                        _bgmax = _bgmin + 1.0
                else:
                    _zbg[_siy[bg_mask_arr], _six[bg_mask_arr]] = 0.0
                    _bgcs = [[0.0, bg_gray], [1.0, bg_gray]]
                    _bgmin, _bgmax = 0.0, 1.0
                fig.add_trace(_raster.heatmap_trace(
                    _zbg, _sxc, _syc, _bgcs, _bgmin, _bgmax, showscale=False))
            # --- 前景: highlight クラスタ（単色）を上に重ねる ---
            if hl_mask_arr.any():
                _zhl = np.full((_ny, _nx), np.nan, dtype=float)
                _zhl[_siy[hl_mask_arr], _six[hl_mask_arr]] = 1.0
                fig.add_trace(_raster.heatmap_trace(
                    _zhl, _sxc, _syc, [[0.0, cl_color], [1.0, cl_color]],
                    0.0, 1.0, showscale=False))
            done = True
    if not done:
        msize = _calc_zero_gap_marker_size(tx, ty, render_height=300)
        if bg_mask_arr.any():
            if "TotalCount" in df_s.columns and df_s["TotalCount"].notna().any():
                fig.add_trace(go.Scattergl(
                    x=tx[bg_mask_arr], y=ty[bg_mask_arr], mode="markers",
                    marker=dict(
                        color=df_s["TotalCount"].values[bg_mask_arr],
                        colorscale="Greys", size=msize,
                        symbol="square", opacity=0.5, showscale=False),
                    showlegend=False, hoverinfo="skip"))
            else:
                fig.add_trace(go.Scattergl(
                    x=tx[bg_mask_arr], y=ty[bg_mask_arr], mode="markers",
                    marker=dict(color="lightgray", size=msize,
                                symbol="square", opacity=0.2),
                    showlegend=False, hoverinfo="skip"))
        if hl_mask_arr.any():
            fig.add_trace(go.Scattergl(
                x=tx[hl_mask_arr], y=ty[hl_mask_arr], mode="markers",
                marker=dict(color=cl_color, size=msize, symbol="square"),
                showlegend=False, hoverinfo="skip"))
    fig.update_xaxes(showticklabels=False, showgrid=False, zeroline=False,
                     visible=False)
    # y 軸は非反転（App と一致 → ①解消）。データ座標 heatmap なので向きは保たれる。
    fig.update_yaxes(showticklabels=False, showgrid=False, zeroline=False,
                     scaleanchor="x", visible=False)
    fig.update_layout(
        showlegend=False, plot_bgcolor="white", paper_bgcolor="white",
        margin=dict(l=5, r=5, t=(44 if title else 5), b=5))
    if title:
        fig.update_layout(
            title=dict(text=str(title), x=0.5, font=dict(size=28)))
    return fig


def _add_conditions_slide(prs, conditions, rows_per_slide=16):
    """「解析条件」スライドを追加する（論文の Methods 用）。

    図と条件が同じファイルに入るので、後から「この図はどの設定か」が
    辿れなくなることがない。機械可読な全量はスピーカーノートに JSON で入れる
    （見た目を変えずに済み、python-pptx / PowerPoint から取り出せる）。
    """
    import json

    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    if not conditions:
        return
    try:
        from app.services.methods_text import render_conditions_rows
        rows = render_conditions_rows(conditions, lang="ja")
    except Exception:
        logger.warning("解析条件スライドの行生成に失敗", exc_info=True)
        return

    missing = conditions.get("_missing") or []
    warnings = conditions.get("warnings") or []
    col_w = [Inches(3.6), Inches(8.8)]
    n_pages = max(1, (len(rows) + rows_per_slide - 1) // rows_per_slide)
    first_slide = None
    for pi in range(n_pages):
        chunk = rows[pi * rows_per_slide:(pi + 1) * rows_per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if first_slide is None:
            first_slide = slide
        title = "解析条件 (Analysis Conditions)"
        if n_pages > 1:
            title += f"  {pi + 1}/{n_pages}"
        _pptx_add_title_bar(slide, title)
        n_rows = len(chunk) + 1
        tbl = slide.shapes.add_table(
            n_rows, 2, Inches(0.3), Inches(0.9),
            sum(int(w) for w in col_w),
            Inches(min(6.0, 0.32 * n_rows + 0.3))).table
        for j, w in enumerate(col_w):
            tbl.columns[j].width = w
        for j, h in enumerate(("項目", "値")):
            c = tbl.cell(0, j)
            c.text = h
            pr = c.text_frame.paragraphs[0]
            pr.font.size = Pt(11)
            pr.font.bold = True
        for i, (label, value) in enumerate(chunk):
            for j, val in enumerate((label, value)):
                c = tbl.cell(i + 1, j)
                c.text = str(val)
                c.text_frame.paragraphs[0].font.size = Pt(9)

    # 未記録・警告は最後のスライドの下部に明示する（黙って欠けさせない）
    if first_slide is not None and (missing or warnings):
        box = prs.slides[-1].shapes.add_textbox(
            Inches(0.3), Inches(6.4), Inches(12.4), Inches(0.9))
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        notes = []
        if missing:
            notes.append(f"未記録の項目 {len(missing)} 件: " + ", ".join(missing[:6])
                         + ("…" if len(missing) > 6 else ""))
        notes.extend(warnings)
        para.text = " / ".join(notes)
        para.font.size = Pt(9)
        para.font.color.rgb = RGBColor(0xCC, 0x44, 0x00)

    # 機械可読な全量をスピーカーノートへ
    if first_slide is not None:
        try:
            first_slide.notes_slide.notes_text_frame.text = json.dumps(
                conditions, ensure_ascii=False, indent=2, default=str)
        except Exception:
            logger.warning("スピーカーノートへの条件埋め込みに失敗", exc_info=True)


def _add_marker_table_slide(prs, title, headers, rows, rows_per_slide=18):
    """headers/rows を python-pptx のテーブルとしてスライドへ描画する（再利用ヘルパー）。

    列: クラスタ / m/z / 化合物名 / 方向(▲Up/▼Down) / log2FC / 調整p値。
    行数が rows_per_slide を超えると複数スライドへ自動改ページ（ヘッダー再掲）。
    行データ生成は呼び出し側で deg_utils.build_marker_rows を用いる。
    """
    from pptx.util import Inches, Pt

    if not rows:
        return
    col_w = [Inches(1.4), Inches(1.6), Inches(4.6), Inches(1.1),
             Inches(1.5), Inches(2.0)]
    total_w = sum(int(w) for w in col_w)
    n_pages = max(1, (len(rows) + rows_per_slide - 1) // rows_per_slide)
    for pi in range(n_pages):
        chunk = rows[pi * rows_per_slide:(pi + 1) * rows_per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _title = title
        if n_pages > 1:
            _title += f"  {pi + 1}/{n_pages}"
        _pptx_add_title_bar(slide, _title)
        n_rows = len(chunk) + 1
        tbl = slide.shapes.add_table(
            n_rows, len(headers), Inches(0.3), Inches(0.9),
            total_w, Inches(min(6.4, 0.3 * n_rows + 0.3))).table
        for j, w in enumerate(col_w):
            if j < len(tbl.columns):
                tbl.columns[j].width = w
        for j, h in enumerate(headers):
            c = tbl.cell(0, j)
            c.text = str(h)
            pr = c.text_frame.paragraphs[0]
            pr.font.size = Pt(11)
            pr.font.bold = True
        for i, row in enumerate(chunk):
            for j, val in enumerate(row):
                c = tbl.cell(i + 1, j)
                c.text = str(val)
                c.text_frame.paragraphs[0].font.size = Pt(9)


def _build_pptx(umap_fig, spatial_fig, meta, cluster_stats_data, rds_path,
                 sub_name="", volcano_fig=None, heatmap_fig=None,
                 deg_data=None, top_n=5, df=None, cache_dir=None,
                 custom_colors=None, rotation_store=None, name_map=None,
                 set_progress=None, mrm_path=None,
                 existing_prs=None, progress_offset=0, progress_total=None,
                 saved_positions=None, cluster_name_map=None,
                 include_deg=True, deadline=None, conditions=None,
                 display_settings=None):
    """グローバル概要 + クラスターごとの詳細スライドを含む PPTX を生成し bytes を返す。

    グローバルセクション:
        1. タイトル  2. UMAP+Spatial (統合)  3. クラスタ統計
    クラスターセクション:
        A. UMAP (ハイライト) + Spatial (ハイライト)  … 常設
        B. Volcano Plot + Top N Feature Plots        … include_deg=True の時のみ
        C. Heatmap (Top N, Z-score)                  … 常設
    include_deg=False の場合、B は出さず、代わりに全クラスタの上位 marker（m/z・化合物名・
    log2FC・調整p値）を集約した表スライド（行数に応じ自動改ページ）を末尾に追加する。

    existing_prs: 既存のPresentationオブジェクト。指定時はそこにスライドを追加し、
                  bytes は返さず None を返す。
    progress_offset: 進捗計算のオフセット（複数手法ループ時に使用）
    progress_total: 進捗計算の全体ステップ数（複数手法ループ時に使用）
    display_settings: `_display_settings(conditions)` の戻り値。図の閾値・スケール・
                  強度レンジを「解析条件」スライドと揃えるために使う (ver51.9)。
                  None なら conditions から自前で導く（従来の呼び出しも壊さない）。
    """
    import time

    # ★ ver51.9: 条件スライドと図が同じ設定を見るようにする。
    #   呼び出し側が渡さなかった場合も conditions から導けるようにしておく
    #   （どちらも無ければ既定 = 従来の図と同じ）。
    if display_settings is None:
        display_settings = _display_settings(conditions)

    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    # 進捗計算用
    _clusters_for_progress = []
    if df is not None:
        try:
            _clusters_for_progress = sorted(
                set(str(c) for c in df["Cluster"].unique()),
                key=_cluster_sort_key)
        except Exception:
            pass
    # 各クラスタ 3 ステップ: A(UMAP/Spatial) + [B(DEG) or per-cluster marker 表] + C(Heatmap)
    _per_cluster = 3
    _local_steps = 3 + len(_clusters_for_progress) * _per_cluster
    _total_steps = progress_total if progress_total else _local_steps
    _current_step = [progress_offset]  # mutable for nested function

    def _progress(label=""):
        _current_step[0] += 1
        # サーバログにもハートビートを残す（ブラウザ進捗バーだけだと停止箇所が追えないため）
        logger.info("[PPTX] %s (%d/%d)", label, _current_step[0], _total_steps)
        if set_progress:
            pct = int(_current_step[0] / _total_steps * 100)
            set_progress((min(pct, 99), 100, label))
        # 全体 watchdog: 期限超過なら中断（1 枚のタイムアウトだけでは総時間が伸び得るため）
        if deadline is not None and time.monotonic() > deadline:
            raise TimeoutError(
                f"PPTX export exceeded overall timeout "
                f"(step {_current_step[0]}/{_total_steps})")

    if existing_prs is not None:
        prs = existing_prs
    else:
        prs = Presentation()
        # 16:9 ワイドスクリーン
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # =====================================================================
    # グローバルセクション
    # =====================================================================

    # --- スライド 1: タイトル ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True

    if sub_name:
        p_title = tf.paragraphs[0]
        p_title.text = sub_name
        p_title.font.size = Pt(40)
        p_title.font.bold = True
        p_title.alignment = PP_ALIGN.CENTER

        p_sub = tf.add_paragraph()
        p_sub.text = "MSI Interactive Analysis Report"
        p_sub.font.size = Pt(20)
        p_sub.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p_sub.alignment = PP_ALIGN.CENTER
    else:
        p_title = tf.paragraphs[0]
        p_title.text = "MSI Interactive Analysis Report"
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p2.alignment = PP_ALIGN.CENTER

    if rds_path:
        p3 = tf.add_paragraph()
        p3.text = f"Source: {Path(rds_path).name}"
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p3.alignment = PP_ALIGN.CENTER

    if meta:
        p4 = tf.add_paragraph()
        samples_str = ", ".join(meta.get("samples", []))
        p4.text = (
            f"Cells: {meta.get('n_cells', '?')} | "
            f"Clusters: {meta.get('n_clusters', '?')} | "
            f"Samples: {samples_str}"
        )
        p4.font.size = Pt(12)
        p4.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p4.alignment = PP_ALIGN.CENTER

    _progress("タイトルスライド")

    # --- 解析条件スライド ---
    # 複数手法パスでは呼び出し元が先頭に 1 回だけ入れるため、ここでは単一手法時のみ。
    if conditions is not None and existing_prs is None:
        _add_conditions_slide(prs, conditions)

    # --- スライド 2: UMAP + Spatial 統合 (サンプル別) ---
    if df is not None and not df.empty and "SpatialX" in df.columns:
        import math
        color_map_global = _get_cluster_color_map(df["Cluster"], custom_colors)
        all_samples = sorted(df["Sample"].unique())
        if not rotation_store:
            rotation_store = {}
        if not name_map:
            name_map = {}

        if all_samples:
            n_sp = len(all_samples)
            # レジェンド幅を確保しつつUMAP/Spatialの横幅を最大化
            _legend_w = 1.3
            avail_w = 13.33 - 0.3 - _legend_w - 0.1  # ≈ 11.93"

            # サンプル数が多い場合はスライドを分割
            _MIN_TILE_W = 1.5  # タイル最小幅(インチ)
            _max_per_slide = max(1, int(avail_w / _MIN_TILE_W))
            if n_sp > _max_per_slide:
                _mid = (n_sp + 1) // 2
                _sample_groups = [all_samples[:_mid], all_samples[_mid:]]
            else:
                _sample_groups = [all_samples]

            # UMAP の全データ軸範囲（サンプル間で統一）
            umap1_min = float(df["UMAP_1"].min())
            umap1_max = float(df["UMAP_1"].max())
            umap2_min = float(df["UMAP_2"].min())
            umap2_max = float(df["UMAP_2"].max())
            umap_pad = max(umap1_max - umap1_min, umap2_max - umap2_min) * 0.05
            umap_xrange = [umap1_min - umap_pad, umap1_max + umap_pad]
            umap_yrange = [umap2_min - umap_pad, umap2_max + umap_pad]

            # レジェンド画像を事前生成（各スライドで共有）
            _n_cl = len(df["Cluster"].unique())
            _legend_h_est = min(_n_cl * 0.35 + 0.2, 6.0)
            legend_fig = _build_cluster_legend_fig(
                df["Cluster"].unique(), color_map_global,
                cluster_name_map=cluster_name_map)
            legend_png = _fig_to_png_bytes(
                legend_fig.to_dict(), width=200, height=600, scale=2)

            for _grp_idx, _grp_samples in enumerate(_sample_groups):
                _grp_n = len(_grp_samples)
                _grp_suffix = f" ({_grp_idx + 1}/{len(_sample_groups)})" if len(_sample_groups) > 1 else ""
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _pptx_add_title_bar(slide, f"UMAP & Spatial Mapping{_grp_suffix}")

                # 上段: サンプル別UMAP (y=0.9" 〜 y=4.0", 高さ3.1")
                tile_w_umap = avail_w / _grp_n
                for idx, s in enumerate(_grp_samples):
                    df_s = df[df["Sample"] == s]
                    _umap_pos = (saved_positions or {}).get("umap_integrated", {})
                    umap_s = _build_umap_integrated_fig(
                        df_s, color_by="Cluster", highlight_clusters=None,
                        show_legend=False, show_labels=True,
                        title=_display_name(s, name_map),
                        marker_size=2, custom_colors=custom_colors,
                        title_font_size=40, label_size=24,
                        saved_positions=_umap_pos,
                        cluster_name_map=cluster_name_map)
                    if umap_s is not None:
                        umap_s.update_xaxes(range=umap_xrange)
                        umap_s.update_yaxes(range=umap_yrange)
                        u_dict = (umap_s.to_dict()
                                  if hasattr(umap_s, "to_dict") else umap_s)
                        u_png = _fig_to_png_bytes(
                            u_dict, width=600, height=600, scale=2)
                        _uw, _uh, _uoff = _square_tile_dims(
                            tile_w_umap, 3.0)
                        u_left = Inches(0.3 + idx * tile_w_umap + _uoff)
                        _pptx_add_image(slide, u_png,
                                        int(u_left), Inches(0.9),
                                        Inches(_uw), Inches(_uh))

                # 下段: サンプル別Spatial (y=4.1" 〜 y=7.3", 高さ3.2")
                tile_w_sp = avail_w / _grp_n
                for idx, s in enumerate(_grp_samples):
                    df_s = df[df["Sample"] == s]
                    transform = rotation_store.get(
                        s, rotation_store.get(
                            "__all__",
                            {"angle": 0, "flip_h": False, "flip_v": False}))
                    if isinstance(transform, (int, float)):
                        transform = {"angle": int(transform),
                                     "flip_h": False, "flip_v": False}

                    _sp_pos = (saved_positions or {}).get("spatial", {}).get(s, {})
                    sp_fig = _create_single_spatial_fig(
                        df_s, color_map_global,
                        highlight_clusters=None,
                        selected_cell_ids=set(),
                        rotation_deg=transform.get("angle", 0),
                        show_labels=True,
                        flip_h=transform.get("flip_h", False),
                        flip_v=transform.get("flip_v", False),
                        title=_display_name(s, name_map),
                        marker_size=0,
                        render_height=560,
                        embed_legend=False,
                        title_font_size=40, label_size=24,
                        saved_positions=_sp_pos,
                        cluster_name_map=cluster_name_map)
                    if sp_fig is not None:
                        sp_dict = (sp_fig.to_dict()
                                   if hasattr(sp_fig, "to_dict") else sp_fig)
                        sp_png = _fig_to_png_bytes(
                            sp_dict, width=600, height=600, scale=2)
                        _sw, _sh, _soff = _square_tile_dims(
                            tile_w_sp, 3.1)
                        sp_left = Inches(0.3 + idx * tile_w_sp + _soff)
                        _pptx_add_image(slide, sp_png,
                                        int(sp_left), Inches(4.1),
                                        Inches(_sw), Inches(_sh))

                # クラスタレジェンド（右下角に配置、白余白最小化）
                _legend_x = Inches(0.3 + avail_w + 0.1)
                _legend_y = Inches(7.5 - _legend_h_est)
                _pptx_add_image_preserve_ratio(slide, legend_png,
                                               int(_legend_x), int(_legend_y),
                                               Inches(_legend_w), Inches(_legend_h_est),
                                               png_w=200, png_h=600)

            # --- UMAP (by Sample) スライド --- 削除済み（不要）

    elif spatial_fig:
        # SpatialX がない場合はアプリの図をフォールバック
        fig_check = go.Figure(spatial_fig)
        if fig_check.data:
            png_bytes = _fig_to_png_bytes(spatial_fig, width=1200, height=800)
            if png_bytes:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _pptx_add_title_bar(slide, "Spatial Mapping")
                _pptx_add_image_preserve_ratio(
                    slide, png_bytes,
                    int((prs.slide_width - Inches(12)) / 2), Inches(0.9),
                    Inches(12), Inches(6.3),
                    png_w=1200, png_h=800)

    _progress("UMAP & Spatial")

    # --- スライド 5: クラスタ統計 ---
    if cluster_stats_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _pptx_add_title_bar(slide, "Cluster Statistics")

        n_rows = len(cluster_stats_data) + 1
        n_cols = 3
        table_w = Inches(6)
        table_h = Inches(min(5.5, 0.4 * n_rows))
        left = int((prs.slide_width - table_w) / 2)
        table_shape = slide.shapes.add_table(
            n_rows, n_cols, left, Inches(1.0), table_w, table_h,
        )
        table = table_shape.table

        for j, h in enumerate(["Cluster", "Pixels", "%"]):
            cell = table.cell(0, j)
            cell.text = h
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.bold = True

        for i, row_data in enumerate(cluster_stats_data):
            table.cell(i + 1, 0).text = str(row_data.get("Cluster", ""))
            table.cell(i + 1, 1).text = str(row_data.get("Pixels", ""))
            table.cell(i + 1, 2).text = str(row_data.get("Percent", ""))
            for j in range(n_cols):
                table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(11)

    _progress("クラスタ統計")

    # =====================================================================
    # クラスターセクション（各クラスター × 3 スライド）
    # A. UMAP & Spatial  B. DEG Analysis  C. Heatmap
    # =====================================================================
    if df is not None and not df.empty:
        clusters = sorted(df["Cluster"].unique(), key=_cluster_sort_key)
        color_map = _get_cluster_color_map(df["Cluster"], custom_colors)
        has_spatial = "SpatialX" in df.columns
        samples = sorted(df["Sample"].unique()) if has_spatial else []
        if not rotation_store:
            rotation_store = {}
        if not name_map:
            name_map = {}

        # UMAP 軸範囲（Per-cluster Slide B でサンプル間統一に使用）
        cl_umap1_min = float(df["UMAP_1"].min())
        cl_umap1_max = float(df["UMAP_1"].max())
        cl_umap2_min = float(df["UMAP_2"].min())
        cl_umap2_max = float(df["UMAP_2"].max())
        cl_umap_pad = max(cl_umap1_max - cl_umap1_min,
                          cl_umap2_max - cl_umap2_min) * 0.05
        cl_umap_xrange = [cl_umap1_min - cl_umap_pad,
                          cl_umap1_max + cl_umap_pad]
        cl_umap_yrange = [cl_umap2_min - cl_umap_pad,
                          cl_umap2_max + cl_umap_pad]

        rds_path_str = str(rds_path) if rds_path else None
        cache_dir_path = (
            Path(cache_dir) if isinstance(cache_dir, str) and cache_dir
            else cache_dir
        )

        # DEG 非選択時、per-cluster marker 表で使う {mz: 化合物名}（MRM 近傍一致）。
        # ループ前に一度だけ生成して各クラスタへ渡す。
        _marker_mz_to_comp = {}
        if not include_deg and mrm_path:
            try:
                _marker_mz_to_comp = _build_mz_to_compound_map(
                    mrm_path, tolerance=0.1) or {}
            except Exception:
                _marker_mz_to_comp = {}
        _any_marker_slide = False  # marker 表を1枚でも出したか（末尾の note 判定用）

        # === cluster loop（描画は共有単一プールで逐次実行）===
        # ラスター化で 1 枚 1〜2 秒のため並列は不要。並列 kaleido/Chromium は競合して
        # 1 枚ごとにハング（→タイムアウトでスキップ）＋プロセスリークの原因だったため撤去。
        for cl in clusters:
            cl_str = str(cl)
            _cl_name = _cluster_display_name(cl_str, cluster_name_map)

            # === Slide A: UMAP + Spatial（サンプル毎の個別画像を1スライドへ）===
            # 結合1枚PNGではなく各サンプルの UMAP(上段)/Spatial(下段) を個別オブジェクトで
            # 配置し、PowerPoint 上で移動/リサイズ可能にする（②）。Spatial は非反転軸で
            # 描画するため向きが App と一致する（①解消）。
            slide_a = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_add_title_bar(slide_a, f"{_cl_name} — UMAP & Spatial")

            n_sp_b = len(samples) if has_spatial and samples else 0

            if n_sp_b > 0:
                _cl_color_a = (custom_colors or {}).get(
                    cl_str, color_map.get(cl_str, "#1f77b4"))
                _legend_w_a = 1.3
                _avail_w_a = 13.33 - 0.3 - _legend_w_a - 0.1  # ≈ 11.93"
                # サンプル多数時はスライド分割（タイル最小幅 1.5"）
                _max_per_a = max(1, int(_avail_w_a / 1.5))
                if n_sp_b > _max_per_a:
                    _mid_a = (n_sp_b + 1) // 2
                    _grps_a = [samples[:_mid_a], samples[_mid_a:]]
                else:
                    _grps_a = [samples]
                # 共有凡例（当該クラスタ + Unselected）
                _legend_fig_a = _build_cluster_legend_fig(
                    [cl_str, "Unselected"],
                    {cl_str: _cl_color_a, "Unselected": "rgb(190,190,190)"},
                    cluster_name_map=cluster_name_map)
                _legend_png_a = _fig_to_png_bytes(
                    _legend_fig_a.to_dict(), width=200, height=600, scale=2)
                _lh_a = min(2 * 0.35 + 0.4, 2.0)
                for _gi_a, _grp_a in enumerate(_grps_a):
                    _gn_a = len(_grp_a)
                    if _gi_a == 0:
                        _slide_cur = slide_a
                    else:
                        _slide_cur = prs.slides.add_slide(
                            prs.slide_layouts[6])
                        _pptx_add_title_bar(
                            _slide_cur,
                            f"{_cl_name} — UMAP & Spatial "
                            f"({_gi_a + 1}/{len(_grps_a)})")
                    _tile_w_a = _avail_w_a / max(_gn_a, 1)
                    for _idx_a, s in enumerate(_grp_a):
                        df_s = df[df["Sample"] == s]
                        _cl_mask_s = df_s["Cluster"].astype(str) == cl_str
                        transform = rotation_store.get(
                            s, rotation_store.get(
                                "__all__",
                                {"angle": 0, "flip_h": False,
                                 "flip_v": False}))
                        if isinstance(transform, (int, float)):
                            transform = {"angle": int(transform),
                                         "flip_h": False, "flip_v": False}
                        # 上段: UMAP（y=0.9"〜, 高さ3.0"）
                        _up_fig = _build_cluster_umap_panel_fig(
                            df_s, _cl_mask_s, _cl_color_a,
                            cl_umap_xrange, cl_umap_yrange,
                            title=_display_name(s, name_map))
                        if _up_fig is not None:
                            _up_png = _fig_to_png_bytes(
                                _up_fig.to_dict(), width=600, height=600,
                                scale=2)
                            _uw, _uh, _uoff = _square_tile_dims(_tile_w_a, 3.0)
                            _u_left = Inches(0.3 + _idx_a * _tile_w_a + _uoff)
                            _pptx_add_image(_slide_cur, _up_png,
                                            int(_u_left), Inches(0.9),
                                            Inches(_uw), Inches(_uh))
                        # 下段: Spatial（y=4.1"〜, 高さ3.1"）
                        _sp_fig = _build_cluster_spatial_panel_fig(
                            df_s, _cl_mask_s, _cl_color_a, transform,
                            title=f"{_display_name(s, name_map)} (Cl {cl_str})")
                        if _sp_fig is not None:
                            _sp_png = _fig_to_png_bytes(
                                _sp_fig.to_dict(), width=600, height=600,
                                scale=2)
                            _sw, _sh, _soff = _square_tile_dims(_tile_w_a, 3.1)
                            _s_left = Inches(0.3 + _idx_a * _tile_w_a + _soff)
                            _pptx_add_image(_slide_cur, _sp_png,
                                            int(_s_left), Inches(4.1),
                                            Inches(_sw), Inches(_sh))
                    # クラスタ凡例（右下角）
                    if _legend_png_a:
                        _lx_a = Inches(0.3 + _avail_w_a + 0.1)
                        _ly_a = Inches(7.5 - _lh_a)
                        _pptx_add_image_preserve_ratio(
                            _slide_cur, _legend_png_a,
                            int(_lx_a), int(_ly_a),
                            Inches(_legend_w_a), Inches(_lh_a),
                            png_w=200, png_h=600)
            else:
                # Spatialデータなし → 単一UMAP（従来互換）
                umap_hl = _build_umap_integrated_fig(
                    df, color_by="Cluster", highlight_clusters=[cl_str],
                    show_legend=True, show_labels=False,
                    marker_size=2, custom_colors=custom_colors,
                    bg_opacity=1.0,
                    cluster_name_map=cluster_name_map)
                if umap_hl is not None:
                    umap_dict = (umap_hl.to_dict()
                                 if hasattr(umap_hl, "to_dict") else umap_hl)
                    upng = _fig_to_png_bytes(
                        umap_dict, width=800, height=800, scale=2)
                    _pptx_add_image(slide_a, upng,
                                    Inches(0.3), Inches(0.7),
                                    Inches(4.5), Inches(4.5))

            _progress(f"{_cl_name} — UMAP/Spatial")

            # === Slide A2: この cluster の marker 表（DEG 非選択時のみ）===
            # UMAP & Spatial の次スライドに、当該クラスタの上位 marker を表で出す（④）。
            if not include_deg:
                from app.utils.deg_utils import build_marker_rows
                _mk_headers, _mk_rows = build_marker_rows(
                    [cl_str], deg_data, top_n=top_n,
                    mz_to_compound=_marker_mz_to_comp,
                    cluster_name_map=cluster_name_map)
                if _mk_rows:
                    _add_marker_table_slide(
                        prs, f"{_cl_name} — Markers (Top {top_n})",
                        _mk_headers, _mk_rows)
                    _any_marker_slide = True
                _progress(f"{_cl_name} — Markers")

            # === Slide B: Volcano + Feature Plots ===（DEG 有効時のみ）
            if include_deg:
                slide_b = prs.slides.add_slide(prs.slide_layouts[6])
                _pptx_add_title_bar(slide_b, f"{_cl_name} — DEG Analysis")

                # Top N features (up / down)
                up_features, down_features = _get_top_n_features_for_cluster(
                    deg_data, cl_str, n=top_n)

                # gene→annotation マッピング（Feature Plotタイトル用）
                _gene_ann_map = {}
                if deg_data:
                    for _r in deg_data:
                        _g = _r.get("gene", "")
                        _a = _r.get("annotation", "")
                        if _g and _is_meaningful_annotation(_a, _g):
                            _gene_ann_map[_g] = _a

                # Volcano Plot を生成（レイアウト計算に必要）。描画は下段配置時に逐次実行。
                # ★ ver51.9: 閾値は「解析条件」スライドと同じものを使う。
                #   従来は既定値のままだったため、同じ資料の中で条件と図が食い違っていた。
                volcano_cl = _build_volcano_fig_for_cluster(
                    deg_data, cl_str,
                    fc_thresh=display_settings["volcano_fc"],
                    p_thresh=display_settings["volcano_p"],
                    label_top_n=display_settings["volcano_top_n"])

                # ---- 自動レイアウト計算 ----
                has_up = bool(up_features)
                has_down = bool(down_features)
                has_volcano = volcano_cl is not None

                _avail_top = 0.65      # Feature配置開始Y (タイトルバー下)
                _avail_bottom = 7.35   # スライド下端マージン
                _avail_h = _avail_bottom - _avail_top  # 6.7"
                _label_h = 0.25        # ラベル行の高さ
                _gap = 0.1             # セクション間隙間
                _volcano_h = 2.5       # Volcano固定高さ

                _n_feat_rows = (1 if has_up else 0) + (1 if has_down else 0)

                if _n_feat_rows > 0 and has_volcano:
                    _non_feat = _n_feat_rows * (_label_h + _gap) + _gap + _volcano_h
                    _feat_h_val = (_avail_h - _non_feat) / _n_feat_rows
                elif _n_feat_rows > 0:
                    _non_feat = _n_feat_rows * (_label_h + _gap)
                    _feat_h_val = (_avail_h - _non_feat) / _n_feat_rows
                else:
                    _feat_h_val = 0

                _feat_h_val = max(1.5, min(_feat_h_val, 3.5))

                feat_w_val = min(2.4, 12.0 / max(top_n, 1))
                feat_w = Inches(feat_w_val)
                feat_h = Inches(_feat_h_val)

                # Y座標を順番に計算
                _cur_y = _avail_top
                _up_label_y = _up_plot_y = 0
                _down_label_y = _down_plot_y = 0
                _volcano_y = 0

                if has_up:
                    _up_label_y = _cur_y
                    _cur_y += _label_h + _gap
                    _up_plot_y = _cur_y
                    _cur_y += _feat_h_val + _gap
                if has_down:
                    _down_label_y = _cur_y
                    _cur_y += _label_h + _gap
                    _down_plot_y = _cur_y
                    _cur_y += _feat_h_val + _gap
                if has_volcano:
                    _volcano_y = _cur_y

                # "▲ Up-regulated" ラベル
                if has_up:
                    lbl = slide_b.shapes.add_textbox(
                        Inches(0.3), Inches(_up_label_y), Inches(3), Inches(0.3))
                    lp = lbl.text_frame.paragraphs[0]
                    lp.text = f"▲ Up-regulated (Top {len(up_features)})"
                    lp.font.size = Pt(10)
                    lp.font.color.rgb = RGBColor(0xFF, 0x2D, 0x2D)
                    lp.font.bold = True

                # Feature Plot 画像配置 — Up（逐次: 構築→描画→配置）
                for i, feat in enumerate(up_features):
                    is_last_up = (i == len(up_features) - 1)
                    feat_fig = _build_feature_plot_fig(
                        df, feat, cache_dir_path, rds_path_str,
                        rotation_store, name_map, marker_size=2,
                        show_colorbar_title=is_last_up,
                        show_colorbar=is_last_up,
                        auto_marker_size=True,
                        intensity_min=display_settings["feature_intensity_min"],
                        intensity_max=display_settings["feature_intensity_max"])
                    if not feat_fig:
                        continue
                    _feat_title = feat
                    if feat in _gene_ann_map:
                        _feat_title = f"{feat}\n({_gene_ann_map[feat]})"
                    feat_fig.update_layout(
                        title=dict(text=_feat_title, font=dict(size=14), x=0.5),
                        margin=dict(t=70))
                    png = _fig_to_png_bytes(feat_fig.to_dict(), 400, 280, 2)
                    if not png:
                        continue
                    left = Inches(0.3 + i * (12.5 / max(top_n, 1)))
                    _pptx_add_image_preserve_ratio(
                        slide_b, png,
                        int(left), Inches(_up_plot_y),
                        feat_w, feat_h,
                        png_w=400, png_h=280)

                # Up features 間の縦区切り線
                _tile_w_feat = 12.5 / max(top_n, 1)
                for i in range(1, len(up_features)):
                    _sep_x = Inches(0.3 + i * _tile_w_feat)
                    _sep = slide_b.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        int(_sep_x) - Emu(6350), Inches(_up_plot_y),
                        Emu(6350), feat_h)
                    _sep.fill.solid()
                    _sep.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                    _sep.line.fill.background()

                # "▼ Down-regulated" ラベル
                if has_down:
                    lbl = slide_b.shapes.add_textbox(
                        Inches(0.3), Inches(_down_label_y), Inches(3), Inches(0.3))
                    lp = lbl.text_frame.paragraphs[0]
                    lp.text = f"▼ Down-regulated (Top {len(down_features)})"
                    lp.font.size = Pt(10)
                    lp.font.color.rgb = RGBColor(0x1E, 0x5B, 0xFF)
                    lp.font.bold = True

                # Feature Plot 画像配置 — Down（逐次: 構築→描画→配置）
                for i, feat in enumerate(down_features):
                    is_last_down = (i == len(down_features) - 1)
                    feat_fig = _build_feature_plot_fig(
                        df, feat, cache_dir_path, rds_path_str,
                        rotation_store, name_map, marker_size=2,
                        show_colorbar_title=is_last_down,
                        show_colorbar=is_last_down,
                        auto_marker_size=True,
                        intensity_min=display_settings["feature_intensity_min"],
                        intensity_max=display_settings["feature_intensity_max"])
                    if not feat_fig:
                        continue
                    _feat_title = feat
                    if feat in _gene_ann_map:
                        _feat_title = f"{feat}\n({_gene_ann_map[feat]})"
                    feat_fig.update_layout(
                        title=dict(text=_feat_title, font=dict(size=14), x=0.5),
                        margin=dict(t=70))
                    png = _fig_to_png_bytes(feat_fig.to_dict(), 400, 280, 2)
                    if not png:
                        continue
                    left = Inches(0.3 + i * (12.5 / max(top_n, 1)))
                    _pptx_add_image_preserve_ratio(
                        slide_b, png,
                        int(left), Inches(_down_plot_y),
                        feat_w, feat_h,
                        png_w=400, png_h=280)

                # Down features 間の縦区切り線
                for i in range(1, len(down_features)):
                    _sep_x = Inches(0.3 + i * _tile_w_feat)
                    _sep = slide_b.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        int(_sep_x) - Emu(6350), Inches(_down_plot_y),
                        Emu(6350), feat_h)
                    _sep.fill.solid()
                    _sep.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                    _sep.line.fill.background()

                # Volcano Plot (下段) — 逐次描画して配置
                if has_volcano:
                    v_dict = (volcano_cl.to_dict()
                              if hasattr(volcano_cl, "to_dict") else volcano_cl)
                    vpng = _fig_to_png_bytes(v_dict, 800, 700, 2)
                    v_aspect = 800 / 700  # ≈ 1.14
                    v_height = Inches(_volcano_h)
                    v_width = int(v_height * v_aspect)
                    v_left = int((prs.slide_width - v_width) / 2)
                    _pptx_add_image(slide_b, vpng,
                                    v_left, Inches(_volcano_y), v_width, v_height)
                elif not up_features and not down_features:
                    # DEG データなし → 注釈
                    no_deg = slide_b.shapes.add_textbox(
                        Inches(3), Inches(3), Inches(7), Inches(1))
                    np_ = no_deg.text_frame.paragraphs[0]
                    np_.text = "No DEG data available for this cluster"
                    np_.font.size = Pt(16)
                    np_.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
                    np_.alignment = PP_ALIGN.CENTER

                _progress(f"{_cl_name} — DEG")

            # === Slide C: Per-cluster Heatmap (Top N, Z-score) ===
            # ★ ver51.9: スケール切替 (Z-score / Raw) を画面と揃える。
            #   従来は常に Z-score だったため、Raw を選んで確認した利用者が
            #   Z-score の図を配ることになっていた（見出しは同じ "Heatmap"）。
            _hm_scale = display_settings["heatmap_scale"]
            hm_fig = _build_heatmap_for_cluster(
                deg_data, cl_str, df, cache_dir, top_n, mrm_path=mrm_path,
                scale=_hm_scale)
            if hm_fig:
                hm_png = _fig_to_png_bytes(hm_fig, width=1200, height=800)
                if hm_png:
                    _hm_kind = "Raw" if _hm_scale == "raw" else "Z-score"
                    slide_c = prs.slides.add_slide(prs.slide_layouts[6])
                    _pptx_add_title_bar(
                        slide_c,
                        f"{_cl_name} — Heatmap (Top {top_n}, {_hm_kind})")
                    _pptx_add_image_preserve_ratio(
                        slide_c, hm_png,
                        int((prs.slide_width - Inches(12)) / 2), Inches(0.9),
                        Inches(12), Inches(6.3),
                        png_w=1200, png_h=800)

            _progress(f"{_cl_name} — Heatmap")

        # DEG 非選択時、どの cluster にも marker が無い（DEG 未実行等）場合のみ、
        # 空の note を 1 枚だけ出す（per-cluster 表はクラスタループ内で出力済み）。
        if not include_deg and not _any_marker_slide:
            _note_slide = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_add_title_bar(_note_slide, "Marker 一覧")
            _nbox = _note_slide.shapes.add_textbox(
                Inches(1), Inches(3), Inches(11), Inches(1))
            _np = _nbox.text_frame.paragraphs[0]
            _np.text = "DEG データがありません。"
            _np.font.size = Pt(16)
            _np.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # existing_prs が渡された場合は呼び出し元がまとめて保存するため
    # ここでは保存しない (現在のステップ数を返す)
    if existing_prs is not None:
        return _current_step[0]

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


# ---------------------------------------------------------------------------
# PPTX エクスポート コールバック
# ---------------------------------------------------------------------------

@callback(
    Output("export_top_n_store", "data"),
    Input("input_export_top_n", "value"),
    prevent_initial_call=True,
)
def sync_export_top_n(value):
    """dbc.Input → dcc.Store ブリッジ (Top N)"""
    # ★ ver52.3 ⑤: 既定値をここに直接書くと、同じ入力の既定が 2 箇所になる。
    #   `input_export_top_n` は min=1 なので 0 は範囲外だが、
    #   「不正な入力を黙って既定値にする」形は同じなので出典を 1 つにする。
    return coerce_count(value, "input_export_top_n")


@callback(
    [Output("export_method_selector", "options"),
     Output("export_method_selector", "value")],
    Input("interactive_rds_map", "data"),
    prevent_initial_call=True,
)
def update_export_method_options(rds_map):
    """rds_map の変更に応じてエクスポート対象手法セレクタ（複数選択）を更新する。

    Checklist 化に伴い「Both」は廃止。既定は実在する全手法をチェック済みにする
    （全部欲しければ全部にマークが付いた状態）。value はリスト。
    """
    if not rds_map or not isinstance(rds_map, dict):
        return [], []

    methods = list(rds_map.keys())
    options = [{"label": m, "value": m} for m in methods]
    return options, methods  # 既定で全手法チェック


@callback(
    [Output("dl_report_pptx", "data"),
     Output("div_export_status", "children")],
    Input("btn_export_report", "n_clicks"),
    [State("interactive_umap_plot", "figure"),
     State("last_spatial_figure_store", "data"),
     State("seurat_rds_path_store", "data"),
     State("cluster_stats_table", "data"),
     State("interactive_sub_project_select", "options"),
     State("interactive_sub_project_select", "value"),
     State("volcano_plot", "figure"),
     State("heatmap_plot", "figure"),
     State("deg_data_store", "data"),
     State("custom_color_map_store", "data"),
     State("spatial_rotation_store", "data"),
     State("sample_name_map_store", "data"),
     State("export_top_n_store", "data"),
     State("seurat_cache_dir_store", "data"),
     State("annotation_path", "value"),
     State("interactive_rds_map", "data"),
     State("interactive_result_folder", "value"),
     State("interactive_integration_method", "value"),
     State("export_method_selector", "value"),
     State("export_include_deg", "value"),
     State("cluster_name_map_store", "data"),
     State("accumulated_label_positions", "data")],
    background=True,
    running=[
        (Output("btn_export_report", "disabled"), True, False),
        (Output("export_progress_container", "style"),
         {"display": "block"}, {"display": "none"}),
    ],
    progress=[
        Output("export_progress_bar", "value"),
        Output("export_progress_bar", "max"),
        Output("export_progress_label", "children"),
    ],
    prevent_initial_call=True,
)
def cb_export_report(set_progress, n_clicks, umap_fig, spatial_fig, rds_path,
                     cluster_stats_data, sub_options, sub_value,
                     volcano_fig, heatmap_fig, deg_data, custom_colors,
                     rotation_store, name_map, top_n, cache_dir_str,
                     mrm_path_str, rds_map, result_folder, current_method,
                     export_method_selection, include_deg, cluster_name_map,
                     accumulated_positions):
    """PPTX レポートをバックグラウンド生成してダウンロード。

    export_method_selection: 選択された手法名のリスト（Checklist）。空なら中止。
    include_deg: DEG(Volcano+Feature)スライドを含めるか。False の場合は代わりに
        m/z・化合物名の marker 一覧表スライドを出力する。UMAP/Spatial・Heatmap は常時。
    """
    if not n_clicks:
        raise PreventUpdate

    set_progress((0, 100, "準備中..."))

    # ★ ver51.8: この回のエクスポートで何枚の図が描画に失敗したかを数えるため、
    #   開始時にカウンタを空にする（完了メッセージで枚数を伝える）。
    try:
        from app.utils.pptx_helpers import reset_render_failures
        reset_render_failures()
    except Exception:
        pass

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        return no_update, (
            "python-pptx がインストールされていません。"
            "pip install python-pptx を実行してください。"
        )

    try:
        import kaleido  # noqa: F401
    except ImportError:
        return no_update, (
            "kaleido がインストールされていません。"
            "pip install kaleido を実行してください。"
        )

    if not umap_fig:
        return no_update, "UMAPプロットが見つかりません。データを読み込んでください。"

    try:
        # サブプロジェクト名を取得
        sub_name = ""
        if sub_options and sub_value:
            for opt in sub_options:
                if opt.get("value") == sub_value:
                    sub_name = opt.get("label", "")
                    break

        # ファイル名を決定
        if sub_name:
            safe_name = re.sub(r'[\\/*?:"<>|]', '_', sub_name)
            filename = f"{safe_name}.pptx"
        else:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"MSI_Report_{timestamp}.pptx"

        top_n = top_n or 5
        # ★ ver51.9: rds_path / method を明示する。この callback は
        #   `background=True` で **DiskcacheManager が fork した別プロセス**で動くため、
        #   `_interactive_data` は空。引数なしだと JSON の解決先が None になり、
        #   **現在の手法だけラベル位置が既定に戻る**（他手法は下の
        #   `_load_label_positions_util(rds, method)` で正しく出るので、
        #   1 つの資料の中でラベルの付き方が不揃いになる）。
        #   正しい形は同ファイル :2020/:2178 と interactive_spatial.py:1161 にある。
        saved_positions = _get_merged_label_positions(
            accumulated_positions, rds_path=rds_path, method=current_method)

        # 解析条件は **クリック時点** で確定させる。PPTX 生成は数分かかることが
        # あり、その間にユーザーが設定を変えると、出力と条件がずれてしまう。
        conditions = None
        try:
            from app.services.provenance import (collect_conditions,
                                                 results_dir_for_rds,
                                                 write_export_record)
            conditions = collect_conditions(
                rds_path=rds_path, result_folder=result_folder,
                integration_method=current_method,
                extra={"exported_file": filename,
                       "report_top_n": top_n,
                       "report_methods": export_method_selection,
                       "report_include_deg": bool(include_deg)})
            write_export_record(results_dir_for_rds(rds_path, result_folder),
                                "pptx_report", conditions)
        except Exception as _e:
            logger.warning("PPTX の条件記録に失敗: %s", _e)

        # ★ ver51.9: 図に使う表示設定は「解析条件」スライドと **同じ dict** から
        #   取り出す。別々に読むと片方だけ直し忘れて、1 つの資料の中で
        #   条件と図が食い違う（従来がまさにそれで、Volcano の破線は常に
        #   既定の 0.5 / 1.3 なのに条件スライドには利用者の閾値が出ていた）。
        display_settings = _display_settings(conditions)

        # 全体 watchdog の期限（既定 45 分, PPTX_EXPORT_TIMEOUT_SEC で調整, 0 で無効）。
        # 1 枚単位のタイムアウトだけでは、多数がタイムアウトすると総時間が伸び得るため。
        import os as _os
        import time as _time
        _export_timeout = float(_os.environ.get("PPTX_EXPORT_TIMEOUT_SEC", "2700"))
        _deadline = (_time.monotonic() + _export_timeout) if _export_timeout > 0 else None

        # expression_matrix.parquet を必要時に on-demand 生成（feature plot / heatmap が利用）
        set_progress((1, 100, "発現データ準備中（初回は数十秒かかります）..."))
        try:
            if rds_path:
                _bridge.ensure_expression_matrix(rds_path)
        except Exception as e:
            logger.warning(f"発現データ準備失敗、feature レンダリングは fallback or スキップ: {e}")

        # ------------------------------------------------------------------
        # 出力対象手法リストの決定（export_method_selector に基づく）
        # ------------------------------------------------------------------
        selected = export_method_selection or []
        if isinstance(selected, str):  # 旧 str 値 ("all"/手法名) との後方互換
            selected = [] if selected in ("", "all") else [selected]

        has_methods = bool(rds_map and isinstance(rds_map, dict))
        if has_methods and not selected:
            return no_update, "出力する手法を1つ以上選択してください。"

        methods_to_export = []
        if has_methods:
            valid = [m for m in selected if m in rds_map]
            if not valid:
                return no_update, "選択した手法が見つかりません。"
            # current_method を先頭に（比較スライドの基準）
            if current_method and current_method in valid:
                methods_to_export = [current_method] + [
                    m for m in valid if m != current_method]
            else:
                methods_to_export = valid

        # ------------------------------------------------------------------
        # 単一手法の場合（従来の動作と完全互換）
        # ------------------------------------------------------------------
        if not methods_to_export:
            cache_dir_path = Path(cache_dir_str) if cache_dir_str else None
            df = None
            meta = {}
            if cache_dir_path:
                plot_parquet = cache_dir_path / "plot_data.parquet"
                plot_csv = cache_dir_path / "plot_data.csv"
                if plot_parquet.exists():
                    df = pd.read_parquet(plot_parquet)
                elif plot_csv.exists():
                    df = pd.read_csv(plot_csv)

                meta_file = cache_dir_path / "extraction_meta.json"
                if meta_file.exists():
                    with open(meta_file, "r", encoding="utf-8") as f:
                        meta = json.load(f)

            # ★ ver51.9 / B-7: 画面が「マージ統合」表示ならそれを反映する。
            df, custom_colors, _merged = _apply_merge_view(
                df, display_settings, custom_colors)

            pptx_bytes = _build_pptx(
                umap_fig, spatial_fig, meta, cluster_stats_data, rds_path,
                sub_name=sub_name, volcano_fig=volcano_fig,
                heatmap_fig=heatmap_fig,
                deg_data=deg_data, top_n=top_n, df=df,
                cache_dir=str(cache_dir_path) if cache_dir_path else None,
                custom_colors=custom_colors, rotation_store=rotation_store,
                name_map=name_map, set_progress=set_progress,
                mrm_path=mrm_path_str,
                saved_positions=saved_positions,
                cluster_name_map=cluster_name_map,
                include_deg=include_deg,
                deadline=_deadline,
                conditions=conditions,
                display_settings=display_settings,
            )

            return (
                dcc.send_bytes(pptx_bytes, filename=filename),
                f"✓ PPTXファイルを出力しました: {filename}",
            )

        # ------------------------------------------------------------------
        # 複数手法 or セレクタで指定された手法 → 1つの PPTX に結合
        # ------------------------------------------------------------------
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 解析条件は手法ごとではなく資料の先頭に 1 回だけ置く
        if conditions is not None:
            _add_conditions_slide(prs, conditions)

        is_multi = len(methods_to_export) > 1

        # 全手法のステップ数（Phase 1 後に再計算する）
        total_steps = len(methods_to_export) * 2  # 暫定: データ読込用
        progress_offset = 0
        exported_methods = []
        section_map = []  # (name, start_idx, end_idx) — ①-5

        # ==================================================================
        # Phase 1: 全手法のデータを事前抽出（キャッシュして二重抽出を防止）
        # ==================================================================
        extracted_data = {}  # method_name → dict
        skipped_methods = []  # 出力できなかった手法（最終ステータスで可視化）

        for method_name in methods_to_export:
            method_rds = rds_map.get(method_name)
            # 派生PCA（未補正）は専用RDSがディスク未生成のことがある。UI 読込
            # (load_stage_b_extract) と同じく Harmony から遅延生成してから出力する。
            if (method_name == "PCA" and method_rds
                    and not Path(method_rds).exists()):
                harmony_rds = rds_map.get("Harmony")
                if harmony_rds and Path(harmony_rds).exists():
                    try:
                        set_progress((
                            min(int(progress_offset / total_steps * 100), 99),
                            100,
                            "PCA(未補正)を生成中（初回は数分かかります）...",
                        ))
                        _bridge.derive_uncorrected_pca(harmony_rds, method_rds)
                    except Exception as e:
                        logger.warning(f"PCA(未補正)の派生生成に失敗: {e}")
            if not method_rds or not Path(method_rds).exists():
                logger.info(f"{method_name}: RDS ファイルが見つかりません → スキップ")
                skipped_methods.append(method_name)
                continue

            set_progress((
                min(int(progress_offset / total_steps * 100), 99), 100,
                f"{method_name} のデータを読み込み中..."
            ))

            try:
                result = _bridge.extract_data(method_rds)
            except Exception as e:
                logger.error(f"{method_name}: データ抽出エラー: {e}")
                skipped_methods.append(method_name)
                continue

            method_df = result["plot_data"]
            method_meta = result["meta"]
            method_cache_dir = result.get("cache_dir")

            # DEG 結果読み込み
            method_deg_data = None
            if result_folder:
                result_base = Path(result_folder)
                method_deg_data = _load_deg_results(
                    result_base, method_name)
            else:
                rds_dir = Path(method_rds).parent
                result_base = (rds_dir.parent
                               if rds_dir.name == "RDS_Files"
                               else rds_dir)
                method_deg_data = _load_deg_results(
                    result_base, method_name)

            # ★ ver51.9 / B-7: 画面が「マージ統合」表示ならそれを反映する。
            #   色も統合後の体系 (親クラスタ濃淡 / 独立色) に切り替わる。
            method_df, method_colors, _merged = _apply_merge_view(
                method_df, display_settings, custom_colors)

            extracted_data[method_name] = {
                "df": method_df,
                "custom_colors": method_colors,
                "merged": _merged,
                "meta": method_meta,
                "cache_dir": method_cache_dir,
                "deg_data": method_deg_data,
                "rds_path": method_rds,
                # ★ ver51.9: クラスタの改名は**手法ごとに独立**
                #   (`cluster_name_map::<method>`)。従来は `cluster_name_map_store`
                #   を State から 1 回読み、その **現在表示中の手法の分** を
                #   全手法のスライドへ渡していた。クラスタ ID は手法間で無関係なので、
                #   Harmony でクラスタ 3 を「腫瘍」にすると RPCA の別物のクラスタ 3 も
                #   「腫瘍」になる。手法比較のための資料でラベルが汚染される。
                #   ラベル位置は既に手法別に読み直している(下記 _method_positions)。
                #   表示中の手法だけ Store の値を使う (未保存の改名を落とさないため)。
                "cluster_name_map": (
                    (cluster_name_map or {})
                    if method_name == current_method
                    else _load_cluster_name_map(method_rds, method_name)),
            }
            progress_offset += 1

        if not extracted_data:
            return no_update, "エクスポート可能な手法がありません。"

        # Phase 1 完了: 実際のクラスタ数に基づいて total_steps を再計算
        total_steps = progress_offset                    # Phase 1 消費済み
        if is_multi:
            total_steps += len(extracted_data) * 2        # Phase 2: 比較 + sample UMAP
        for _ed in extracted_data.values():
            total_steps += 1                              # セパレータスライド
            _method_df = _ed["df"]
            if _method_df is not None:
                _n_clusters = len(_method_df["Cluster"].unique())
                total_steps += 3 + _n_clusters * 3       # _build_pptx 実ステップ (title, UMAP&Spatial, stats + 3/cluster)
            else:
                total_steps += 4

        # ==================================================================
        # Phase 2: 比較セクション — 全手法の UMAP & Spatial を先頭に配置
        #          (①-3: 複数手法の場合のみ)
        # ==================================================================
        if is_multi and len(extracted_data) > 1:
            comparison_start = len(prs.slides)

            for method_name, ed in extracted_data.items():
                m_df = ed["df"]
                if m_df is None or m_df.empty:
                    continue
                has_spatial_cmp = "SpatialX" in m_df.columns

                # 手法別ラベル位置を取得（現在のメソッドはメモリ蓄積をマージ）
                if method_name == current_method:
                    _method_positions = saved_positions
                else:
                    _method_positions = _load_label_positions_util(
                        ed["rds_path"], method_name)
                # ver51.9: 改名も色もマージ表示も手法別（Phase 1 で解決済み）
                _method_names_cmp = ed["cluster_name_map"]
                _method_colors_cmp = ed["custom_colors"]

                color_map_cmp = _get_cluster_color_map(
                    m_df["Cluster"], _method_colors_cmp)
                all_samples_cmp = sorted(m_df["Sample"].unique())
                n_sp_cmp = len(all_samples_cmp)
                _legend_w_cmp = 1.3
                avail_w_cmp = 13.33 - 0.3 - _legend_w_cmp - 0.1  # ≈ 11.93"

                # サンプル数が多い場合はスライドを分割
                _max_per_slide_cmp = max(1, int(avail_w_cmp / 1.5))
                if n_sp_cmp > _max_per_slide_cmp:
                    _mid_cmp = (n_sp_cmp + 1) // 2
                    _grp_cmp = [all_samples_cmp[:_mid_cmp], all_samples_cmp[_mid_cmp:]]
                else:
                    _grp_cmp = [all_samples_cmp]

                # UMAP axis ranges (全サンプル統一)
                u1_min = float(m_df["UMAP_1"].min())
                u1_max = float(m_df["UMAP_1"].max())
                u2_min = float(m_df["UMAP_2"].min())
                u2_max = float(m_df["UMAP_2"].max())
                u_pad = max(u1_max - u1_min,
                            u2_max - u2_min) * 0.05
                u_xrange = [u1_min - u_pad, u1_max + u_pad]
                u_yrange = [u2_min - u_pad, u2_max + u_pad]

                # レジェンド画像を事前生成
                _n_cl_cmp = len(m_df["Cluster"].unique())
                _lh_cmp = min(_n_cl_cmp * 0.35 + 0.2, 6.0)
                legend_fig_cmp = _build_cluster_legend_fig(
                    m_df["Cluster"].unique(), color_map_cmp,
                    cluster_name_map=_method_names_cmp)
                legend_png_cmp = _fig_to_png_bytes(
                    legend_fig_cmp.to_dict(),
                    width=200, height=600, scale=2)

                for _gi, _gs in enumerate(_grp_cmp):
                    _gs_n = len(_gs)
                    _gs_sfx = f" ({_gi + 1}/{len(_grp_cmp)})" if len(_grp_cmp) > 1 else ""
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    _pptx_add_title_bar(
                        slide,
                        f"UMAP & Spatial Mapping \u2014 {method_name}{_gs_sfx}")

                    # 上段: サンプル別 UMAP
                    tile_w_cmp = avail_w_cmp / max(_gs_n, 1)
                    _umap_pos_cmp = (_method_positions or {}).get("umap_integrated", {})
                    for idx_s, s in enumerate(_gs):
                        df_s = m_df[m_df["Sample"] == s]
                        umap_s = _build_umap_integrated_fig(
                            df_s, color_by="Cluster",
                            highlight_clusters=None,
                            show_legend=False, show_labels=True,
                            title=_display_name(s, name_map),
                            marker_size=3, custom_colors=_method_colors_cmp,
                            title_font_size=40, label_size=24,
                            saved_positions=_umap_pos_cmp,
                            cluster_name_map=_method_names_cmp)
                        if umap_s is not None:
                            umap_s.update_xaxes(range=u_xrange)
                            umap_s.update_yaxes(range=u_yrange)
                            u_dict = (umap_s.to_dict()
                                      if hasattr(umap_s, "to_dict")
                                      else umap_s)
                            u_png = _fig_to_png_bytes(
                                u_dict, width=600, height=600, scale=2)
                            _cw, _ch, _coff = _square_tile_dims(
                                tile_w_cmp, 3.0)
                            u_left = Inches(
                                0.3 + idx_s * tile_w_cmp + _coff)
                            _pptx_add_image(slide, u_png,
                                            int(u_left), Inches(0.9),
                                            Inches(_cw), Inches(_ch))

                    # 下段: サンプル別 Spatial
                    if has_spatial_cmp:
                        _rot_store = rotation_store or {}
                        for idx_s, s in enumerate(_gs):
                            df_s = m_df[m_df["Sample"] == s]
                            transform = _rot_store.get(
                                s, _rot_store.get(
                                    "__all__",
                                    {"angle": 0, "flip_h": False,
                                     "flip_v": False}))
                            if isinstance(transform, (int, float)):
                                transform = {
                                    "angle": int(transform),
                                    "flip_h": False, "flip_v": False}
                            _sp_pos_cmp = (_method_positions or {}).get("spatial", {}).get(s, {})
                            sp_fig_cmp = _create_single_spatial_fig(
                                df_s, color_map_cmp,
                                highlight_clusters=None,
                                selected_cell_ids=set(),
                                rotation_deg=transform.get("angle", 0),
                                show_labels=True,
                                flip_h=transform.get("flip_h", False),
                                flip_v=transform.get("flip_v", False),
                                title=_display_name(s, name_map),
                                marker_size=0, render_height=560,
                                embed_legend=False,
                                title_font_size=40, label_size=24,
                                saved_positions=_sp_pos_cmp,
                                cluster_name_map=_method_names_cmp)
                            if sp_fig_cmp is not None:
                                sp_dict = (sp_fig_cmp.to_dict()
                                           if hasattr(sp_fig_cmp, "to_dict")
                                           else sp_fig_cmp)
                                sp_png = _fig_to_png_bytes(
                                    sp_dict, width=600, height=600,
                                    scale=2)
                                _csw, _csh, _csoff = _square_tile_dims(
                                    tile_w_cmp, 3.1)
                                sp_left = Inches(
                                    0.3 + idx_s * tile_w_cmp + _csoff)
                                _pptx_add_image(
                                    slide, sp_png,
                                    int(sp_left), Inches(4.1),
                                    Inches(_csw), Inches(_csh))

                    # クラスタレジェンド（右下角に配置）
                    _lx_cmp = Inches(0.3 + avail_w_cmp + 0.1)
                    _ly_cmp = Inches(7.5 - _lh_cmp)
                    _pptx_add_image_preserve_ratio(slide, legend_png_cmp,
                                                   int(_lx_cmp), int(_ly_cmp),
                                                   Inches(_legend_w_cmp), Inches(_lh_cmp),
                                                   png_w=200, png_h=600)

                # --- UMAP (by Sample) comparison slide --- 削除済み（不要）

                progress_offset += 1

            comparison_end = len(prs.slides) - 1
            if comparison_end >= comparison_start:
                section_map.append(
                    ("Comparison", comparison_start, comparison_end))

        # ==================================================================
        # Phase 3: 各手法のフルセクション
        # ==================================================================
        for method_name in methods_to_export:
            if method_name not in extracted_data:
                continue

            ed = extracted_data[method_name]
            method_df = ed["df"]
            method_meta = ed["meta"]
            method_cache_dir = ed["cache_dir"]
            method_deg_data = ed["deg_data"]
            method_rds = ed["rds_path"]
            # ver51.9: 改名・色・マージ表示は手法別（Phase 1 で解決済み）
            method_name_map = ed["cluster_name_map"]
            method_colors = ed["custom_colors"]
            method_merged = ed["merged"]

            method_start_idx = len(prs.slides)

            # 手法別ラベル位置を取得
            if method_name == current_method:
                method_saved_positions = saved_positions
            else:
                method_saved_positions = _load_label_positions_util(
                    method_rds, method_name)

            set_progress((
                min(int(progress_offset / total_steps * 100), 99),
                100,
                f"{method_name} のスライドを生成中..."
            ))

            # --- セパレータスライド ---
            sep_slide = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_add_title_bar(sep_slide, f"═══ {method_name} ═══")
            txBox = sep_slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(11), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"Integration Method: {method_name}"
            p.font.size = Pt(28)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            p2 = tf.add_paragraph()
            p2.text = (
                f"Cells: {method_meta.get('n_cells', '?')} | "
                f"Clusters: {method_meta.get('n_clusters', '?')}"
            )
            p2.font.size = Pt(18)
            p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            p2.alignment = PP_ALIGN.CENTER
            if method_merged:
                # ★ ver51.9 / B-7: どちらのクラスタリングを出したのかを明記する。
                #   DEG は元クラスタで計算されているため、マージ後のクラスタには
                #   対応する Volcano / Heatmap が存在しない（スライドが出ない）。
                #   黙って欠けるのではなく理由を書く。
                p3 = tf.add_paragraph()
                p3.text = ("クラスタ表示: マージ統合 "
                           "（DEG は元クラスタで算出のため Volcano / Heatmap は非掲載）")
                p3.font.size = Pt(14)
                p3.font.color.rgb = RGBColor(0x85, 0x64, 0x04)
                p3.alignment = PP_ALIGN.CENTER
            progress_offset += 1

            # --- UMAP 図を生成 ---
            _umap_pos_m = (method_saved_positions or {}).get("umap_integrated", {})
            method_umap_fig = _build_umap_integrated_fig(
                method_df, color_by="Cluster",
                highlight_clusters=None,
                show_legend=True, show_labels=True,
                custom_colors=method_colors,
                saved_positions=_umap_pos_m,
                cluster_name_map=method_name_map,
            )

            # --- クラスタ統計の生成 ---
            method_cluster_stats = []
            try:
                clusters_sorted = sorted(
                    method_df["Cluster"].unique(),
                    key=_cluster_sort_key)
                n_total = len(method_df)
                for c in clusters_sorted:
                    n_c = int((method_df["Cluster"] == c).sum())
                    pct = (f"{n_c / n_total * 100:.1f}"
                           if n_total else "0")
                    # ★ ver51.9 / B-6: 表は生のクラスタ ID、各スライドの見出しは
                    #   改名後の表示名、という不整合があった。対応表が無いので
                    #   「腫瘍」の統計がどの行か分からない。見出しと同じ関数
                    #   (_cluster_display_name) で揃える。
                    #   改名していないクラスタは従来どおり ID がそのまま出る。
                    method_cluster_stats.append({
                        "Cluster": _cluster_display_name(str(c), method_name_map),
                        "Pixels": n_c,
                        "Percent": pct,
                    })
            except Exception:
                method_cluster_stats = []

            # --- _build_pptx でフルセットを追加 ---
            method_sub_name = (
                f"{sub_name} [{method_name}]"
                if sub_name else method_name
            )
            returned = _build_pptx(
                method_umap_fig, None, method_meta,
                method_cluster_stats, method_rds,
                sub_name=method_sub_name,
                volcano_fig=None, heatmap_fig=None,
                deg_data=method_deg_data, top_n=top_n,
                df=method_df,
                cache_dir=(str(method_cache_dir)
                           if method_cache_dir else None),
                custom_colors=method_colors,
                rotation_store=rotation_store,
                name_map=name_map,
                set_progress=set_progress,
                mrm_path=mrm_path_str,
                existing_prs=prs,
                progress_offset=progress_offset,
                progress_total=total_steps,
                saved_positions=method_saved_positions,
                cluster_name_map=method_name_map,
                include_deg=include_deg,
                deadline=_deadline,
                display_settings=display_settings,
            )
            if isinstance(returned, int):
                progress_offset = returned

            method_end_idx = len(prs.slides) - 1
            section_map.append(
                (method_name, method_start_idx, method_end_idx))
            exported_methods.append(method_name)

        if not exported_methods:
            return no_update, "エクスポート可能な手法がありません。"

        # ==================================================================
        # Phase 4: セクション情報を PPTX XML に追加（①-5）
        # ==================================================================
        if section_map:
            try:
                _pptx_add_sections(prs, section_map)
            except Exception as e:
                logger.warning(f"PPTX セクション追加エラー（無視して継続）: {e}")

        # --- 結合 PPTX をバイト列に変換 ---
        output = BytesIO()
        prs.save(output)
        output.seek(0)

        methods_str = " + ".join(exported_methods)
        status_msg = f"✓ PPTXファイルを出力しました ({methods_str}): {filename}"
        if skipped_methods:
            status_msg += (
                f"（スキップ: {', '.join(dict.fromkeys(skipped_methods))}"
                f" — RDS が見つからない/抽出失敗）"
            )
        # ★ ver51.8: 図の PNG 変換に失敗した枚数を必ず伝える。
        #   従来は失敗が無言で握りつぶされ、**図の無いスライドがそのまま出力**
        #   されていた（ログにも残らなかった）。kaleido の不整合などで全滅しても
        #   「✓ 出力しました」としか表示されない状態だった。
        try:
            from app.utils.pptx_helpers import render_failure_count
            _n_failed = render_failure_count()
            if _n_failed:
                status_msg += (
                    f"　⚠ {_n_failed} 枚の図を描画できずスライドから省きました"
                    "（ログに理由を記録しています）"
                )
        except Exception:
            pass
        return (
            dcc.send_bytes(output.getvalue(), filename=filename),
            status_msg,
        )

    except Exception as e:
        return no_update, f"エクスポートエラー: {e}"
    finally:
        # 描画に使った共有プールと配下の Chromium を必ず後始末（PIDS 増殖を防ぐ）
        try:
            _shutdown_render_pool()
        except Exception:
            pass
