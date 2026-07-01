"""PPTX エクスポート用ヘルパー関数群。

interactive_callbacks.py から抽出した PPTX 生成ユーティリティ。
"""

import copy
import logging
import os
import threading
import uuid
from concurrent.futures import Future, ProcessPoolExecutor
from concurrent.futures import TimeoutError as FuturesTimeout
from io import BytesIO
from typing import Optional

import plotly.graph_objects as go
import plotly.io as pio
from lxml import etree
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.utils.color_utils import cluster_display_name, cluster_sort_key

logger = logging.getLogger(__name__)

# 1 枚あたりの描画タイムアウト（秒）。kaleido 0.2.1 は Docker/ヘッドレスで to_image() が
# 無言ハングし得るため、必ず打ち切れるようにする。
_RENDER_TIMEOUT_SEC = float(os.environ.get("PPTX_RENDER_TIMEOUT_SEC", "60"))
# 0 以下で無効。>0 の場合、この点数を超える散布図は SVG 化の過負荷を防ぐため均等間引きする。
_MAX_SCATTER_POINTS = int(os.environ.get("PPTX_MAX_SCATTER_POINTS", "0"))


# ---------------------------------------------------------------------------
# PNG conversion
# ---------------------------------------------------------------------------

def _maybe_downsample_scatter(tr):
    """散布図 trace の点数が _MAX_SCATTER_POINTS を超える場合、均等間引きする（in-place）。

    x / y と、配列で与えられた marker.color/size・text・customdata を同じ間隔で間引き整合を保つ。
    既定 (_MAX_SCATTER_POINTS<=0) では何もしない。失敗しても描画は継続。
    """
    if _MAX_SCATTER_POINTS <= 0:
        return
    try:
        x = tr.get("x")
        if x is None:
            return
        n = len(x)
        if n <= _MAX_SCATTER_POINTS:
            return
        step = (n // _MAX_SCATTER_POINTS) + 1

        def _stride(seq):
            if isinstance(seq, (list, tuple)) and len(seq) == n:
                return seq[::step]
            return seq

        tr["x"] = _stride(tr.get("x"))
        tr["y"] = _stride(tr.get("y"))
        if "text" in tr:
            tr["text"] = _stride(tr.get("text"))
        if "customdata" in tr:
            tr["customdata"] = _stride(tr.get("customdata"))
        marker = tr.get("marker")
        if isinstance(marker, dict):
            if isinstance(marker.get("color"), (list, tuple)):
                marker["color"] = _stride(marker.get("color"))
            if isinstance(marker.get("size"), (list, tuple)):
                marker["size"] = _stride(marker.get("size"))
    except Exception:
        pass


def _sanitize_fig_dict_for_export(fig_dict):
    """kaleido/ヘッドレス Chromium 向けに図 dict を安全化する（in-place）。

    - WebGL(scattergl) → SVG(scatter): GPU 無しコンテナの SoftGL(SwiftShader) では
      scattergl の静的描画が無言ハングし得るため、SVG に変換する（ハングの直接対処）。
    - 過大な点数の散布図はダウンサンプル（既定は無効）。
    """
    if not isinstance(fig_dict, dict):
        return fig_dict
    data = fig_dict.get("data")
    if not isinstance(data, list):
        return fig_dict
    for tr in data:
        if not isinstance(tr, dict):
            continue
        if tr.get("type") == "scattergl":
            tr["type"] = "scatter"
            _maybe_downsample_scatter(tr)
    return fig_dict


def fig_to_png_bytes(fig_dict, width=1200, height=800, scale=2):
    """Plotly figure (dict または go.Figure) を PNG バイト列に変換する。

    - kaleido 未インストール時や描画失敗時は None を返す。
    - WebGL(scattergl) を SVG(scatter) に変換し、ヘッドレス Chromium のハングを回避する。
    - 呼び出し側の dict を壊さないよう deepcopy してから加工する。
    """
    try:
        if hasattr(fig_dict, "to_dict"):
            d = fig_dict.to_dict()
        else:
            d = copy.deepcopy(fig_dict)
        _sanitize_fig_dict_for_export(d)
        fig = go.Figure(d)
        fig.update_layout(paper_bgcolor="white", plot_bgcolor="white")
        return pio.to_image(fig, format="png", width=width, height=height, scale=scale)
    except Exception:
        return None


# ---------------------------------------------------------------------------
# タイムアウト付き描画（ハング根絶 + プロセスリーク防止）
# ---------------------------------------------------------------------------
# kaleido の描画は必ず worker プロセスで実行し fut.result(timeout) で打ち切る。固まった
# worker とその配下の Chromium はプロセスツリーごと kill してリーク(PIDS 増殖)を防ぐ。

_shared_lock = threading.Lock()
_shared_pool = None  # ProcessPoolExecutor(max_workers=1) — 単発描画の共有プール


def _kill_process_trees(procs, grace=3.0):
    """worker プロセスとその子孫(kaleido/Chromium)を確実に終了させる。"""
    if not procs:
        return
    try:
        import psutil
    except Exception:
        for p in procs:
            try:
                p.terminate()
            except Exception:
                pass
        return
    targets = []
    for p in procs:
        pid = getattr(p, "pid", None)
        if not pid:
            continue
        try:
            parent = psutil.Process(pid)
        except Exception:
            continue
        try:
            targets.extend(parent.children(recursive=True))
        except Exception:
            pass
        targets.append(parent)
    for t in targets:
        try:
            t.terminate()
        except Exception:
            pass
    try:
        _gone, alive = psutil.wait_procs(targets, timeout=grace)
    except Exception:
        alive = targets
    for t in alive:
        try:
            t.kill()
        except Exception:
            pass


def _shutdown_pool_hard(pool):
    """pool の worker ツリーを kill してから shutdown する（wait でブロックしない）。"""
    if pool is None:
        return
    try:
        procs = list(getattr(pool, "_processes", {}).values())
    except Exception:
        procs = []
    try:
        pool.shutdown(wait=False, cancel_futures=True)
    except TypeError:  # Python < 3.9 は cancel_futures 無し
        try:
            pool.shutdown(wait=False)
        except Exception:
            pass
    except Exception:
        pass
    _kill_process_trees(procs)


def _recycle_shared_pool():
    global _shared_pool
    with _shared_lock:
        pool = _shared_pool
        _shared_pool = None
    _shutdown_pool_hard(pool)


def _kill_lingering_kaleido():
    """残存する kaleido/Chromium プロセスを一掃するバックストップ。

    ProcessPool worker の子として掴めず（別プロセスグループ化/再親化により）残った kaleido を
    確実に回収する。PPTX 出力は事実上直列で、対話 UI 側は kaleido を使わないため、
    出力終了時の一括 kill は安全。
    """
    try:
        import psutil
    except Exception:
        return
    for p in psutil.process_iter(["name", "cmdline"]):
        try:
            name = (p.info.get("name") or "").lower()
            cmd = " ".join(p.info.get("cmdline") or []).lower()
            if "kaleido" in name or "kaleido" in cmd:
                p.kill()
        except Exception:
            pass


def shutdown_shared_queue():
    """共有描画プールを破棄する（エクスポート終了時に必ず呼ぶ）。"""
    _recycle_shared_pool()
    _kill_lingering_kaleido()


def render_png(fig_dict, width=1200, height=800, scale=2, timeout=None):
    """1 枚をタイムアウト付きで描画する。タイムアウト/失敗時は None を返す（スキップ可能）。

    固まった場合は worker+Chromium を kill してプールを作り直す。
    """
    global _shared_pool
    t = timeout if timeout is not None else _RENDER_TIMEOUT_SEC
    with _shared_lock:
        if _shared_pool is None:
            _shared_pool = ProcessPoolExecutor(max_workers=1)
        pool = _shared_pool
    try:
        fut = pool.submit(fig_to_png_bytes, fig_dict, width, height, scale)
    except Exception as e:
        logger.warning("[PPTX] 描画プールへの投入に失敗、作り直します: %s", e)
        _recycle_shared_pool()
        return None
    try:
        return fut.result(timeout=t)
    except FuturesTimeout:
        logger.warning(
            "[PPTX] 画像描画が %.0fs を超過。この図をスキップし描画プロセスを回収します。", t)
        _recycle_shared_pool()
        return None
    except Exception as e:
        logger.warning("[PPTX] 画像描画に失敗、この図をスキップ: %s", e)
        return None


class RenderQueue:
    """kaleido PNG レンダリングを ProcessPool で並列化するキュー。

    kaleido は Chromium プロセスを介したレンダリング (1 回 200-500ms) で、
    内部状態がプロセス単位のため ThreadPool では並列化できない。
    ProcessPoolExecutor で各 worker が独立した Chromium を持つことで真の並列化を実現する。

    使い方:
        with RenderQueue(max_workers=4) as q:
            fut1 = q.submit(fig_dict1, 800, 600, 2)
            fut2 = q.submit(fig_dict2, 1200, 800, 2)
            # ... 他の処理を進める間にバックグラウンドでレンダリング
            png1 = fut1.result()  # ここで block (すでに並列で進んでいる)
            png2 = fut2.result()

    注意:
        - 各 worker は独立した kaleido / Chromium を起動するため初期化 1-2 秒のオーバーヘッド
          (worker あたり 50-100MB の RAM を消費)
        - fig_dict は pickle 可能な dict であること (plotly Figure を .to_dict() したもの)
        - context manager 外で submit() を呼ぶと RuntimeError
    """

    def __init__(self, max_workers: int = 4):
        self._pool: Optional[ProcessPoolExecutor] = None
        self._max_workers = max_workers

    def __enter__(self) -> "RenderQueue":
        self._pool = ProcessPoolExecutor(max_workers=self._max_workers)
        return self

    def __exit__(self, exc_type, exc_val, exc_tb) -> None:
        # wait=True は固まった worker を待って二次ブロックするため使わない。
        # worker ツリー(kaleido/Chromium)を kill してから shutdown し、リークを防ぐ。
        pool = self._pool
        self._pool = None
        _shutdown_pool_hard(pool)

    def submit(self, fig_dict, width: int = 1200, height: int = 800,
               scale: int = 2) -> Future:
        """fig_dict のレンダリングをバックグラウンドで開始。Future を返す。"""
        if self._pool is None:
            raise RuntimeError(
                "RenderQueue は context manager (with) として使ってください"
            )
        return self._pool.submit(fig_to_png_bytes, fig_dict, width, height, scale)

    def result(self, fut, timeout=None):
        """Future をタイムアウト付きで取得。タイムアウト/失敗時は None（スキップ）。"""
        if fut is None:
            return None
        t = timeout if timeout is not None else _RENDER_TIMEOUT_SEC
        try:
            return fut.result(timeout=t)
        except FuturesTimeout:
            logger.warning("[PPTX] 画像描画が %.0fs を超過。この図をスキップします。", t)
            try:
                fut.cancel()
            except Exception:
                pass
            return None
        except Exception as e:
            logger.warning("[PPTX] 画像描画に失敗、この図をスキップ: %s", e)
            return None

    def render(self, fig_dict, width: int = 1200, height: int = 800,
               scale: int = 2, timeout=None):
        """submit + タイムアウト付き result のワンショット。"""
        return self.result(self.submit(fig_dict, width, height, scale), timeout=timeout)


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def pptx_add_title_bar(slide, title_text):
    """PPTX スライドにタイトルバーを追加する。"""
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(8), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def pptx_add_image(slide, png_bytes, left, top, width, height):
    """PNG バイト列を PPTX スライドの指定位置に配置する。"""
    if png_bytes:
        img_stream = BytesIO(png_bytes)
        slide.shapes.add_picture(img_stream, left, top, width, height)


def pptx_add_image_preserve_ratio(slide, png_bytes, left, top,
                                   max_width, max_height,
                                   png_w=None, png_h=None):
    """PNG 元比率を保持して max_width × max_height 内に収まるよう配置する。

    png_w, png_h: PNG 生成時のピクセル寸法（比率計算に使用）。
    未指定時は max_width × max_height をそのまま使用（従来互換）。
    """
    if not png_bytes:
        return
    if png_w and png_h:
        aspect = png_w / png_h
        max_aspect = max_width / max_height if max_height else 1.0
        if max_aspect > aspect:
            # 高さ制約: height = max_height, width = height * aspect
            h = max_height
            w = int(h * aspect)
        else:
            # 幅制約: width = max_width, height = width / aspect
            w = max_width
            h = int(w / aspect)
        # 枠内中央寄せ
        left = left + int((max_width - w) / 2)
        top = top + int((max_height - h) / 2)
    else:
        w, h = max_width, max_height
    img_stream = BytesIO(png_bytes)
    slide.shapes.add_picture(img_stream, left, top, w, h)


# ---------------------------------------------------------------------------
# Layout helpers
# ---------------------------------------------------------------------------

def square_tile_dims(tile_w, row_h, margin=0.95):
    """正方形画像のタイル配置寸法を計算する。

    Returns:
        (side, side, offset): 正方形の辺長と、タイル内の左オフセット
    """
    avail_w = tile_w * margin
    side = min(avail_w, row_h)
    offset = (tile_w - side) / 2
    return side, side, offset


# ---------------------------------------------------------------------------
# Legend figures
# ---------------------------------------------------------------------------

def build_cluster_legend_fig(cluster_list, color_map, font_size=18,
                             marker_size=14, cluster_name_map=None):
    """クラスタレジェンド専用のPlotly Figureを生成"""
    fig = go.Figure()
    for cl in sorted(cluster_list, key=cluster_sort_key):
        fig.add_trace(go.Scatter(
            x=[None], y=[None], mode="markers",
            marker=dict(size=marker_size,
                        color=color_map.get(str(cl), "#999")),
            name=cluster_display_name(cl, cluster_name_map), showlegend=True,
        ))
    fig.update_layout(
        showlegend=True,
        legend=dict(
            font=dict(size=font_size),
            itemsizing="constant",
            yanchor="middle", y=0.5,
            xanchor="center", x=0.5,
            tracegroupgap=4,
        ),
        margin=dict(l=0, r=0, t=0, b=0),
        xaxis=dict(visible=False), yaxis=dict(visible=False),
        plot_bgcolor="white", paper_bgcolor="white",
        width=200, height=600,
    )
    return fig


def build_sample_legend_fig(sample_list, sample_color_map, name_map=None,
                             font_size=18, marker_size=14):
    """サンプルレジェンド専用のPlotly Figureを生成"""
    fig = go.Figure()
    for s in sorted(sample_list):
        display = (name_map or {}).get(s, s) if name_map else s
        fig.add_trace(go.Scatter(
            x=[None], y=[None], mode="markers",
            marker=dict(size=marker_size,
                        color=sample_color_map.get(str(s), "#999")),
            name=str(display), showlegend=True,
        ))
    fig.update_layout(
        showlegend=True,
        legend=dict(
            font=dict(size=font_size),
            itemsizing="constant",
            yanchor="middle", y=0.5,
            xanchor="center", x=0.5,
            tracegroupgap=4,
        ),
        margin=dict(l=0, r=0, t=0, b=0),
        xaxis=dict(visible=False), yaxis=dict(visible=False),
        plot_bgcolor="white", paper_bgcolor="white",
        width=200, height=600,
    )
    return fig


# ---------------------------------------------------------------------------
# Presentation sections
# ---------------------------------------------------------------------------

def pptx_add_sections(prs, section_map):
    """PPTX にセクション情報を追加する（PowerPoint のセクションパネルに表示される）。

    Args:
        prs: python-pptx Presentation オブジェクト
        section_map: list of (section_name, start_slide_idx, end_slide_idx)
            slide indices は 0-based、両端を含む。
    """
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    p14_ns = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    ext_uri = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"

    prs_elem = prs.element  # CT_Presentation lxml element

    # p14 名前空間プレフィックスを登録（PowerPoint互換性のため）
    etree.register_namespace("p14", p14_ns)

    # sldIdLst から全スライド ID を取得
    sldIdLst = prs_elem.find(f"{{{p_ns}}}sldIdLst")
    if sldIdLst is None:
        return

    sld_ids = []
    for sldId in sldIdLst:
        sid = sldId.get("id")
        if sid:
            sld_ids.append(sid)
    if not sld_ids:
        return

    # extLst を検索または作成
    extLst = prs_elem.find(f"{{{p_ns}}}extLst")
    if extLst is None:
        extLst = etree.SubElement(prs_elem, f"{{{p_ns}}}extLst")

    # ext 要素を作成（セクション用 URI）
    ext = etree.SubElement(extLst, f"{{{p_ns}}}ext")
    ext.set("uri", ext_uri)

    # sectionLst を作成
    sectionLst = etree.SubElement(ext, f"{{{p14_ns}}}sectionLst")

    for sec_name, start_idx, end_idx in section_map:
        section = etree.SubElement(sectionLst, f"{{{p14_ns}}}section")
        section.set("name", sec_name)
        section.set("id", "{" + str(uuid.uuid4()).upper() + "}")

        sldIdLst_sec = etree.SubElement(section, f"{{{p14_ns}}}sldIdLst")

        for i in range(start_idx, min(end_idx + 1, len(sld_ids))):
            sldId_ref = etree.SubElement(sldIdLst_sec, f"{{{p14_ns}}}sldId")
            sldId_ref.set("id", sld_ids[i])
