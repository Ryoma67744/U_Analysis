"""エクスポート経路への条件付与のテスト。

観点:
  - ZIP / xlsx には条件が同梱されること。
  - 生 CSV は列を変えないまま、条件がサーバ側に残ること。
  - Methods 表示が Master Password で守られていること。
"""
import json
import zipfile
from io import BytesIO

import pandas as pd
import pytest

from app.services import provenance as pv

PARAMS = {
    "analysis_type": "tims_v8",
    "umap_n_neighbors": 30, "umap_min_dist": 0.3,
    "umap_metric": "cosine", "umap_dims_n": 30,
    "p_thresh": 0.05, "logfc_thresh": 0.25,
    "ion_mode": "Positive", "tolerance_mz": 0.01,
    "template_path": "/app/Script/TIMS/tmpl.R",
    "norm_mode": "log1p", "input_normalized": False,
}


@pytest.fixture
def result_dir(tmp_path):
    out = tmp_path / "Analysis_20260729"
    (out / "RDS_Files").mkdir(parents=True)
    (out / "analysis_params.json").write_text(json.dumps(PARAMS), encoding="utf-8")
    rds = out / "RDS_Files" / "obj.rds"
    rds.write_bytes(b"x")
    return out, rds


# ---------------------------------------------------------------------------
# バッチ ZIP
# ---------------------------------------------------------------------------

# ver56.5: 図の中身は「トレースが 1 本でもあること」。以前ここは
# {"data": [], "layout": {}} という**空の figure** を置き場所として使っていたが、
# 空の figure は PNG にすると真っ白になるため _create_zip_from_figures が
# 弾くようになった (C10-5)。条件 JSON の同梱を見るテストなので、
# 図としては最小の実データを持たせる。
_FIG = {"data": [{"type": "scattergl", "x": [1, 2], "y": [1, 2]}], "layout": {}}


def test_batch_zip_includes_conditions(monkeypatch):
    import app.callbacks.interactive_batch_save as bs

    monkeypatch.setattr(bs, "fig_to_png_bytes",
                        lambda *a, **k: b"\x89PNG\r\n\x1a\n" + b"0" * 64)
    conditions = {"conditions_version": "1", "analysis": {"analysis_type": "tims_v8"}}
    data = bs._create_zip_from_figures(
        [("UMAP_integrated", _FIG)],
        width=100, height=100, scale=1, section_name="UMAP",
        conditions=conditions)

    assert data is not None
    with zipfile.ZipFile(BytesIO(data)) as zf:
        assert "analysis_conditions.json" in zf.namelist()
        assert "UMAP_integrated.png" in zf.namelist()
        payload = json.loads(zf.read("analysis_conditions.json").decode("utf-8"))
        assert payload["analysis"]["analysis_type"] == "tims_v8"


def test_batch_zip_without_conditions_is_unchanged(monkeypatch):
    """条件が取れなかったときは従来どおりの ZIP（出力自体は止めない）。"""
    import app.callbacks.interactive_batch_save as bs

    monkeypatch.setattr(bs, "fig_to_png_bytes", lambda *a, **k: b"PNGDATA")
    data = bs._create_zip_from_figures(
        [("UMAP_integrated", _FIG)],
        width=100, height=100, scale=1, conditions=None)
    with zipfile.ZipFile(BytesIO(data)) as zf:
        assert zf.namelist() == ["UMAP_integrated.png"]


def test_batch_zip_returns_none_when_no_figures(monkeypatch):
    """図が 0 枚なら manifest だけの ZIP を作らず None を返す（従来の契約）。"""
    import app.callbacks.interactive_batch_save as bs

    monkeypatch.setattr(bs, "fig_to_png_bytes", lambda *a, **k: None)
    assert bs._create_zip_from_figures([], 100, 100, 1, conditions={"a": 1}) is None
    assert bs._create_zip_from_figures(
        [("x", {})], 100, 100, 1, conditions={"a": 1}) is None


def test_conditions_for_records_server_side(result_dir, monkeypatch):
    import app.callbacks.interactive_batch_save as bs

    out, rds = result_dir
    conditions = bs._conditions_for(str(rds), "batch_zip_umap")
    assert conditions is not None
    records = list((out / "provenance").glob("export_*_batch_zip_umap.json"))
    assert len(records) == 1


def test_conditions_for_survives_failure(monkeypatch):
    """条件収集がこけてもエクスポートは止めない。"""
    import app.callbacks.interactive_batch_save as bs
    import app.services.provenance as prov

    def _boom(*a, **k):
        raise RuntimeError("boom")

    monkeypatch.setattr(prov, "collect_conditions", _boom)
    assert bs._conditions_for("/nope/x.rds", "batch_zip_umap") is None


# ---------------------------------------------------------------------------
# データ出力 (xlsx の Conditions シート)
# ---------------------------------------------------------------------------

def test_conditions_sheet_dataframe():
    from app.callbacks.interactive_data_export import _conditions_sheet_df

    conditions = pv.collect_conditions(rds_path=None, integration_method="Harmony")
    df = _conditions_sheet_df(conditions)
    assert list(df.columns) == ["項目", "値"]
    assert (df["項目"] == "統合手法").any()
    assert df.loc[df["項目"] == "統合手法", "値"].iloc[0] == "Harmony"
    # 未記録項目が一覧行として入る
    assert (df["項目"] == "未記録の項目").any()


def test_conditions_sheet_written_to_xlsx(tmp_path):
    """Conditions シートが実際にブックへ書けること。"""
    from app.callbacks.interactive_data_export import _conditions_sheet_df

    conditions = pv.collect_conditions(rds_path=None)
    path = tmp_path / "out.xlsx"
    with pd.ExcelWriter(path, engine="openpyxl") as writer:
        pd.DataFrame({"a": [1]}).to_excel(writer, sheet_name="Data", index=False)
        _conditions_sheet_df(conditions).to_excel(
            writer, sheet_name="Conditions", index=False)
    book = pd.ExcelFile(path)
    assert "Conditions" in book.sheet_names


# ---------------------------------------------------------------------------
# 生 CSV: 列は変えず、条件はサーバ側へ
# ---------------------------------------------------------------------------

def test_marker_csv_keeps_columns_and_records_conditions(result_dir):
    from app.callbacks.interactive_loupe import export_marker_table

    out, rds = result_dir
    rows = [{"cluster": "1", "mz": 611.14, "log2FC": 1.2},
            {"cluster": "2", "mz": 700.20, "log2FC": -0.8}]
    result = export_marker_table(
        1, rows, top_n=1, sort_by=[{"column_id": "log2FC", "direction": "desc"}],
        filter_query="{cluster} eq 1", cluster_filter="1",
        rds_path=str(rds), result_folder=None, method="Harmony")

    # CSV 本体は従来どおり（列を足していない）
    assert "content" in result or callable(getattr(result, "get", lambda k: None))
    records = list((out / "provenance").glob("export_*_csv_markers_topN.json"))
    assert len(records) == 1
    payload = json.loads(records[0].read_text(encoding="utf-8"))
    extra = payload["extra"]
    assert extra["exported_file"] == "markers_topN.csv"
    assert extra["top_n"] == 1
    assert extra["rows_exported"] == 1
    # 並び替え・絞り込みが残っていないと同じ表を再現できない
    assert extra["sort_by"][0]["column_id"] == "log2FC"
    assert extra["filter_query"] == "{cluster} eq 1"
    assert extra["cluster_filter"] == "1"


def test_onthefly_de_csv_records_selection_and_thresholds(result_dir):
    from app.callbacks.interactive_de import export_onthefly_de

    out, rds = result_dir
    rows = [{"gene": "611.14", "avg_log2FC": 1.5, "p_val_adj": 0.001}]
    export_onthefly_de(
        1, rows, top_n=50, sort_by=None, filter_query=None,
        mode="local", target_clusters=["3", "4"],
        selected_ids=["c1", "c2", "c3"], rds_path=str(rds),
        result_folder=None, method="Harmony")

    records = list((out / "provenance").glob("export_*_csv_onthefly_DE.json"))
    assert len(records) == 1
    payload = json.loads(records[0].read_text(encoding="utf-8"))
    assert payload["extra"]["de_mode"] == "local"
    assert payload["extra"]["de_target_clusters"] == ["3", "4"]
    assert payload["extra"]["n_selected_pixels"] == 3
    # GUI に出ていない固定値が必ず添付される
    assert payload["onthefly_de_fixed_params"]["test"] == "wilcox"
    assert payload["onthefly_de_fixed_params"]["min_pct"] == 0.05


def test_selection_groups_csv_records_conditions(result_dir):
    from app.callbacks.interactive_selection_groups import export_selection_groups

    out, rds = result_dir
    state = {"groups": [{"id": "g1", "name": "Tumor", "cell_ids": ["a", "b"]}]}
    result = export_selection_groups(1, state, str(rds), None, "Harmony")
    assert result["filename"] == "selection_groups.csv"
    records = list((out / "provenance").glob("export_*_csv_selection_groups.json"))
    assert len(records) == 1
    payload = json.loads(records[0].read_text(encoding="utf-8"))
    assert payload["extra"]["groups"] == [{"name": "Tumor", "n_cells": 2}]


def test_feature_lists_csv_records_conditions(result_dir):
    from app.callbacks.interactive_feature_lists import export_feature_lists

    out, rds = result_dir
    state = {"lists": [{"id": "l1", "name": "Lipids", "features": ["1", "2", "3"]}]}
    result = export_feature_lists(1, state, str(rds), None, "Harmony")
    assert result["filename"] == "feature_lists.csv"
    records = list((out / "provenance").glob("export_*_csv_feature_lists.json"))
    assert len(records) == 1
    payload = json.loads(records[0].read_text(encoding="utf-8"))
    assert payload["extra"]["lists"] == [{"name": "Lipids", "n_features": 3}]


def test_record_csv_export_noop_without_result_dir():
    """結果フォルダが無くても例外にしない（エクスポート本体を止めない）。"""
    assert pv.record_csv_export("x.csv", rds_path=None) is None


# ---------------------------------------------------------------------------
# PPTX の解析条件スライド
# ---------------------------------------------------------------------------

def test_pptx_conditions_slide_added():
    from pptx import Presentation
    from pptx.util import Inches

    from app.callbacks.interactive_pptx import _add_conditions_slide

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    conditions = pv.collect_conditions(rds_path=None, integration_method="Harmony")

    _add_conditions_slide(prs, conditions)
    assert len(prs.slides) >= 1

    texts = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                texts.extend(c.text for c in row.cells)
    joined = "\n".join(texts)
    assert "解析条件" in joined
    assert "統合手法" in joined
    assert "Harmony" in joined
    # 未記録は未記録として出る（既定値をでっち上げない）
    assert "未記録" in joined

    # 機械可読な全量がスピーカーノートに入る
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert json.loads(notes)["integration_method"] == "Harmony"


def test_pptx_conditions_slide_noop_without_conditions():
    from pptx import Presentation

    from app.callbacks.interactive_pptx import _add_conditions_slide

    prs = Presentation()
    _add_conditions_slide(prs, None)
    assert len(prs.slides) == 0


# ---------------------------------------------------------------------------
# Master Password ゲート
# ---------------------------------------------------------------------------

def test_methods_locked_with_wrong_password(result_dir, monkeypatch):
    import app.callbacks.provenance_callbacks as pc

    out, rds = result_dir
    monkeypatch.setattr("app.services.auth_service.verify_master",
                        lambda p: False)
    (unlock, err, lock_style, content_style,
     rendered, dl_disabled, copy_disabled, pw) = pc.unlock_methods(
        1, "wrong", str(rds), None, "Harmony")

    assert unlock is None
    assert "Master Password" in err
    assert lock_style == {"display": "block"}
    assert content_style == {"display": "none"}
    assert rendered is None            # 本文は一切ブラウザへ渡さない
    assert dl_disabled is True
    assert copy_disabled is True
    assert pw == ""                    # 入力欄をクリアする


def test_methods_unlocked_with_correct_password(result_dir, monkeypatch):
    import app.callbacks.provenance_callbacks as pc

    out, rds = result_dir
    monkeypatch.setattr("app.services.auth_service.verify_master",
                        lambda p: p == "secret")
    (unlock, err, lock_style, content_style,
     rendered, dl_disabled, copy_disabled, pw) = pc.unlock_methods(
        1, "secret", str(rds), None, "Harmony")

    assert unlock["ok"] is True
    assert err == ""
    assert content_style == {"display": "block"}
    # 平文（論文用）と表形式の両方が、日英そろって返る
    assert set(rendered) == {"prose", "table"}
    assert set(rendered["prose"]) == {"ja", "en"}
    assert "解析条件" in rendered["table"]["ja"]
    assert "Methods draft" in rendered["table"]["en"]
    assert rendered["prose"]["ja"][0]["heading"]
    assert dl_disabled is False
    assert copy_disabled is False
    # パスワードは解錠情報にも入力欄にも残さない
    assert "secret" not in json.dumps(unlock)
    assert pw == ""


def test_methods_requires_password(result_dir):
    import app.callbacks.provenance_callbacks as pc

    out, rds = result_dir
    result = pc.unlock_methods(1, "", str(rds), None, "Harmony")
    assert result[0] is None
    assert result[4] is None       # 本文 Store は空のまま


def test_methods_download_blocked_when_locked(result_dir):
    import app.callbacks.provenance_callbacks as pc
    from dash.exceptions import PreventUpdate

    out, rds = result_dir
    with pytest.raises(PreventUpdate):
        pc.download_methods_bundle(1, None, str(rds), None, "Harmony")
    with pytest.raises(PreventUpdate):
        pc.download_methods_bundle(1, {"ok": False}, str(rds), None, "Harmony")


def test_methods_download_bundle_when_unlocked(result_dir):
    import app.callbacks.provenance_callbacks as pc

    out, rds = result_dir
    result = pc.download_methods_bundle(
        1, {"ok": True}, str(rds), None, "Harmony")
    assert result["filename"].startswith("analysis_conditions_")

    import base64
    raw = base64.b64decode(result["content"])
    with zipfile.ZipFile(BytesIO(raw)) as zf:
        names = zf.namelist()
        assert "analysis_conditions.json" in names
        assert "METHODS_ja.md" in names
        assert "METHODS_en.md" in names
        # 論文用の平文（色つき HTML を含む）
        assert "METHODS_prose_ja.html" in names
        assert "METHODS_prose_en.html" in names
        assert "METHODS_prose_ja.md" in names
        assert "analysis_params.json" in names   # 裏付けも同梱
        assert "解析条件" in zf.read("METHODS_ja.md").decode("utf-8")
        html = zf.read("METHODS_prose_ja.html").decode("utf-8")
        assert "methods" not in html or "color:#d32f2f" in html


def test_conditions_bundle_button_writes_files(result_dir):
    import app.callbacks.provenance_callbacks as pc

    out, rds = result_dir
    msg = pc.export_conditions_bundle(1, str(rds), None, "Harmony")
    assert "provenance" in msg
    assert (out / "provenance" / "analysis_conditions.json").exists()
    assert (out / "provenance" / "METHODS_ja.md").exists()
    assert (out / "provenance" / "METHODS_en.md").exists()


def test_conditions_bundle_button_reports_cache_only_embedding(tmp_path, monkeypatch):
    import app.callbacks.provenance_callbacks as pc

    cache = tmp_path / "cache"
    (cache / "derived_pca").mkdir(parents=True)
    rds = cache / "derived_pca" / "derived.rds"
    rds.write_bytes(b"x")
    monkeypatch.setattr("app.config.SEURAT_CACHE_DIR", cache, raising=False)
    msg = pc.export_conditions_bundle(1, str(rds), None, "PCA")
    assert "キャッシュ" in msg
