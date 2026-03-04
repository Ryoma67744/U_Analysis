# =============================================================================
# MSI Analysis Application - Interactive Analysis Callbacks
# インタラクティブ解析 コールバック
# =============================================================================

import base64
import io
from io import BytesIO
import json
import re
import tempfile
from datetime import datetime
from pathlib import Path

import numpy as np
import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
import dash_bootstrap_components as dbc
from dash import (Input, Output, State, callback, ctx, no_update, html, dcc,
                  dash_table, ALL, clientside_callback, ClientsideFunction)
from dash.exceptions import PreventUpdate

from app.config import (
    CLUSTER_PRESET_COLORS, DESI_COLORS_50, HIGHLIGHT_GRAY,
    DEFAULT_ADDUCT_POSITIVE,
)
from app.services.seurat_bridge import SeuratBridge

# Seuratブリッジのシングルトン
_bridge = SeuratBridge()

# モジュールレベルのデータキャッシュ
_interactive_data = {
    "plot_data": None,
    "cluster_stats": None,
    "features_list": None,
    "meta": None,
    "rds_path": None,
    "cache_dir": None,
}


def _cluster_sort_key(x):
    """クラスタIDの統一ソートキー（数値優先、文字列は末尾）"""
    s = str(x)
    return (int(s) if s.isdigit() else float("inf"), s)


def _get_cluster_color_map(clusters, custom_colors=None):
    """クラスタ値のリストから、ソート済みの色マップ dict を返す。
    custom_colors が指定された場合、デフォルト色をカスタム色で上書きする。"""
    str_cls = list(set(str(c) for c in clusters))
    str_cls.sort(key=_cluster_sort_key)
    cmap = {cl: CLUSTER_PRESET_COLORS[i % len(CLUSTER_PRESET_COLORS)] for i, cl in enumerate(str_cls)}
    if custom_colors:
        cmap.update(custom_colors)
    return cmap


def _get_sample_color_map(samples):
    """サンプル名のリストからサンプル色マップを生成（DESI_COLORS_50 使用）"""
    sorted_samples = sorted(set(str(s) for s in samples))
    return {s: DESI_COLORS_50[i % len(DESI_COLORS_50)]
            for i, s in enumerate(sorted_samples)}


def _is_meaningful_annotation(ann, gene=""):
    """annotation値が表示価値のある化合物名かどうかを判定。
    数値のみのアノテーション(例: "240.984")や空文字は除外する。"""
    if not ann or not isinstance(ann, str):
        return False
    ann = ann.strip()
    if not ann:
        return False
    # 数値のみ（小数点含む）は無意味なアノテーション
    if re.match(r'^[\d.]+$', ann):
        return False
    # geneと同一テキストも除外
    if ann == gene:
        return False
    return True


def _extract_mz_numeric(f):
    """フィーチャー名から数値部分(m/z値)を抽出してソート用floatを返す"""
    match = re.search(r"(\d+\.?\d*)", f)
    return float(match.group(1)) if match else float("inf")


# ---------------------------------------------------------------------------
# m/z キャリブレーション
# ---------------------------------------------------------------------------
def _calibrate_mz(features_list, expression_df, reference_mz,
                  search_window=0.5, min_peaks=2, regression_mode="linear"):
    """全ピクセル平均スペクトルから参照ピークのppmずれを計算し、回帰で補正値を返す。

    Args:
        features_list: Feature名リスト ("m/z 123.45678" 形式)
        expression_df: DataFrame (cells×features, 列名=feature名)
        reference_mz: list[float] — 参照m/z理論値
        search_window: float — 検索ウィンドウ(Da)
        min_peaks: int — 最低マッチ数
        regression_mode: str — "linear", "poly2", "poly3"

    Returns:
        dict: calibrated, corrected_mz_map, report, regression_mode, r_squared, ...
    """
    from scipy.stats import linregress

    # Feature名 → m/z数値マッピング
    mz_values = {}
    for f in features_list:
        mz = _extract_mz_numeric(f)
        if mz != float("inf"):
            mz_values[f] = mz

    if not mz_values:
        return {"calibrated": False, "corrected_mz_map": {}, "report": []}

    mz_array = np.array(list(mz_values.values()))
    feature_names = list(mz_values.keys())

    # 平均スペクトル算出
    avg_spectrum = {}
    for f in feature_names:
        if f in expression_df.columns:
            avg_spectrum[f] = expression_df[f].mean()
        else:
            avg_spectrum[f] = 0.0

    # 参照ピークとのマッチング
    matched = []
    for ref in reference_mz:
        within_window = np.where(np.abs(mz_array - ref) <= search_window)[0]
        if len(within_window) == 0:
            continue
        # ウィンドウ内で最大強度のピークを選択
        best_idx = None
        best_intensity = -1
        for idx in within_window:
            fname = feature_names[idx]
            intensity = avg_spectrum.get(fname, 0.0)
            if intensity > best_intensity:
                best_intensity = intensity
                best_idx = idx
        if best_idx is not None:
            obs = mz_array[best_idx]
            ppm = (obs - ref) / ref * 1e6
            matched.append({
                "ref_mz": ref, "obs_mz": float(obs),
                "ppm_drift": float(ppm), "avg_intensity": float(best_intensity),
            })

    if len(matched) < min_peaks:
        return {"calibrated": False, "corrected_mz_map": {}, "report": matched}

    # 回帰: ppm_drift を obs_mz の関数としてフィッティング
    obs_arr = np.array([m["obs_mz"] for m in matched])
    ppm_arr = np.array([m["ppm_drift"] for m in matched])

    if regression_mode in ("poly2", "poly3"):
        degree = 2 if regression_mode == "poly2" else 3
        degree = min(degree, len(obs_arr) - 1)  # 過学習防止
        coeffs = np.polyfit(obs_arr, ppm_arr, degree)
        predicted_ppm = np.polyval(coeffs, mz_array)
        # R² 算出
        fitted_ppm = np.polyval(coeffs, obs_arr)
        ss_res = np.sum((ppm_arr - fitted_ppm) ** 2)
        ss_tot = np.sum((ppm_arr - np.mean(ppm_arr)) ** 2)
        r_squared = 1 - ss_res / ss_tot if ss_tot > 0 else 0.0
    else:
        slope, intercept, r_value, _, _ = linregress(obs_arr, ppm_arr)
        predicted_ppm = slope * mz_array + intercept
        coeffs = None
        r_squared = float(r_value ** 2)

    # 全m/zに補正適用
    corrected_mz = mz_array / (1 + predicted_ppm / 1e6)
    corrected_mz_map = {f: float(c) for f, c in zip(feature_names, corrected_mz)}

    result = {
        "calibrated": True,
        "corrected_mz_map": corrected_mz_map,
        "report": matched,
        "regression_mode": regression_mode,
        "r_squared": r_squared,
    }
    if coeffs is not None:
        result["poly_coeffs"] = [float(c) for c in coeffs]
    else:
        result["slope"] = float(slope)
        result["intercept"] = float(intercept)
    return result


def _calibrate_mz_from_pairs(features_list, matched_pairs,
                             regression_mode="linear"):
    """ユーザーが明示的に対応付けたref/obsペアから直接キャリブレーション回帰を行う。

    Args:
        features_list: Feature名リスト ("m/z 123.45678" 形式)
        matched_pairs: list[dict] — {ref_mz, obs_mz, ppm_drift} のリスト
        regression_mode: str — "linear", "poly2", "poly3"

    Returns:
        dict: _calibrate_mz() と同一形式
    """
    from scipy.stats import linregress

    if len(matched_pairs) < 2:
        return {"calibrated": False, "corrected_mz_map": {}, "report": matched_pairs}

    obs_arr = np.array([p["obs_mz"] for p in matched_pairs])
    ppm_arr = np.array([p["ppm_drift"] for p in matched_pairs])

    if regression_mode in ("poly2", "poly3"):
        degree = 2 if regression_mode == "poly2" else 3
        degree = min(degree, len(obs_arr) - 1)  # 過学習防止
        coeffs = np.polyfit(obs_arr, ppm_arr, degree)
        # R² 算出
        fitted_ppm = np.polyval(coeffs, obs_arr)
        ss_res = np.sum((ppm_arr - fitted_ppm) ** 2)
        ss_tot = np.sum((ppm_arr - np.mean(ppm_arr)) ** 2)
        r_squared = 1 - ss_res / ss_tot if ss_tot > 0 else 0.0
    else:
        slope, intercept, r_value, _, _ = linregress(obs_arr, ppm_arr)
        coeffs = None
        r_squared = float(r_value ** 2)

    # Feature名 → m/z数値マッピング
    mz_values = {}
    for f in features_list:
        mz = _extract_mz_numeric(f)
        if mz != float("inf"):
            mz_values[f] = mz

    if not mz_values:
        return {"calibrated": False, "corrected_mz_map": {}, "report": matched_pairs}

    mz_array = np.array(list(mz_values.values()))
    feature_names = list(mz_values.keys())

    # 全m/zに補正適用
    if coeffs is not None:
        predicted_ppm = np.polyval(coeffs, mz_array)
    else:
        predicted_ppm = slope * mz_array + intercept
    corrected_mz = mz_array / (1 + predicted_ppm / 1e6)
    corrected_mz_map = {f: float(c) for f, c in zip(feature_names, corrected_mz)}

    result = {
        "calibrated": True,
        "corrected_mz_map": corrected_mz_map,
        "report": matched_pairs,
        "regression_mode": regression_mode,
        "r_squared": r_squared,
    }
    if coeffs is not None:
        result["poly_coeffs"] = [float(c) for c in coeffs]
    else:
        result["slope"] = float(slope)
        result["intercept"] = float(intercept)
    return result


def _reannotate_with_calibration(deg_data, corrected_mz_map, mrm_path, tolerance=0.1,
                                  annotation_csv_path=None, ion_mode="Positive",
                                  adduct_patterns=None):
    """補正済みm/zでアノテーションDBと照合し、DEGデータのannotation列を更新する。

    Args:
        deg_data: list[dict] — DEGレコードリスト
        corrected_mz_map: dict[str, float] — feature名 → 補正済みm/z
        mrm_path: str — MRM/TraceFinder DBファイルパス
        tolerance: float — マッチング許容誤差(Da)
        annotation_csv_path: str — アノテーションCSV（TraceFinder/HMDB形式）パス
        ion_mode: str — "Positive" or "Negative"
        adduct_patterns: list[str] — 付加イオンフィルター

    Returns:
        list[dict] — annotation更新済みDEGデータ
    """
    # アノテーションCSV（TraceFinder/HMDB）を優先、MRMをフォールバック
    mz_to_compound = _build_mz_to_compound_map(mrm_path, tolerance=tolerance)
    ann_map = _build_annotation_csv_map(
        annotation_csv_path, ion_mode=ion_mode,
        adduct_patterns=adduct_patterns, tolerance=tolerance,
    )
    mz_to_compound.update(ann_map)  # ann_map を優先
    if not mz_to_compound:
        return deg_data

    mrm_mz_values = np.array(sorted(mz_to_compound.keys()))
    if len(mrm_mz_values) == 0:
        return deg_data

    updated = []
    for row in deg_data:
        row = dict(row)  # コピー
        gene = row.get("gene", "")
        corrected = corrected_mz_map.get(gene)
        if corrected is not None and len(mrm_mz_values) > 0:
            idx = np.argmin(np.abs(mrm_mz_values - corrected))
            if abs(mrm_mz_values[idx] - corrected) <= tolerance:
                compound = mz_to_compound[mrm_mz_values[idx]]
                row["annotation"] = compound
        updated.append(row)
    return updated


def _get_label_positions_path():
    """label_positions.json のパスを返す（RDSと同ディレクトリ）"""
    rds_path = _interactive_data.get("rds_path")
    if not rds_path:
        return None
    return Path(rds_path).parent / "label_positions.json"


def _load_label_positions():
    """label_positions.json を読み込んで dict を返す。ファイルなし or エラー時は空dict"""
    path = _get_label_positions_path()
    if path and path.exists():
        try:
            return json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


def _get_interactive_settings_path():
    """interactive_settings.json のパスを返す（RDSと同ディレクトリ）"""
    rds_path = _interactive_data.get("rds_path")
    if not rds_path:
        return None
    return Path(rds_path).parent / "interactive_settings.json"


def _load_interactive_settings():
    """interactive_settings.json を読み込み。ファイルなし/エラー時は空dict"""
    path = _get_interactive_settings_path()
    if path and path.exists():
        try:
            return json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


def _save_interactive_settings(key, value):
    """interactive_settings.json の指定キーを更新して書き込む"""
    path = _get_interactive_settings_path()
    if not path:
        return
    try:
        existing = _load_interactive_settings()
        existing[key] = value
        existing["_saved_at"] = datetime.now().strftime("%Y-%m-%dT%H:%M:%S")
        path.write_text(
            json.dumps(existing, indent=2, ensure_ascii=False),
            encoding="utf-8",
        )
    except Exception:
        pass


def _extract_annotation_positions_by_name(relayout_data, clusters):
    """relayoutData の annotations[N].x/y をクラスタ名ベースで抽出。

    Parameters
    ----------
    relayout_data : dict  – Dash relayoutData
    clusters : list       – ソート済みクラスタ名リスト（annotation 追加順と一致）

    Returns
    -------
    dict – {"クラスタ名": {"x": v, "y": v}}  （更新があったもののみ）
    """
    if not relayout_data or not clusters:
        return {}
    positions = {}
    sorted_clusters = list(clusters)
    for key, val in relayout_data.items():
        m = re.match(r"annotations\[(\d+)\]\.([xy])", key)
        if m:
            idx = int(m.group(1))
            attr = m.group(2)
            if idx < len(sorted_clusters):
                cl = str(sorted_clusters[idx])
                if cl not in positions:
                    positions[cl] = {}
                positions[cl][attr] = val
    return positions


def _merge_label_positions(base, overlay):
    """base dict に overlay dict をディープマージ（overlay が優先）。

    base/overlay は {"key": {"x": v, "y": v}} 形式。
    base を破壊的に更新して返す。
    """
    for k, v in (overlay or {}).items():
        if k not in base:
            base[k] = {}
        if isinstance(v, dict):
            base[k].update(v)
        else:
            base[k] = v
    return base


def _get_cluster_colorscale(clusters, custom_colors=None):
    """Scattergl用: 数値インデックスベースのcolorscale情報を返す。

    HEX文字列配列をmarker.colorに渡すとWebGL内部処理で色ミスマッチが
    生じるため、数値+colorscaleで確実に色を指定する。

    Returns:
        cluster_to_idx: dict[str, int] — クラスタ文字列→0-based数値インデックス
        discrete_colorscale: list — Plotly colorscale形式
    """
    str_cls = list(set(str(c) for c in clusters))
    str_cls.sort(key=_cluster_sort_key)
    n = max(len(str_cls), 1)
    cluster_to_idx = {cl: i for i, cl in enumerate(str_cls)}

    # discrete colorscale: 各色が均等な範囲を占める
    colorscale = []
    for i, cl in enumerate(str_cls):
        low = i / n
        high = (i + 1) / n
        color = CLUSTER_PRESET_COLORS[i % len(CLUSTER_PRESET_COLORS)]
        if custom_colors and cl in custom_colors:
            color = custom_colors[cl]
        colorscale.append([low, color])
        colorscale.append([high, color])

    return cluster_to_idx, colorscale


# ---------------------------------------------------------------------------
# 統合手法検出ヘルパー
# ---------------------------------------------------------------------------

def _detect_integration_methods(folder_path: str) -> dict:
    """結果フォルダ内のRDSファイルを検出し、統合手法→パスのマッピングを返す。

    Returns:
        {"Harmony": "path/to/seu_harmony.rds", "RPCA": "path/to/seu_rpca.rds", ...}
    """
    rds_map = {}
    base = Path(folder_path)
    if not base.is_dir():
        return rds_map

    # RDS_Files/ フォルダ内を検索
    rds_dir = base / "RDS_Files"
    search_dirs = [rds_dir, base] if rds_dir.is_dir() else [base]

    # data.frame 型 RDS を除外するプレフィックス
    _EXCLUDE_PREFIXES = ("umap_", "deg_", "plotdata_", "feature_")

    # 第1段階: TIMS ver13 の Step2/Step3 ファイルを優先マッチ
    for search_dir in search_dirs:
        for rds_file in search_dir.glob("*.rds"):
            name_lower = rds_file.name.lower()
            if "step2" in name_lower and "harmony" in name_lower:
                rds_map["Harmony"] = str(rds_file)
            elif "step3" in name_lower and "rpca" in name_lower:
                rds_map["RPCA"] = str(rds_file)
            elif "single" in name_lower and "PCA" not in rds_map:
                rds_map["PCA"] = str(rds_file)

    # 第2段階: 第1段階で見つからなかったキーのみ、従来マッチ（data.frame除外）
    for search_dir in search_dirs:
        for rds_file in search_dir.glob("*.rds"):
            name_lower = rds_file.name.lower()
            if any(name_lower.startswith(p) for p in _EXCLUDE_PREFIXES):
                continue
            if "harmony" in name_lower and "Harmony" not in rds_map:
                rds_map["Harmony"] = str(rds_file)
            elif "rpca" in name_lower and "RPCA" not in rds_map:
                rds_map["RPCA"] = str(rds_file)

    # rglob でサブフォルダも検索（上記で見つからない場合のフォールバック）
    if not rds_map:
        # 第1段階: Step2/Step3 優先
        for rds_file in base.rglob("*.rds"):
            name_lower = rds_file.name.lower()
            if "step2" in name_lower and "harmony" in name_lower:
                rds_map["Harmony"] = str(rds_file)
            elif "step3" in name_lower and "rpca" in name_lower:
                rds_map["RPCA"] = str(rds_file)
            elif "single" in name_lower and "PCA" not in rds_map:
                rds_map["PCA"] = str(rds_file)

        # 第2段階: 従来マッチ（data.frame除外）
        if "Harmony" not in rds_map or "RPCA" not in rds_map:
            for rds_file in base.rglob("*.rds"):
                name_lower = rds_file.name.lower()
                if any(name_lower.startswith(p) for p in _EXCLUDE_PREFIXES):
                    continue
                if "harmony" in name_lower and "Harmony" not in rds_map:
                    rds_map["Harmony"] = str(rds_file)
                elif "rpca" in name_lower and "RPCA" not in rds_map:
                    rds_map["RPCA"] = str(rds_file)

    return rds_map


def _build_msi_samples_ui(folder_path: str):
    """MSIフォルダ内のサンプル一覧UIを生成"""
    if not folder_path or not Path(folder_path).is_dir():
        return html.Div("MSIフォルダが見つかりません", className="text-muted")

    from app.services.data_manager import list_msi_files, list_tims_files
    samples = list_msi_files(folder_path)
    # .txtが見つからない場合はTIMS形式(.parquet等)を試行
    if not samples:
        samples = list_tims_files(folder_path)
    if not samples:
        return html.Div("MSIファイルが見つかりません", className="text-warning")

    import dash_bootstrap_components as dbc
    return dbc.Checklist(
        id="interactive_msi_sample_checks",
        options=[{"label": s, "value": s} for s in samples],
        value=samples,
    )


# ---------------------------------------------------------------------------
# 結果フォルダスキャン → 統合手法検出
# ---------------------------------------------------------------------------

@callback(
    [Output("interactive_integration_method", "options"),
     Output("interactive_integration_method", "value"),
     Output("interactive_rds_map", "data")],
    Input("scan_result_folder", "n_clicks"),
    State("interactive_result_folder", "value"),
    prevent_initial_call=True,
)
def scan_rds_files(n_clicks, folder_path):
    if not folder_path or not Path(folder_path).is_dir():
        return [], None, None

    rds_map = _detect_integration_methods(folder_path)
    if not rds_map:
        return [], None, None

    options = [{"label": k, "value": k} for k in rds_map.keys()]
    # Harmony を優先デフォルト、なければ最初の手法
    default = "Harmony" if "Harmony" in rds_map else list(rds_map.keys())[0]

    return options, default, rds_map


# ---------------------------------------------------------------------------
# 結果フォルダ変更時 → 自動スキャン
# ---------------------------------------------------------------------------

@callback(
    [Output("interactive_integration_method", "options", allow_duplicate=True),
     Output("interactive_integration_method", "value", allow_duplicate=True),
     Output("interactive_rds_map", "data", allow_duplicate=True)],
    Input("interactive_result_folder", "value"),
    prevent_initial_call=True,
)
def auto_scan_rds_files(folder_path):
    """結果フォルダのパスが設定された時、自動で統合手法を検出"""
    if not folder_path or not Path(folder_path).is_dir():
        return no_update, no_update, no_update

    rds_map = _detect_integration_methods(folder_path)
    if not rds_map:
        return no_update, no_update, no_update

    options = [{"label": k, "value": k} for k in rds_map.keys()]
    default = "Harmony" if "Harmony" in rds_map else list(rds_map.keys())[0]

    return options, default, rds_map


# ---------------------------------------------------------------------------
# MSIフォルダスキャン（ボタンクリック）
# ---------------------------------------------------------------------------

@callback(
    Output("interactive_msi_samples", "children"),
    Input("scan_msi_folder", "n_clicks"),
    State("interactive_msi_folder", "value"),
    prevent_initial_call=True,
)
def scan_msi_files(n_clicks, folder_path):
    return _build_msi_samples_ui(folder_path)


# ---------------------------------------------------------------------------
# MSIフォルダ変更時 → 自動スキャン
# ---------------------------------------------------------------------------

@callback(
    Output("interactive_msi_samples", "children", allow_duplicate=True),
    Input("interactive_msi_folder", "value"),
    prevent_initial_call=True,
)
def auto_scan_msi_files(folder_path):
    """MSIフォルダのパスが設定された時、自動でサンプル一覧を表示"""
    if not folder_path or not Path(folder_path).is_dir():
        return no_update
    return _build_msi_samples_ui(folder_path)


# ---------------------------------------------------------------------------
# 解析手法セクション 展開/折りたたみ
# ---------------------------------------------------------------------------

@callback(
    [Output("integration_method_collapse", "is_open"),
     Output("toggle_integration_method", "children")],
    Input("toggle_integration_method", "n_clicks"),
    State("integration_method_collapse", "is_open"),
    prevent_initial_call=True,
)
def toggle_integration_method(n, is_open):
    """解析手法セクションの展開/折りたたみを切り替え"""
    new_state = not is_open
    label = "解析手法 \u25bc" if new_state else "解析手法 \u25b6"
    return new_state, label


# ---------------------------------------------------------------------------
# データ読み込み（Seuratブリッジ経由）
# ---------------------------------------------------------------------------

@callback(
    [Output("interactive_data_info", "children"),
     Output("interactive_viz_container", "style"),
     Output("umap_highlight_cluster", "options"),
     Output("interactive_sample", "options"),
     Output("feature_select", "options"),
     Output("seurat_rds_path_store", "data"),
     Output("seurat_cache_dir_store", "data"),
     Output("deg_data_store", "data", allow_duplicate=True),
     Output("deg_results_section", "style"),
     Output("spatial_exclude_cluster", "options"),
     Output("spatial_highlight_cluster", "options"),
     Output("umap_exclude_cluster", "options"),
     Output("feature_sample_select", "options"),
     Output("deg_no_data_message", "style"),
     Output("feature_cluster_filter", "options"),
     Output("sample_name_map_store", "data", allow_duplicate=True),
     Output("spatial_rotation_store", "data", allow_duplicate=True),
     Output("custom_color_map_store", "data", allow_duplicate=True),
     Output("feature_history_store", "data", allow_duplicate=True),
     # キャリブレーション設定復元 (11個)
     Output("int_cal_table_data", "data", allow_duplicate=True),
     Output("int_cal_enable", "value"),
     Output("int_cal_ion_mode", "value"),
     Output("int_cal_matrix", "value"),
     Output("int_cal_adduct_filter", "value", allow_duplicate=True),
     Output("int_cal_mrm_path", "value"),
     Output("int_cal_search_window", "value"),
     Output("int_cal_min_peaks", "value"),
     Output("int_cal_regression_mode", "value"),
     Output("int_cal_ms_instrument", "data", allow_duplicate=True),
     Output("int_cal_restore_pending", "data"),
     Output("sap_btn_wrapper", "style")],
    [Input("load_interactive_data", "n_clicks"),
     Input("interactive_integration_method", "value"),
     Input("interactive_rds_map", "data")],
    [State("interactive_result_folder", "value"),
     State("calibration_enable", "value"),
     State("calibration_matrix", "value"),
     State("calibration_table_data", "data"),
     State("calibration_search_window", "value"),
     State("calibration_min_peaks", "value"),
     State("calibration_regression_mode", "value"),
     State("ion_mode", "value"),
     State("mrm_path", "value"),
     State("tolerance_mz", "value"),
     State("adduct_filter", "value"),
     State("interactive_project_select", "value"),
     State("interactive_sub_project_select", "value"),
     State("default_annotation_csv", "value")],
    prevent_initial_call=True,
)
def load_interactive_data(n_clicks, integration_method, rds_map, result_folder,
                          cal_enable, cal_matrix, cal_table_data,
                          cal_search_window, cal_min_peaks,
                          cal_regression_mode,
                          ion_mode, mrm_path, tolerance_mz,
                          adduct_filter, project_id, sub_project_id,
                          annotation_csv):
    _n_out = 31
    _no_cal = (no_update,) * 11  # キャリブレーション設定復元用
    _sap_hide = ({"display": "none"},)  # sap_btn_wrapper 非表示
    if not integration_method or not rds_map:
        return (
            "統合手法を選択してください（結果フォルダをスキャンしてください）",
            {"display": "none"}, [], [], [], None, None, None,
            {"display": "none"}, [], [], [], [],
            {"display": "none"}, [],
            no_update, no_update, no_update, no_update,
        ) + _no_cal + _sap_hide

    rds_path = rds_map.get(integration_method)
    if not rds_path or not Path(rds_path).exists():
        return (
            f"RDSファイルが見つかりません: {integration_method}",
            {"display": "none"}, [], [], [], None, None, None,
            {"display": "none"}, [], [], [], [],
            {"display": "none"}, [],
            no_update, no_update, no_update, no_update,
        ) + _no_cal + _sap_hide

    try:
        result = _bridge.extract_data(rds_path)

        _interactive_data["plot_data"] = result["plot_data"]
        _interactive_data["cluster_stats"] = result["cluster_stats"]
        _interactive_data["features_list"] = result["features_list"]
        _interactive_data["meta"] = result["meta"]
        _interactive_data["rds_path"] = rds_path
        _interactive_data["cache_dir"] = result.get("cache_dir")

        meta = result["meta"]
        info_text = (
            f"読み込み完了 [{integration_method}]: "
            f"{meta.get('n_cells', '?')} cells, "
            f"{meta.get('n_clusters', '?')} clusters, "
            f"samples: {', '.join(meta.get('samples', []))}"
        )

        # クラスタ選択肢
        clusters = sorted(_interactive_data["plot_data"]["Cluster"].unique(), key=_cluster_sort_key)
        cluster_options = [
            {"label": f"Cluster {c}", "value": str(c)} for c in clusters
        ]

        # サンプル選択肢
        samples = sorted(_interactive_data["plot_data"]["Sample"].unique())
        sample_options = [{"label": s, "value": s} for s in samples]

        # Feature選択肢（全件、dcc.Dropdownのsearch機能でフィルタ）
        features = result["features_list"]
        feature_options = [{"label": f, "value": f} for f in features]

        # DEG 結果を探す（選択した統合手法のフォルダを優先）
        deg_data = None
        if result_folder:
            result_base = Path(result_folder)
            deg_data = _load_deg_results(result_base, integration_method)
        else:
            rds_dir = Path(rds_path).parent
            result_base = rds_dir.parent if rds_dir.name == "RDS_Files" else rds_dir
            deg_data = _load_deg_results(result_base, integration_method)

        # --- m/z キャリブレーション（有効時のみ） ---
        if cal_enable and deg_data and mrm_path and cal_table_data:
            try:
                # テーブルから use="Yes" のペアを抽出
                matched_pairs = []
                ref_only_mz = []
                for row in cal_table_data:
                    if row.get("use") != "Yes":
                        continue
                    ref = row.get("ref_mz")
                    obs = row.get("obs_mz")
                    if ref and obs and str(ref).strip() and str(obs).strip():
                        ref_f = float(ref)
                        obs_f = float(obs)
                        ppm = (obs_f - ref_f) / ref_f * 1e6
                        matched_pairs.append({
                            "ref_mz": ref_f, "obs_mz": obs_f,
                            "ppm_drift": ppm,
                        })
                    elif ref and str(ref).strip():
                        ref_only_mz.append(float(ref))

                mp = int(cal_min_peaks or 2)
                reg_mode = cal_regression_mode or "linear"
                cal_result = None

                if len(matched_pairs) >= mp:
                    # 明示的ペアが十分 → 直接回帰
                    cal_result = _calibrate_mz_from_pairs(
                        result["features_list"], matched_pairs,
                        regression_mode=reg_mode,
                    )
                elif ref_only_mz:
                    # obs未入力行あり → 自動検出にフォールバック
                    cache_dir = result.get("cache_dir")
                    expr_path = (Path(cache_dir) / "expression_matrix.parquet"
                                 if cache_dir else None)
                    if expr_path and expr_path.exists():
                        expr_df = pd.read_parquet(expr_path)
                        sw = float(cal_search_window or 0.5)
                        cal_result = _calibrate_mz(
                            result["features_list"], expr_df, ref_only_mz,
                            search_window=sw, min_peaks=mp,
                            regression_mode=reg_mode,
                        )

                if cal_result and cal_result.get("calibrated"):
                    tol = float(tolerance_mz or 0.1)
                    deg_data = _reannotate_with_calibration(
                        deg_data, cal_result["corrected_mz_map"],
                        mrm_path, tolerance=tol,
                        annotation_csv_path=annotation_csv,
                        ion_mode=ion_mode,
                        adduct_patterns=adduct_filter,
                    )
                    _interactive_data["_calibration_result"] = cal_result
            except Exception:
                pass  # キャリブレーション失敗時はそのまま続行

        # DEGセクションは常に表示、データ有無でメッセージ切替
        deg_section_style = {}  # 常に表示
        deg_no_data_style = {"display": "none"} if deg_data else {}
        # クラスタフィルタ選択肢（DEGデータから生成）
        cluster_filter_opts = []
        if deg_data:
            deg_clusters = sorted(set(str(r.get("cluster", "")) for r in deg_data), key=_cluster_sort_key)
            cluster_filter_opts = [
                {"label": f"Cluster {c}", "value": c} for c in deg_clusters
            ]

        # 保存済み設定を読み込み（初回ロード時にStoreへ復元）
        saved = _load_interactive_settings()
        restored_name_map = saved.get("sample_name_map", {})
        restored_rotation = saved.get("spatial_rotation", {})
        restored_colors = saved.get("custom_color_map", {})
        restored_bookmarks = saved.get("feature_bookmarks", [])
        if restored_name_map:
            _interactive_data["_name_map"] = restored_name_map

        # --- キャリブレーション設定の復元 ---
        int_cal = saved.get("int_calibration", {})
        if int_cal:
            r_table = int_cal.get("table_data", [])
            r_enable = int_cal.get("enable", False)
            r_ion_mode = int_cal.get("ion_mode", ion_mode or "Positive")
            r_matrix = int_cal.get("matrix", cal_matrix or "DHB")
            r_adduct = int_cal.get("adduct_filter",
                                   adduct_filter or DEFAULT_ADDUCT_POSITIVE)
            r_mrm = int_cal.get("mrm_path", mrm_path or "")
            r_sw = int_cal.get("search_window", cal_search_window or 0.5)
            r_mp = int_cal.get("min_peaks", cal_min_peaks or 2)
            r_reg = int_cal.get("regression_mode", cal_regression_mode or "poly3")
        else:
            # フォールバック: 解析設定タブの値
            r_table = cal_table_data or []
            r_enable = cal_enable or False
            r_ion_mode = ion_mode or "Positive"
            r_matrix = cal_matrix or "DHB"
            r_adduct = adduct_filter or DEFAULT_ADDUCT_POSITIVE
            r_mrm = mrm_path or ""
            r_sw = cal_search_window or 0.5
            r_mp = cal_min_peaks or 2
            r_reg = cal_regression_mode or "poly3"

        # ms_instrument をサブプロジェクトから取得
        r_instrument = "TIMS"
        if sub_project_id and project_id:
            try:
                from app.services.project_manager import get_sub_project
                sub = get_sub_project(project_id, sub_project_id)
                if sub:
                    r_instrument = sub.get("ms_instrument", "TIMS")
            except Exception:
                pass

        return (
            info_text,
            {},  # 可視化コンテナ表示
            cluster_options,
            sample_options,
            feature_options,
            rds_path,
            str(result.get("cache_dir", "")),
            deg_data,
            deg_section_style,
            cluster_options,  # spatial_exclude_cluster用
            cluster_options,  # spatial_highlight_cluster用
            cluster_options,  # umap_exclude_cluster用
            sample_options,   # feature_sample_select用
            deg_no_data_style,
            cluster_filter_opts,
            restored_name_map,
            restored_rotation,
            restored_colors,
            restored_bookmarks,
            # キャリブレーション設定復元 (11個)
            r_table, r_enable, r_ion_mode, r_matrix, r_adduct, r_mrm,
            r_sw, r_mp, r_reg, r_instrument, True,
            {},  # sap_btn_wrapper 表示
        )

    except Exception as e:
        return (
            f"読み込みエラー: {e}",
            {"display": "none"}, [], [], [], None, None, None,
            {"display": "none"}, [], [], [], [],
            {"display": "none"}, [],
            no_update, no_update, no_update, no_update,
        ) + _no_cal + _sap_hide


def _standardize_deg_df(df: pd.DataFrame) -> list[dict] | None:
    """DEG DataFrame の列名を標準化し、dict のリストとして返す。
    CSV / RDS 両方の読み込みから共通で使用する。"""
    try:
        # 列名を標準化
        col_map = {}
        for col in df.columns:
            cl = col.lower().strip()
            if cl in ("gene", "row.names", "x", "...1"):
                col_map[col] = "gene"
            elif "cluster" in cl:
                col_map[col] = "cluster"
            elif "avg_log2fc" in cl or "avg_logfc" in cl:
                col_map[col] = "avg_log2FC"
            elif "p_val_adj" in cl:
                col_map[col] = "p_val_adj"
            elif cl == "pct.1":
                col_map[col] = "pct.1"
            elif cl == "pct.2":
                col_map[col] = "pct.2"

        df = df.rename(columns=col_map)
        # gene列がない場合、最初の列をgeneとする
        if "gene" not in df.columns and len(df.columns) > 0:
            df = df.rename(columns={df.columns[0]: "gene"})

        # 必要な列のみ抽出
        keep = [c for c in ["gene", "cluster", "avg_log2FC", "p_val_adj",
                             "pct.1", "pct.2", "annotation"] if c in df.columns]
        df = df[keep]

        # 数値列を丸める
        for col in ["avg_log2FC", "p_val_adj", "pct.1", "pct.2"]:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors="coerce")
                if col == "p_val_adj":
                    # Volcano Plot用に元の数値を保持
                    df["p_val_adj_raw"] = df[col].copy()
                    # テーブル表示用に科学記数法文字列へ変換
                    df[col] = df[col].map(
                        lambda x: f"{x:.2e}" if pd.notna(x) else ""
                    )
                else:
                    df[col] = df[col].round(4)

        return df.to_dict("records")
    except Exception as e:
        print(f"[DEG] _standardize_deg_df エラー: {e}")
        import traceback; traceback.print_exc()
        return None


def _read_deg_rds(rds_path: Path) -> list[dict] | None:
    """TIMS ver13 が出力する deg_FindAllMarkers_raw_*.rds を読み込む。
    R subprocess で RDS → 一時CSV → pandas DataFrame に変換。"""
    import subprocess
    import tempfile

    tmp_csv = Path(tempfile.mktemp(suffix=".csv"))
    try:
        r_cmd = (
            f'deg <- readRDS("{rds_path.as_posix()}");\n'
            f'write.csv(deg, "{tmp_csv.as_posix()}", row.names=TRUE)'
        )
        result = subprocess.run(
            ["Rscript", "-e", r_cmd],
            capture_output=True, timeout=30,
        )
        if not tmp_csv.exists():
            return None
        df = pd.read_csv(tmp_csv)
        return _standardize_deg_df(df)
    except Exception as e:
        print(f"[DEG] _read_deg_rds error: {e}")
        return None
    finally:
        tmp_csv.unlink(missing_ok=True)


def _load_deg_results(
    result_base: Path, integration_method: str | None = None
) -> list[dict] | None:
    """解析結果フォルダ内の DEG CSV / RDS を読み込む"""
    # 選択した統合手法のフォルダを優先検索
    if integration_method and integration_method in ("Harmony", "RPCA", "PCA"):
        method_dir = integration_method
    else:
        method_dir = "Harmony"

    # --- 1. CSV ファイル検索 ---
    csv_patterns = [
        f"{method_dir}/*deg*markers*.csv",
        f"{method_dir}/*top*markers*.csv",
        f"{method_dir}/markers_annotated*.csv",
        f"{method_dir}/markers_mz_only*.csv",
        "Harmony/*deg*markers*.csv",
        "Harmony/*top*markers*.csv",
        "RPCA/*deg*markers*.csv",
        "RPCA/*top*markers*.csv",
        "PCA/*deg*markers*.csv",
        "PCA/*top*markers*.csv",
        "*deg*markers*.csv",
        "*top*markers*.csv",
        "markers_annotated*.csv",
        "markers_mz_only*.csv",
    ]

    def _try_load_csv(matches):
        """マッチしたCSVファイルの読み込みを試行"""
        for csv_path in matches:
            try:
                df = pd.read_csv(csv_path, encoding="utf-8")
                print(f"[DEG] CSV発見: {csv_path} (列: {list(df.columns)}, 行数: {len(df)})")
                result = _standardize_deg_df(df)
                if result:
                    print(f"[DEG] 読み込み成功: {len(result)} レコード")
                    return result
                else:
                    print(f"[DEG] _standardize_deg_df が None を返しました: {csv_path}")
            except Exception as e:
                print(f"[DEG] CSV読み込みエラー: {csv_path} — {e}")
        return None

    # 第1段階: result_base 直下を glob で検索
    for pattern in csv_patterns:
        matches = sorted(result_base.glob(pattern))
        if matches:
            result = _try_load_csv(matches)
            if result:
                return result

    # 第2段階: result_base 以下を rglob で再帰検索
    # （TIMS ver13 は日付サブフォルダ内に出力するため）
    rglob_csv_names = [
        "markers_annotated*.csv",
        "markers_mz_only*.csv",
        "*deg*markers*.csv",
        "*top*markers*.csv",
    ]
    for name_pattern in rglob_csv_names:
        matches = sorted(result_base.rglob(name_pattern))
        if matches:
            # 選択した統合手法のフォルダ内を優先
            method_lower = method_dir.lower()
            prioritized = [m for m in matches if method_lower in m.parent.name.lower()]
            ordered = prioritized + [m for m in matches if m not in prioritized]
            result = _try_load_csv(ordered)
            if result:
                return result

    # --- 2. DEG RDS ファイル検索（TIMS ver13 が出力する deg_FindAllMarkers_raw_*.rds） ---
    rds_patterns = [
        f"{method_dir}/deg_FindAllMarkers_raw_*.rds",
        "RDS_Files/deg_FindAllMarkers_raw_*.rds",
        "deg_FindAllMarkers_raw_*.rds",
    ]
    for pattern in rds_patterns:
        matches = sorted(result_base.glob(pattern))
        if matches:
            try:
                print(f"[DEG] RDS発見: {matches[0]}")
                result = _read_deg_rds(matches[0])
                if result:
                    print(f"[DEG] RDS読み込み成功: {len(result)} レコード")
                    return result
                else:
                    print(f"[DEG] _read_deg_rds が None を返しました: {matches[0]}")
            except Exception as e:
                print(f"[DEG] RDS読み込みエラー: {matches[0]} — {e}")
                continue

    # 第2段階(RDS): rglob でサブフォルダも再帰検索
    rds_rglob_matches = sorted(result_base.rglob("deg_FindAllMarkers_raw_*.rds"))
    if rds_rglob_matches:
        try:
            print(f"[DEG] RDS発見(rglob): {rds_rglob_matches[0]}")
            result = _read_deg_rds(rds_rglob_matches[0])
            if result:
                print(f"[DEG] RDS読み込み成功: {len(result)} レコード")
                return result
        except Exception as e:
            print(f"[DEG] RDS読み込みエラー(rglob): {rds_rglob_matches[0]} — {e}")

    print(f"[DEG] result_base={result_base}, method_dir={method_dir} — DEGファイル見つからず")
    return None


# ---------------------------------------------------------------------------
# フィーチャー検索（サーバーサイドフィルタ）
# ---------------------------------------------------------------------------

@callback(
    Output("feature_select", "options", allow_duplicate=True),
    [Input("feature_select", "search_value"),
     Input("feature_filter_mode", "value"),
     Input("feature_cluster_filter", "value")],
    [State("feature_mz_filtered_list", "data"),
     State("deg_data_store", "data")],
    prevent_initial_call=True,
)
def filter_features(search_value, filter_mode, cluster_filter,
                    mz_filtered, deg_data):
    """フィーチャードロップダウンの検索値に基づいてオプションをフィルタ"""
    features = _interactive_data.get("features_list")
    if not features:
        return []

    # annotation マッピングを構築（deg_dataから）
    ann_map = {}
    if deg_data:
        for r in deg_data:
            gene = r.get("gene", "")
            ann = r.get("annotation", "")
            if gene and _is_meaningful_annotation(ann, gene):
                ann_map[gene] = ann

    def _make_option(f, rank=None):
        """フィーチャー名からドロップダウン用オプションを生成"""
        prefix = f"★{rank} " if rank is not None else ""
        if f in ann_map:
            return {"label": f"{prefix}{f} ({ann_map[f]})", "value": f}
        return {"label": f"{prefix}{f}", "value": f}

    # DEGマーカーモード: クラスタのマーカーm/zのみ表示
    top15_ordered = None
    rest_ordered = None

    if filter_mode == "deg" and deg_data:
        if cluster_filter:
            # 選択クラスタのDEGレコードを抽出
            cluster_records = [
                r for r in deg_data
                if str(r.get("cluster", "")) == str(cluster_filter)
            ]

            # p値昇順（最も有意が先頭）、|log2FC|降順でタイブレーク
            def _sort_key(r):
                p = r.get("p_val_adj_raw")
                if p is None or (isinstance(p, float) and np.isnan(p)):
                    p = 1.0
                fc = r.get("avg_log2FC", 0)
                if fc is None or (isinstance(fc, float) and np.isnan(fc)):
                    fc = 0.0
                return (float(p), -abs(float(fc)))

            cluster_records_sorted = sorted(cluster_records, key=_sort_key)

            # Top 15 ユニーク遺伝子を抽出
            seen = set()
            top15_genes = []
            for r in cluster_records_sorted:
                g = str(r.get("gene", ""))
                if g and g not in seen:
                    seen.add(g)
                    top15_genes.append(g)
                    if len(top15_genes) >= 15:
                        break

            top15_set = set(top15_genes)
            features_set = set(features)

            # features_listに存在するもののみ保持
            top15_ordered = [f for f in top15_genes if f in features_set]

            # 残りのDEG遺伝子をm/z数値順でソート
            all_deg_set = set(str(r.get("gene", "")) for r in cluster_records)
            rest_ordered = sorted(
                [f for f in features if f in all_deg_set and f not in top15_set],
                key=_extract_mz_numeric,
            )
        else:
            # 全クラスタのDEGマーカー（従来通り）
            deg_genes = [str(r.get("gene", "")) for r in deg_data]
            deg_set = set(deg_genes)
            features = [f for f in features if f in deg_set]
    else:
        # 全m/zモード: m/zフィルタ適用済みリストがあればそれをベースにする
        if mz_filtered:
            features = mz_filtered

    # --- Top15 + 残り の特別表示モード ---
    if top15_ordered is not None:
        if not search_value:
            options = []
            for rank, f in enumerate(top15_ordered, 1):
                options.append(_make_option(f, rank=rank))
            if top15_ordered and rest_ordered:
                options.append({
                    "label": "──── その他の DEG マーカー ────",
                    "value": "__separator__",
                    "disabled": True,
                })
            for f in rest_ordered[:500 - len(top15_ordered) - 1]:
                options.append(_make_option(f))
            return options
        else:
            keyword = search_value.lower()
            options = []
            for rank, f in enumerate(top15_ordered, 1):
                if keyword in f.lower() or keyword in ann_map.get(f, "").lower():
                    options.append(_make_option(f, rank=rank))
            matched_rest = [
                f for f in rest_ordered
                if keyword in f.lower() or keyword in ann_map.get(f, "").lower()
            ]
            if options and matched_rest:
                options.append({
                    "label": "──── その他の DEG マーカー ────",
                    "value": "__separator__",
                    "disabled": True,
                })
            for f in matched_rest[:100]:
                options.append(_make_option(f))
            return options

    # --- 通常モード ---
    if not search_value:
        # 検索なしの場合は全件（最大500件）
        return [_make_option(f) for f in features[:500]]

    # 検索値でフィルタ（大文字小文字区別なし、アノテーションも検索対象）
    keyword = search_value.lower()
    filtered = [
        f for f in features
        if keyword in f.lower() or keyword in ann_map.get(f, "").lower()
    ]
    return [_make_option(f) for f in filtered[:100]]


# ---------------------------------------------------------------------------
# m/z 範囲フィルタ
# ---------------------------------------------------------------------------

@callback(
    Output("feature_mz_filtered_list", "data"),
    Input("apply_feature_mz_filter", "n_clicks"),
    [State("feature_mz_min", "value"),
     State("feature_mz_max", "value")],
    prevent_initial_call=True,
)
def apply_mz_filter(n_clicks, mz_min, mz_max):
    """m/z最小値・最大値で feature リストを絞り込み、Storeに保存"""
    features = _interactive_data.get("features_list")
    if not features:
        return None
    if mz_min is None and mz_max is None:
        return None  # フィルタなし → 全件に戻す

    filtered = []
    for f in features:
        # feature名から数値部分を抽出（例: "mz_123.456" → 123.456）
        match = re.search(r"(\d+\.?\d*)", f)
        if match:
            val = float(match.group(1))
            if mz_min is not None and val < mz_min:
                continue
            if mz_max is not None and val > mz_max:
                continue
        filtered.append(f)
    return filtered


@callback(
    Output("feature_select", "options", allow_duplicate=True),
    Input("feature_mz_filtered_list", "data"),
    prevent_initial_call=True,
)
def update_feature_options_on_mz_filter(mz_filtered):
    """m/zフィルタ適用後、ドロップダウンの選択肢を即時更新"""
    if mz_filtered is None:
        # フィルタ解除 → 全件に戻す
        features = _interactive_data.get("features_list")
        if not features:
            return []
        return [{"label": f, "value": f} for f in features[:500]]

    return [{"label": f, "value": f} for f in mz_filtered[:500]]


# ---------------------------------------------------------------------------
# 表示名ヘルパー
# ---------------------------------------------------------------------------

def _display_name(original: str, name_map: dict | None) -> str:
    """元のサンプル名をユーザー指定の表示名に変換する。マップが空なら元名をそのまま返す。"""
    if not name_map:
        return original
    return name_map.get(original, original)


def _compact_sci(v):
    """数値をコンパクトな指数表記に変換: 280000 → '2.8e5'"""
    if v == 0:
        return "0"
    exp = int(np.floor(np.log10(abs(v))))
    coeff = v / (10 ** exp)
    return f"{coeff:.1f}e{exp}"


def _format_plain_number(v):
    """数値を e 表記なしのプレーンな文字列で返す。
    0.00123 → '0.00123', 280000 → '280000', 3.5 → '3.5'"""
    if v == 0:
        return "0"
    if abs(v) >= 1:
        # 整数表示可能ならそうする
        if v == int(v):
            return str(int(v))
        return f"{v:.1f}"
    # 小数の場合、有効数字を維持しつつ e なし
    return f"{v:.6g}"


def _generate_umap_arrow_image():
    """参照画像と同じスタイルのL字型UMAP軸画像をbase64 PNGで生成する（キャッシュ付き）"""
    if hasattr(_generate_umap_arrow_image, "_cache"):
        return _generate_umap_arrow_image._cache

    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig_mpl, ax = plt.subplots(figsize=(1.6, 1.6), dpi=150)
    ax.set_xlim(-0.35, 1.15)
    ax.set_ylim(-0.35, 1.15)
    ax.set_aspect("equal")
    ax.axis("off")
    fig_mpl.patch.set_alpha(0.0)
    ax.patch.set_alpha(0.0)

    lw = 3.5
    hl = 0.08   # 矢印ヘッドの長さ
    hw = 0.06   # 矢印ヘッドの幅

    # 水平矢印（UMAP1）
    ax.annotate("", xy=(1.0, 0.0), xytext=(0.0, 0.0),
                arrowprops=dict(arrowstyle="->, head_length={}, head_width={}".format(hl * 3, hw * 3),
                                lw=lw, color="black"))
    # 垂直矢印（UMAP2）
    ax.annotate("", xy=(0.0, 1.0), xytext=(0.0, 0.0),
                arrowprops=dict(arrowstyle="->, head_length={}, head_width={}".format(hl * 3, hw * 3),
                                lw=lw, color="black"))

    # ラベル
    ax.text(1.0, -0.18, "UMAP1", fontsize=14, fontweight="bold",
            ha="center", va="top", color="black")
    ax.text(-0.18, 1.0, "UMAP2", fontsize=14, fontweight="bold",
            ha="center", va="center", rotation=90, color="black")

    buf = io.BytesIO()
    fig_mpl.savefig(buf, format="png", bbox_inches="tight",
                    transparent=True, pad_inches=0.05)
    plt.close(fig_mpl)
    buf.seek(0)
    b64 = base64.b64encode(buf.read()).decode()
    data_uri = f"data:image/png;base64,{b64}"

    _generate_umap_arrow_image._cache = data_uri
    return data_uri


def _add_umap_arrows(fig):
    """UMAPプロットの左下にL字型 UMAP1/UMAP2 軸画像を埋め込む"""
    arrow_src = _generate_umap_arrow_image()
    fig.add_layout_image(
        source=arrow_src,
        x=-0.02, y=-0.02,
        xref="paper", yref="paper",
        sizex=0.22, sizey=0.22,
        xanchor="left", yanchor="bottom",
        layer="above",
        opacity=1.0,
    )


# ---------------------------------------------------------------------------
# UMAP プロット — ヘルパー関数
# ---------------------------------------------------------------------------

def _build_umap_integrated_fig(df, color_by, highlight_clusters,
                                show_legend, show_labels, title=None,
                                marker_size=2, exclude_clusters=None,
                                label_size=14, saved_positions=None,
                                custom_colors=None, bg_opacity=0.1,
                                title_font_size=None):
    """統合UMAPのgo.Figureを生成（メイン/フルスクリーン共用）"""
    fig = go.Figure()

    # 除外クラスタのフィルタリング
    if exclude_clusters:
        exclude_set = set(str(c) for c in exclude_clusters)
        df = df[~df["Cluster"].astype(str).isin(exclude_set)]
        if df.empty:
            fig.add_annotation(text="全クラスタが除外されています", showarrow=False,
                               xref="paper", yref="paper", x=0.5, y=0.5)
            return fig

    color_map = _get_cluster_color_map(df["Cluster"], custom_colors)

    if highlight_clusters and len(highlight_clusters) > 0:
        highlight_set = set(str(c) for c in highlight_clusters)
        mask_bg = ~df["Cluster"].astype(str).isin(highlight_set)
        if mask_bg.any():
            fig.add_trace(go.Scattergl(
                x=df.loc[mask_bg, "UMAP_1"],
                y=df.loc[mask_bg, "UMAP_2"],
                mode="markers",
                marker=dict(size=max(1, marker_size - 1), color=HIGHLIGHT_GRAY, opacity=bg_opacity),
                name="Other", showlegend=False, hoverinfo="skip",
            ))
        for cl in highlight_clusters:
            mask = df["Cluster"].astype(str) == str(cl)
            if mask.any():
                fig.add_trace(go.Scattergl(
                    x=df.loc[mask, "UMAP_1"],
                    y=df.loc[mask, "UMAP_2"],
                    mode="markers",
                    marker=dict(size=marker_size + 1, color=color_map.get(str(cl), "#999999")),
                    name=f"Cluster {cl}",
                    text=df.loc[mask, "CellID"],
                    hovertemplate="Cluster: %{meta}<br>%{text}<extra></extra>",
                    meta=[str(cl)] * mask.sum(),
                ))
    else:
        color_col = color_by if color_by in df.columns else "Cluster"
        categories = sorted(df[color_col].unique(), key=_cluster_sort_key)
        cat_color_map = _get_cluster_color_map(categories, custom_colors)
        for cat in categories:
            mask = df[color_col] == cat
            rank = _cluster_sort_key(cat)[0] if str(cat).isdigit() else 1000
            fig.add_trace(go.Scattergl(
                x=df.loc[mask, "UMAP_1"],
                y=df.loc[mask, "UMAP_2"],
                mode="markers",
                marker=dict(size=marker_size, color=cat_color_map.get(str(cat), "#999999")),
                name=f"Cluster {cat}",
                legendrank=rank,
                text=df.loc[mask, "CellID"],
                hovertemplate=f"{color_col}: {cat}<br>" + "%{text}<extra></extra>",
            ))

    if show_labels:
        centroids = df.groupby("Cluster").agg(
            cx=("UMAP_1", "mean"), cy=("UMAP_2", "mean"),
        ).reset_index()
        centroids = centroids.sort_values(
            "Cluster", key=lambda col: col.map(_cluster_sort_key)
        )
        for _, row in centroids.iterrows():
            cl_str = str(row["Cluster"])
            pos = (saved_positions or {}).get(cl_str, {})
            fig.add_annotation(
                x=pos.get("x", row["cx"]),
                y=pos.get("y", row["cy"]),
                text=cl_str,
                showarrow=False,
                font=dict(size=label_size, color="black", family="Arial Black"),
            )

    layout_opts = dict(
        dragmode="select",
        showlegend=bool(show_legend),
        legend=dict(itemsizing="constant", font=dict(size=12), tracegroupgap=2),
        margin=dict(l=60, r=10,
                    t=max(40, (title_font_size or 14) + 15) if title else 30,
                    b=60),
        xaxis=dict(showgrid=False, showline=False, zeroline=False,
                   showticklabels=False, title=""),
        yaxis=dict(scaleanchor="x", showgrid=False, showline=False,
                   zeroline=False, showticklabels=False, title=""),
        plot_bgcolor="white",
    )
    if title:
        layout_opts["title"] = dict(
            text=title, font=dict(size=title_font_size or 14), x=0.5)
    fig.update_layout(**layout_opts)
    _add_umap_arrows(fig)
    return fig


def _build_umap_per_sample_graphs(df, color_map, highlight_clusters,
                                   show_labels, graph_height="300px",
                                   marker_size=2, exclude_clusters=None,
                                   label_size=11, saved_positions=None,
                                   show_legend=True, name_map=None,
                                   columns_per_row=0):
    """サンプル別UMAPのhtml.Divリストを生成（メイン/フルスクリーン共用）"""
    # 除外クラスタのフィルタリング
    if exclude_clusters:
        exclude_set = set(str(c) for c in exclude_clusters)
        df = df[~df["Cluster"].astype(str).isin(exclude_set)]
        if df.empty:
            return [html.Div("全クラスタが除外されています", className="text-muted small mt-2")]

    samples = sorted(df["Sample"].unique())
    if len(samples) <= 1:
        return [html.Div("サンプルが1つのみです", className="text-muted small mt-2")]

    graphs = []
    for s in samples:
        fig = go.Figure()
        mask_sample = df["Sample"] == s
        df_s = df[mask_sample]
        if highlight_clusters and len(highlight_clusters) > 0:
            hl_set = set(str(c) for c in highlight_clusters)
            mask_hl = df_s["Cluster"].astype(str).isin(hl_set)
            mask_bg_s = ~mask_hl
            if mask_bg_s.any():
                fig.add_trace(go.Scattergl(
                    x=df_s.loc[mask_bg_s, "UMAP_1"],
                    y=df_s.loc[mask_bg_s, "UMAP_2"],
                    mode="markers",
                    marker=dict(size=marker_size, color=HIGHLIGHT_GRAY, opacity=0.3),
                    name="Other", showlegend=False, hoverinfo="skip",
                ))
            for cl in highlight_clusters:
                mask_cl = df_s["Cluster"].astype(str) == str(cl)
                if mask_cl.any():
                    fig.add_trace(go.Scattergl(
                        x=df_s.loc[mask_cl, "UMAP_1"],
                        y=df_s.loc[mask_cl, "UMAP_2"],
                        mode="markers",
                        marker=dict(size=marker_size + 1, color=color_map.get(str(cl), "#999999")),
                        name=f"Cluster {cl}", showlegend=False,
                        legendgroup=f"Cluster {cl}",
                    ))
        else:
            for cl in sorted(df_s["Cluster"].unique(), key=_cluster_sort_key):
                mask_cl = df_s["Cluster"] == cl
                fig.add_trace(go.Scattergl(
                    x=df_s.loc[mask_cl, "UMAP_1"],
                    y=df_s.loc[mask_cl, "UMAP_2"],
                    mode="markers",
                    marker=dict(size=marker_size, color=color_map.get(str(cl), "#999999")),
                    name=str(cl), showlegend=False,
                    legendgroup=f"Cluster {cl}",
                ))

        # 凡例用ダミートレース（全クラスタ共通で統一凡例を表示）
        if show_legend:
            for cl in sorted(df["Cluster"].unique(), key=_cluster_sort_key):
                rank = _cluster_sort_key(cl)[0] if str(cl).isdigit() else 1000
                fig.add_trace(go.Scattergl(
                    x=[None], y=[None], mode="markers",
                    marker=dict(size=10, color=color_map.get(str(cl), "#999999")),
                    name=f"Cluster {cl}", showlegend=True, legendrank=rank,
                    legendgroup=f"Cluster {cl}",
                ))

        if show_labels:
            sample_pos = (saved_positions or {}).get(s, {})
            centroids = df_s.groupby("Cluster").agg(
                cx=("UMAP_1", "mean"), cy=("UMAP_2", "mean"),
            ).reset_index()
            centroids = centroids.sort_values(
                "Cluster", key=lambda col: col.map(_cluster_sort_key)
            )
            for _, row in centroids.iterrows():
                cl_str = str(row["Cluster"])
                pos = sample_pos.get(cl_str, {})
                fig.add_annotation(
                    x=pos.get("x", row["cx"]),
                    y=pos.get("y", row["cy"]),
                    text=cl_str,
                    showarrow=False,
                    font=dict(size=label_size, color="black", family="Arial Black"),
                )

        display_s = _display_name(s, name_map)
        fig.update_layout(
            margin=dict(l=50, r=10, t=30, b=50),
            title=dict(text=display_s, font=dict(size=12), x=0.5),
            xaxis=dict(showgrid=False, showline=False, zeroline=False,
                       showticklabels=False, title=""),
            yaxis=dict(scaleanchor="x", showgrid=False, showline=False,
                       zeroline=False, showticklabels=False, title=""),
            plot_bgcolor="white",
            showlegend=bool(show_legend),
            legend=dict(itemsizing="constant", font=dict(size=9), tracegroupgap=1),
        )
        _add_umap_arrows(fig)

        cfg = dict(_UMAP_PER_SAMPLE_CONFIG)
        cfg["toImageButtonOptions"] = dict(cfg["toImageButtonOptions"],
                                           filename=f"UMAP_{display_s}")
        if columns_per_row:
            n_cols = columns_per_row
            gap_total = (n_cols - 1) * 15
            flex_basis = f"calc({100 / n_cols:.2f}% - {gap_total / n_cols:.1f}px)"
            min_w = "0"
        else:
            n_cols = len(samples)
            flex_basis = f"{max(20, 90 // n_cols)}%"
            min_w = "300px"
        graphs.append(
            html.Div(
                style={"flex": f"1 1 {flex_basis}", "minWidth": min_w,
                        "border": "1px solid #dee2e6", "borderRadius": "6px",
                        "padding": "5px", "backgroundColor": "#fff"},
                children=[
                    dcc.Graph(id={"type": "umap_per_sample_graph", "index": s},
                              figure=fig, style={"height": graph_height}, config=cfg),
                ],
            )
        )
    return graphs


# ---------------------------------------------------------------------------
# UMAP プロット — コールバック
# ---------------------------------------------------------------------------

def _get_merged_label_positions(accumulated_positions=None):
    """JSON ファイル + 蓄積 Store からマージしたラベル位置を返す。

    蓄積データは JSON より新しいため、蓄積データで JSON をオーバーライドする。
    """
    all_pos = _load_label_positions()
    acc = accumulated_positions or {}
    for section in ("umap_integrated", "umap_per_sample", "spatial"):
        acc_section = acc.get(section)
        if not acc_section:
            continue
        saved_section = all_pos.get(section, {})
        if section == "umap_integrated":
            _merge_label_positions(saved_section, acc_section)
        else:
            for sample_name, pos_dict in acc_section.items():
                sample_saved = saved_section.get(sample_name, {})
                _merge_label_positions(sample_saved, pos_dict)
                saved_section[sample_name] = sample_saved
        all_pos[section] = saved_section
    return all_pos


@callback(
    Output("interactive_umap_plot", "figure"),
    [Input("umap_color_by", "value"),
     Input("umap_highlight_cluster", "value"),
     Input("umap_show_legend", "value"),
     Input("umap_show_labels", "value"),
     Input("umap_display_mode", "value"),
     Input("umap_marker_size", "value"),
     Input("umap_exclude_cluster", "value"),
     Input("umap_label_size", "value"),
     Input("seurat_rds_path_store", "data"),
     Input("fullscreen_closed_trigger", "data"),
     Input("custom_color_map_store", "data")],
    State("accumulated_label_positions", "data"),
)
def update_umap_plot(color_by, highlight_clusters, show_legend, show_labels,
                     display_mode, marker_size, exclude_clusters, label_size,
                     rds_path, _fs_trigger, custom_colors, accumulated_positions):
    if display_mode == "per_sample":
        return go.Figure()
    df = _interactive_data.get("plot_data")
    if df is None:
        return go.Figure()
    all_pos = _get_merged_label_positions(accumulated_positions)
    return _build_umap_integrated_fig(df, color_by, highlight_clusters,
                                       show_legend, show_labels,
                                       marker_size=marker_size or 2,
                                       exclude_clusters=exclude_clusters,
                                       label_size=label_size or 14,
                                       saved_positions=all_pos.get("umap_integrated"),
                                       custom_colors=custom_colors)


# ---------------------------------------------------------------------------
# UMAP 統合/サンプル別 表示切替
# ---------------------------------------------------------------------------

@callback(
    Output("umap_integrated_wrapper", "style"),
    Input("umap_display_mode", "value"),
)
def toggle_umap_integrated_visibility(mode):
    """「統合」選択時のみ統合UMAPを表示、「サンプル別」では非表示"""
    if mode == "per_sample":
        return {"display": "none"}
    return {"display": "block"}


# ---------------------------------------------------------------------------
# サンプル別 UMAP 表示
# ---------------------------------------------------------------------------

_UMAP_PER_SAMPLE_CONFIG = {
    "scrollZoom": True,
    "edits": {"annotationPosition": True},
    "toImageButtonOptions": {"format": "png", "scale": 3},
}


@callback(
    Output("umap_per_sample_container", "children"),
    [Input("umap_display_mode", "value"),
     Input("umap_highlight_cluster", "value"),
     Input("umap_show_labels", "value"),
     Input("umap_marker_size", "value"),
     Input("umap_exclude_cluster", "value"),
     Input("umap_label_size", "value"),
     Input("seurat_rds_path_store", "data"),
     Input("umap_show_legend", "value"),
     Input("sample_name_map_store", "data"),
     Input("fullscreen_closed_trigger", "data"),
     Input("custom_color_map_store", "data"),
     Input("umap_columns_per_row", "value")],
    State("accumulated_label_positions", "data"),
)
def update_umap_per_sample(display_mode, highlight_clusters, show_labels,
                            marker_size, exclude_clusters, label_size, rds_path,
                            show_legend, name_map, _fs_trigger, custom_colors,
                            columns_per_row, accumulated_positions):
    """表示モード「サンプル別」の場合、各サンプルのUMAPを並列表示"""
    if display_mode != "per_sample":
        return ""
    df = _interactive_data.get("plot_data")
    if df is None:
        return ""
    color_map = _get_cluster_color_map(df["Cluster"], custom_colors)
    all_pos = _get_merged_label_positions(accumulated_positions)
    graphs = _build_umap_per_sample_graphs(df, color_map, highlight_clusters,
                                            show_labels, graph_height="300px",
                                            marker_size=marker_size or 2,
                                            exclude_clusters=exclude_clusters,
                                            label_size=label_size or 11,
                                            saved_positions=all_pos.get("umap_per_sample"),
                                            show_legend=bool(show_legend),
                                            name_map=name_map,
                                            columns_per_row=columns_per_row or 0)
    return html.Div(
        style={"display": "flex", "flexWrap": "wrap", "gap": "15px", "marginTop": "10px"},
        children=graphs,
    )


# ---------------------------------------------------------------------------
# クラスタ統計テーブル
# ---------------------------------------------------------------------------

@callback(
    Output("cluster_stats_table", "data"),
    Input("seurat_rds_path_store", "data"),
)
def update_cluster_stats(rds_path):
    df = _interactive_data.get("plot_data")
    if df is None:
        return []

    total = len(df)
    stats = df["Cluster"].value_counts()
    stats = stats.reindex(sorted(stats.index, key=_cluster_sort_key))
    return [
        {"Cluster": str(c), "Pixels": int(n), "Percent": f"{n / total * 100:.1f}%"}
        for c, n in stats.items()
    ]


# ---------------------------------------------------------------------------
# クラスタ情報テキスト
# ---------------------------------------------------------------------------

@callback(
    Output("cluster_info_text", "children"),
    [Input("cluster_stats_table", "selected_rows"),
     Input("umap_highlight_cluster", "value")],
    State("cluster_stats_table", "data"),
)
def update_cluster_info(selected_rows, highlight, table_data):
    df = _interactive_data.get("plot_data")
    if df is None:
        return "データを読み込んでください"

    cluster_id = None
    if selected_rows and table_data:
        cluster_id = table_data[selected_rows[0]].get("Cluster")
    elif highlight and len(highlight) == 1:
        cluster_id = str(highlight[0])

    if cluster_id is None:
        meta = _interactive_data.get("meta", {})
        return (
            f"Total cells: {meta.get('n_cells', '?')}\n"
            f"Clusters: {meta.get('n_clusters', '?')}\n"
            f"Samples: {', '.join(meta.get('samples', []))}"
        )

    mask = df["Cluster"].astype(str) == str(cluster_id)
    n = mask.sum()
    total = len(df)
    samples = df.loc[mask, "Sample"].value_counts()
    sample_info = "\n".join(f"  {s}: {c} pixels" for s, c in samples.items())

    return f"Cluster {cluster_id}: {n} pixels ({n / total * 100:.1f}%)\n{sample_info}"


# ---------------------------------------------------------------------------
# Spatial Mapping プロット (UMAP選択連動)
# ---------------------------------------------------------------------------

def _transform_coords(x, y, angle_deg, flip_h=False, flip_v=False):
    """中心基準で座標を反転+2D回転"""
    cx, cy = x.mean(), y.mean()
    # 反転（回転の前に適用）
    if flip_h:
        x = 2 * cx - x
    if flip_v:
        y = 2 * cy - y
    # 回転
    if angle_deg == 0:
        return x, y
    # 反転後の中心を再計算
    cx, cy = x.mean(), y.mean()
    rad = np.radians(angle_deg)
    cos_a, sin_a = np.cos(rad), np.sin(rad)
    x_rot = cos_a * (x - cx) - sin_a * (y - cy) + cx
    y_rot = sin_a * (x - cx) + cos_a * (y - cy) + cy
    return x_rot, y_rot


def _calc_zero_gap_marker_size(plot_x, plot_y, render_height=310):
    """点間距離ゼロ（隣接点が接する）のマーカーサイズを計算（scale_factor=1.0）"""
    if len(plot_x) <= 1:
        return 4
    sorted_ux = np.sort(np.unique(plot_x))
    if len(sorted_ux) <= 1:
        return 4
    min_spacing = float(np.min(np.diff(sorted_ux)))
    x_range = float(plot_x.max() - plot_x.min())
    y_range = float(plot_y.max() - plot_y.min()) if len(plot_y) > 1 else 1.0
    if y_range > 0 and x_range > 0:
        effective_w = render_height * (x_range / y_range)
        return max(2, round(min_spacing * effective_w / x_range))
    return 4


def _create_single_spatial_fig(df_sample, color_map, highlight_clusters,
                               selected_cell_ids, rotation_deg=0,
                               show_labels=False, flip_h=False, flip_v=False,
                               title=None, embed_legend=False,
                               cluster_to_idx=None, discrete_cscale=None,
                               marker_size=4, exclude_clusters=None,
                               label_size=10, saved_positions=None,
                               title_font_size=None, render_height=None):
    """単一サンプルのSpatial Mapping figureを生成"""
    # 除外クラスタのフィルタリング
    if exclude_clusters:
        exclude_set = set(str(c) for c in exclude_clusters)
        df_sample = df_sample[~df_sample["Cluster"].astype(str).isin(exclude_set)]
        if df_sample.empty:
            fig = go.Figure()
            fig.add_annotation(text="全クラスタが除外されています", showarrow=False,
                               xref="paper", yref="paper", x=0.5, y=0.5)
            return fig
    fig = go.Figure()

    # 座標の取得と変換適用（反転+回転）
    raw_x = df_sample["SpatialX"].values
    raw_y = -df_sample["SpatialY"].values  # Y軸反転
    plot_x, plot_y = _transform_coords(raw_x, raw_y, rotation_deg,
                                        flip_h=flip_h, flip_v=flip_v)

    # marker_size=0 の場合、データ密度ベースで自動計算（点間距離ゼロ）
    if marker_size <= 0 and len(plot_x) > 1:
        sorted_ux = np.sort(np.unique(plot_x))
        if len(sorted_ux) > 1:
            min_spacing = float(np.min(np.diff(sorted_ux)))
            x_range = float(plot_x.max() - plot_x.min())
            y_range = float(plot_y.max() - plot_y.min()) if len(plot_y) > 1 else 1.0
            # scaleanchor="x" のため、描画幅は高さ×アスペクト比で決定
            # render_height 指定時はそれを使用、未指定時はWebデフォルト310px
            effective_h = render_height or 310
            _scale_factor = 1.0
            if y_range > 0 and x_range > 0:
                effective_w = effective_h * (x_range / y_range)
                marker_size = max(2, round(min_spacing * effective_w / x_range * _scale_factor))
            else:
                marker_size = 4
        else:
            marker_size = 4
    elif marker_size <= 0:
        marker_size = 4

    if selected_cell_ids:
        mask_selected = df_sample["CellID"].isin(selected_cell_ids).values
        mask_bg = ~mask_selected
        if mask_bg.any():
            if "TotalCount" in df_sample.columns:
                tc_values = df_sample["TotalCount"].values[mask_bg]
                bg_marker = dict(size=marker_size, symbol="square", color=tc_values, colorscale="Greys",
                                 opacity=0.5, showscale=False)
                bg_name = "TIC"
            else:
                bg_marker = dict(size=marker_size, symbol="square", color=HIGHLIGHT_GRAY, opacity=0.2)
                bg_name = "Other"
            fig.add_trace(go.Scattergl(
                x=plot_x[mask_bg],
                y=plot_y[mask_bg],
                mode="markers",
                marker=bg_marker,
                name=bg_name, showlegend=False, hoverinfo="skip",
            ))
        if mask_selected.any():
            fig.add_trace(go.Scattergl(
                x=plot_x[mask_selected],
                y=plot_y[mask_selected],
                mode="markers",
                marker=dict(size=marker_size + 1, symbol="square", color="red"),
                name=f"Selected ({mask_selected.sum()})",
            ))
    elif highlight_clusters and len(highlight_clusters) > 0:
        highlight_set = set(str(c) for c in highlight_clusters)
        # 非ハイライトクラスタをTIC or 灰色で描画
        mask_bg = ~df_sample["Cluster"].astype(str).isin(highlight_set)
        if mask_bg.values.any():
            if "TotalCount" in df_sample.columns:
                tc_values = df_sample["TotalCount"].values[mask_bg.values]
                bg_marker = dict(size=marker_size, symbol="square", color=tc_values, colorscale="Greys",
                                 opacity=0.5, showscale=False)
                bg_name = "TIC"
            else:
                bg_marker = dict(size=marker_size, symbol="square", color=HIGHLIGHT_GRAY, opacity=0.2)
                bg_name = "Other"
            fig.add_trace(go.Scattergl(
                x=plot_x[mask_bg.values],
                y=plot_y[mask_bg.values],
                mode="markers",
                marker=bg_marker,
                name=bg_name, showlegend=False, hoverinfo="skip",
            ))
        # ハイライトクラスタを色付きで描画
        for cl in sorted(highlight_clusters, key=lambda x: _cluster_sort_key(x), reverse=True):
            mask = (df_sample["Cluster"].astype(str) == str(cl)).values
            if mask.any():
                fig.add_trace(go.Scattergl(
                    x=plot_x[mask],
                    y=plot_y[mask],
                    mode="markers",
                    marker=dict(size=marker_size + 1, symbol="square", color=color_map.get(str(cl), "#999999")),
                    name=f"Cluster {cl}",
                    legendgroup=f"Cluster {cl}",
                ))
    else:
        if embed_legend:
            # 凡例リンク用: クラスタ別個別トレース（legendgroup でダミーと連動）
            for cl in sorted(df_sample["Cluster"].unique(), key=_cluster_sort_key):
                mask = (df_sample["Cluster"].astype(str) == str(cl)).values
                if mask.any():
                    fig.add_trace(go.Scattergl(
                        x=plot_x[mask], y=plot_y[mask], mode="markers",
                        marker=dict(size=marker_size, symbol="square",
                                    color=color_map.get(str(cl), "#999999")),
                        text=[f"Cluster {cl}"] * int(mask.sum()),
                        hovertemplate="%{text}<extra></extra>",
                        name=f"Cluster {cl}", showlegend=False,
                        legendgroup=f"Cluster {cl}",
                    ))
        elif cluster_to_idx is not None and discrete_cscale is not None:
            # 数値インデックス + discrete colorscale 方式
            n_clusters = max(len(cluster_to_idx), 1)
            point_values = np.array(
                [cluster_to_idx.get(str(cl), 0) for cl in df_sample["Cluster"]]
            )
            point_normalized = (point_values + 0.5) / n_clusters
            fig.add_trace(go.Scattergl(
                x=plot_x, y=plot_y, mode="markers",
                marker=dict(
                    size=marker_size,
                    symbol="square",
                    color=point_normalized,
                    colorscale=discrete_cscale,
                    cmin=0, cmax=1,
                    showscale=False,
                ),
                text=[f"Cluster {cl}" for cl in df_sample["Cluster"]],
                hovertemplate="%{text}<extra></extra>",
                showlegend=False,
            ))
        else:
            # フォールバック: HEX文字列配列方式
            point_colors = [color_map.get(str(cl), "#999999") for cl in df_sample["Cluster"]]
            fig.add_trace(go.Scattergl(
                x=plot_x, y=plot_y, mode="markers",
                marker=dict(size=marker_size, symbol="square", color=point_colors),
                text=[f"Cluster {cl}" for cl in df_sample["Cluster"]],
                hovertemplate="%{text}<extra></extra>",
                showlegend=False,
            ))
        # 凡例用ダミートレース（大きいマーカーで見やすく）
        if embed_legend:
            for cl in sorted(df_sample["Cluster"].unique(), key=_cluster_sort_key):
                rank = _cluster_sort_key(cl)[0] if str(cl).isdigit() else 1000
                fig.add_trace(go.Scattergl(
                    x=[None], y=[None],
                    mode="markers",
                    marker=dict(size=10, symbol="square", color=color_map.get(str(cl), "#999999")),
                    name=f"Cluster {cl}",
                    showlegend=True,
                    legendrank=rank,
                    legendgroup=f"Cluster {cl}",
                ))

    # クラスタ番号ラベル
    if show_labels:
        for cl in sorted(df_sample["Cluster"].unique(), key=_cluster_sort_key):
            mask = (df_sample["Cluster"] == cl).values
            if mask.any():
                cx_default = plot_x[mask].mean()
                cy_default = plot_y[mask].mean()
                cl_str = str(cl)
                pos = (saved_positions or {}).get(cl_str, {})
                fig.add_annotation(
                    x=pos.get("x", cx_default),
                    y=pos.get("y", cy_default),
                    text=cl_str,
                    showarrow=False,
                    font=dict(size=label_size, color="black"),
                )

    layout_opts = dict(
        xaxis=dict(showgrid=False, showline=False, zeroline=False,
                   showticklabels=False, title="", visible=False),
        yaxis=dict(scaleanchor="x", showgrid=False, showline=False, zeroline=False,
                   showticklabels=False, title="", visible=False),
        margin=dict(l=10, r=10, t=max(30, (title_font_size or 14) + 15) if title else 10, b=10),
        plot_bgcolor="white",
        showlegend=embed_legend,
        legend=dict(itemsizing="constant", font=dict(size=12), tracegroupgap=2),
    )
    if title:
        layout_opts["title"] = dict(text=title, font=dict(size=title_font_size or 14), x=0.5)
    fig.update_layout(**layout_opts)
    return fig


_SPATIAL_IMG_CONFIG = {
    "scrollZoom": True,
    "edits": {"annotationPosition": True},
    "toImageButtonOptions": {
        "format": "png", "scale": 3,
    },
}


# ---------------------------------------------------------------------------
# Spatial サンプル別コントロール生成
# ---------------------------------------------------------------------------

@callback(
    Output("spatial_controls_container", "children"),
    Input("seurat_rds_path_store", "data"),
    [State("spatial_rotation_store", "data"),
     State("sample_name_map_store", "data"),
     State("custom_color_map_store", "data")],
)
def create_spatial_controls(rds_path, rotation_store, name_map, custom_color_map):
    """データ読み込み後、サンプル別の回転/反転 + サンプル名変更コントロールを生成"""
    df = _interactive_data.get("plot_data")
    if df is None or "SpatialX" not in df.columns:
        return ""
    if not rotation_store:
        rotation_store = {}
    if not name_map:
        name_map = {}
    if not custom_color_map:
        custom_color_map = {}

    # ==================== 回転/反転 ====================
    samples = sorted(df["Sample"].unique())
    sample_options = [{"label": s, "value": s} for s in samples]
    first_sample = samples[0] if samples else None

    sample_blocks = []
    for i, s in enumerate(samples):
        transform = rotation_store.get(
            s, rotation_store.get("__all__", {"angle": 0, "flip_h": False, "flip_v": False}))
        if isinstance(transform, (int, float)):
            transform = {"angle": int(transform), "flip_h": False, "flip_v": False}

        display_s = _display_name(s, name_map)
        is_first = (i == 0)
        block = html.Div(
            id={"type": "sample_block", "index": s},
            style={"padding": "4px 8px",
                   "display": "block" if is_first else "none"},
            children=[
                dbc.Row(className="align-items-center mb-1", children=[
                    dbc.Col(width=3, children=[
                        html.Label(s, className="fw-bold small mb-0",
                                   style={"whiteSpace": "nowrap", "overflow": "hidden",
                                          "textOverflow": "ellipsis"}),
                    ]),
                    dbc.Col(width=4, children=[
                        dbc.Input(
                            id={"type": "sample_rename_input", "index": s},
                            value=display_s if display_s != s else "",
                            placeholder=s,
                            size="sm", debounce=True,
                        ),
                    ]),
                    dbc.Col(width=5, children=[
                        html.Div(className="d-flex gap-2 align-items-center", children=[
                            dbc.Checkbox(
                                id={"type": "per_sample_flip_h", "index": s},
                                label="↔ 左右", value=transform.get("flip_h", False),
                            ),
                            dbc.Checkbox(
                                id={"type": "per_sample_flip_v", "index": s},
                                label="↕ 上下", value=transform.get("flip_v", False),
                            ),
                        ]),
                    ]),
                ]),
                dcc.Slider(
                    id={"type": "per_sample_rotation", "index": s},
                    min=0, max=270, step=90,
                    value=transform.get("angle", 0),
                    marks={0: "0°", 90: "90°", 180: "180°", 270: "270°"},
                ),
            ],
        )
        sample_blocks.append(block)

    rotation_section = [
        dbc.Select(
            id="spatial_sample_selector",
            options=sample_options,
            value=first_sample,
            size="sm",
            className="mb-2",
        ),
        *sample_blocks,
    ]

    # ==================== クラスタ色変更 ====================
    clusters = sorted(df["Cluster"].unique(), key=_cluster_sort_key)
    cluster_options = [{"label": f"Cluster {c}", "value": str(c)} for c in clusters]
    first_cluster = str(clusters[0]) if clusters else None

    # 現在使用中の色マップを構築（デフォルト + カスタム）
    current_cmap = _get_cluster_color_map(df["Cluster"], custom_color_map)
    # 各クラスタが使用中の色 → {色: クラスタ} のマップ
    color_usage = {}
    for cl_key, col_val in current_cmap.items():
        upper_col = col_val.upper()
        if upper_col not in color_usage:
            color_usage[upper_col] = cl_key

    cluster_blocks = []
    for idx, cl in enumerate(clusters):
        cl_str = str(cl)
        default_color = current_cmap.get(cl_str, "#999999")
        is_first = (idx == 0)

        swatches = []
        for pc in CLUSTER_PRESET_COLORS:
            # 他のクラスタで使用中ならグレーアウト
            owner = color_usage.get(pc.upper())
            used_by_other = (owner is not None and owner != cl_str)
            swatch_style = {
                "width": "18px", "height": "18px",
                "backgroundColor": pc,
                "border": "2px solid #aaa",
                "borderRadius": "3px",
                "display": "inline-block",
            }
            if used_by_other:
                swatch_style.update({
                    "opacity": "0.25",
                    "cursor": "not-allowed",
                    "pointerEvents": "none",
                })
            else:
                swatch_style["cursor"] = "pointer"

            swatches.append(
                html.Div(
                    style=swatch_style,
                    id={"type": "cluster_color_swatch",
                        "index": cl_str, "color": pc},
                    n_clicks=0,
                )
            )

        block = html.Div(
            id={"type": "cluster_block", "index": cl_str},
            style={"display": "block" if is_first else "none"},
            className="mb-2",
            children=[
                html.Label(f"Cluster {cl}", className="small mb-1 fw-bold"),
                html.Div(
                    style={"display": "flex", "alignItems": "center", "gap": "6px"},
                    children=[
                        html.Div(
                            style={"display": "flex", "flexWrap": "wrap", "gap": "3px"},
                            children=swatches,
                        ),
                        dbc.Input(
                            type="color",
                            id={"type": "cluster_color_picker", "index": cl_str},
                            value=default_color,
                            style={"width": "32px", "height": "24px", "padding": "1px",
                                   "border": "1px solid #ccc", "cursor": "pointer",
                                   "flexShrink": "0"},
                        ),
                    ],
                ),
            ],
        )
        cluster_blocks.append(block)

    color_section = [
        dbc.Select(
            id="spatial_cluster_selector",
            options=cluster_options,
            value=first_cluster,
            size="sm",
            className="mb-2",
        ),
        *cluster_blocks,
    ]

    return dbc.Accordion(
        [
            dbc.AccordionItem(title="回転/反転", children=rotation_section),
            dbc.AccordionItem(title="クラスタ色変更", children=color_section),
        ],
        start_collapsed=True,
        flush=True,
        always_open=True,
        style={"marginBottom": "8px"},
    )


# ---------------------------------------------------------------------------
# サンプル/クラスタ ブロック表示切替
# ---------------------------------------------------------------------------

@callback(
    Output({"type": "sample_block", "index": ALL}, "style"),
    Input("spatial_sample_selector", "value"),
    prevent_initial_call=True,
)
def toggle_sample_rotation_visibility(selected):
    """ドロップダウンで選択されたサンプルのみ表示"""
    styles = []
    for item in ctx.outputs_list:
        idx = item["id"]["index"]
        vis = "block" if idx == selected else "none"
        styles.append({"padding": "4px 8px", "display": vis})
    return styles


@callback(
    Output({"type": "cluster_block", "index": ALL}, "style"),
    Input("spatial_cluster_selector", "value"),
    prevent_initial_call=True,
)
def toggle_cluster_color_visibility(selected):
    """ドロップダウンで選択されたクラスタのみ表示"""
    styles = []
    for item in ctx.outputs_list:
        idx = item["id"]["index"]
        vis = "block" if idx == selected else "none"
        styles.append({"display": vis})
    return styles


# ---------------------------------------------------------------------------
# スウォッチの使用済み色グレーアウト
# ---------------------------------------------------------------------------

@callback(
    Output({"type": "cluster_color_swatch", "index": ALL, "color": ALL}, "style"),
    Input("custom_color_map_store", "data"),
    prevent_initial_call=True,
)
def update_swatch_disabled_state(custom_colors):
    """色マップ変更時に、他クラスタで使用中のスウォッチをグレーアウトする"""
    if not custom_colors:
        custom_colors = {}
    df = _interactive_data.get("plot_data")
    if df is None:
        raise PreventUpdate

    # 現在の全色マップ（デフォルト＋カスタム）
    current_cmap = _get_cluster_color_map(df["Cluster"], custom_colors)
    # {色(大文字) → クラスタID} の逆引き
    color_usage = {}
    for cl_key, col_val in current_cmap.items():
        upper_col = col_val.upper()
        if upper_col not in color_usage:
            color_usage[upper_col] = cl_key

    styles = []
    for item in ctx.outputs_list:
        swatch_cluster = item["id"]["index"]
        swatch_color = item["id"]["color"]

        base_style = {
            "width": "18px", "height": "18px",
            "backgroundColor": swatch_color,
            "border": "2px solid #aaa",
            "borderRadius": "3px",
            "display": "inline-block",
        }

        owner = color_usage.get(swatch_color.upper())
        if owner is not None and owner != swatch_cluster:
            base_style.update({
                "opacity": "0.25",
                "cursor": "not-allowed",
                "pointerEvents": "none",
            })
        else:
            base_style["cursor"] = "pointer"

        styles.append(base_style)
    return styles


# ---------------------------------------------------------------------------
# クラスタ色 Store 管理 (パターンマッチング)
# ---------------------------------------------------------------------------

@callback(
    [Output("custom_color_map_store", "data"),
     Output({"type": "cluster_color_picker", "index": ALL}, "value")],
    [Input({"type": "cluster_color_picker", "index": ALL}, "value"),
     Input({"type": "cluster_color_swatch", "index": ALL, "color": ALL}, "n_clicks")],
    State("custom_color_map_store", "data"),
    prevent_initial_call=True,
)
def update_custom_color_map(picker_values, swatch_clicks, current_store):
    """カラーピッカーまたはスウォッチクリックでカスタム色マップStoreを更新する"""
    current_store = current_store or {}

    # トリガーされたコンポーネントの判定
    triggered = ctx.triggered_id
    if not triggered:
        raise PreventUpdate

    # カラーピッカーのIDリストを取得（Output用の順序と一致させる）
    picker_ids = ctx.inputs_list[0]
    picker_cluster_order = [item["id"]["index"] for item in picker_ids]

    if isinstance(triggered, dict) and triggered.get("type") == "cluster_color_swatch":
        # --- スウォッチがクリックされた場合 ---
        cl = str(triggered["index"])
        color = triggered["color"]
        current_store[cl] = color
        # カラーピッカーの表示色も同期
        updated_picker_values = []
        for c_id in picker_cluster_order:
            if c_id == cl:
                updated_picker_values.append(color)
            else:
                # 既存のピッカー値を維持
                idx = picker_cluster_order.index(c_id)
                updated_picker_values.append(picker_values[idx])
        _save_interactive_settings("custom_color_map", current_store)
        return current_store, updated_picker_values
    else:
        # --- カラーピッカーが変更された場合 ---
        custom = {}
        for item in picker_ids:
            cl = item["id"]["index"]
            val = item.get("value")
            if val:
                custom[str(cl)] = val
        _save_interactive_settings("custom_color_map", custom)
        return custom, picker_values


# ---------------------------------------------------------------------------
# Spatial 回転 Store 管理 (パターンマッチング)
# ---------------------------------------------------------------------------

@callback(
    Output("spatial_rotation_store", "data"),
    [Input({"type": "per_sample_rotation", "index": ALL}, "value"),
     Input({"type": "per_sample_flip_h", "index": ALL}, "value"),
     Input({"type": "per_sample_flip_v", "index": ALL}, "value")],
    State("spatial_rotation_store", "data"),
    prevent_initial_call=True,
)
def update_rotation_store_from_per_sample(rotations, flip_hs, flip_vs, current_store):
    """各プロットのコントロール変更時に Store を更新"""
    if current_store is None:
        current_store = {}

    triggered = ctx.triggered_id
    if not triggered or not isinstance(triggered, dict):
        return no_update

    sample_name = triggered["index"]

    # 現在のサンプルの設定を取得
    transform = current_store.get(sample_name, {"angle": 0, "flip_h": False, "flip_v": False})
    if isinstance(transform, (int, float)):
        transform = {"angle": int(transform), "flip_h": False, "flip_v": False}

    # ctx.inputs_list から現在の全値を取得してトリガーサンプルの値を特定
    all_inputs = ctx.inputs_list
    rotation_inputs = all_inputs[0]
    flip_h_inputs = all_inputs[1]
    flip_v_inputs = all_inputs[2]

    for r_input in rotation_inputs:
        if r_input.get("id", {}).get("index") == sample_name:
            transform["angle"] = r_input.get("value", 0) or 0
    for fh_input in flip_h_inputs:
        if fh_input.get("id", {}).get("index") == sample_name:
            transform["flip_h"] = bool(fh_input.get("value", False))
    for fv_input in flip_v_inputs:
        if fv_input.get("id", {}).get("index") == sample_name:
            transform["flip_v"] = bool(fv_input.get("value", False))

    current_store[sample_name] = transform
    _save_interactive_settings("spatial_rotation", current_store)
    return current_store


# ---------------------------------------------------------------------------
# サンプル名マップ管理 (パターンマッチング)
# ---------------------------------------------------------------------------

@callback(
    Output("sample_name_map_store", "data"),
    [Input({"type": "sample_rename_input", "index": ALL}, "value"),
     Input({"type": "umap_sample_rename_input", "index": ALL}, "value")],
    prevent_initial_call=True,
)
def update_sample_name_map(spatial_values, umap_values):
    """サンプル名変更入力（Spatial側 + UMAP側）から表示名マッピングを更新。
    両側の値をマージし、トリガーされた側の値を高優先とする。"""
    triggered = ctx.triggered_id
    triggered_type = triggered.get("type", "") if isinstance(triggered, dict) else ""

    name_map = {}

    def _collect(inputs_list_idx):
        result = {}
        for inp in ctx.inputs_list[inputs_list_idx]:
            original = inp.get("id", {}).get("index", "")
            display_val = inp.get("value", "") or ""
            display_val = display_val.strip()
            if display_val and display_val != original:
                result[original] = display_val
        return result

    if triggered_type == "umap_sample_rename_input":
        # 非トリガー側（Spatial）を先に読み込み（低優先）
        name_map.update(_collect(0))
        # トリガー側（UMAP）で上書き（高優先）
        name_map.update(_collect(1))
    else:
        # 非トリガー側（UMAP）を先に読み込み（低優先）
        name_map.update(_collect(1))
        # トリガー側（Spatial）で上書き（高優先）
        name_map.update(_collect(0))

    # フルスクリーン用にも参照できるようモジュール変数にも保持
    _interactive_data["_name_map"] = name_map
    _save_interactive_settings("sample_name_map", name_map)
    return name_map


@callback(
    [Output("interactive_sample", "options", allow_duplicate=True),
     Output("feature_sample_select", "options", allow_duplicate=True)],
    Input("sample_name_map_store", "data"),
    prevent_initial_call=True,
)
def update_sample_dropdown_labels(name_map):
    """サンプル名変更時にSpatial Mapping & Feature Plotサンプルドロップダウンのラベルを更新"""
    df = _interactive_data.get("plot_data")
    if df is None:
        return no_update, no_update
    if not name_map:
        name_map = {}
    samples = sorted(df["Sample"].unique())
    opts = [{"label": _display_name(s, name_map), "value": s} for s in samples]
    return opts, opts


# ---------------------------------------------------------------------------
# UMAP側サンプル名変更コントロール生成
# ---------------------------------------------------------------------------

@callback(
    Output("umap_name_controls_container", "children"),
    Input("seurat_rds_path_store", "data"),
    State("sample_name_map_store", "data"),
)
def create_umap_name_controls(rds_path, name_map):
    """データ読み込み後、UMAP側にサンプル名変更UIを生成"""
    df = _interactive_data.get("plot_data")
    if df is None:
        return ""
    if not name_map:
        name_map = {}

    samples = sorted(df["Sample"].unique())
    if len(samples) <= 1:
        return ""

    controls = []
    for i, s in enumerate(samples):
        display_s = _display_name(s, name_map)
        controls.append(
            html.Div(
                style={"padding": "2px 8px"},
                children=[
                    dbc.Row(className="align-items-center", children=[
                        dbc.Col(width=4, children=[
                            html.Label(s, className="fw-bold small mb-0",
                                       style={"whiteSpace": "nowrap", "overflow": "hidden",
                                              "textOverflow": "ellipsis"}),
                        ]),
                        dbc.Col(width=8, children=[
                            dbc.Input(
                                id={"type": "umap_sample_rename_input", "index": s},
                                value=display_s if display_s != s else "",
                                placeholder=s,
                                size="sm", debounce=True,
                            ),
                        ]),
                    ]),
                    html.Hr(className="my-1") if i < len(samples) - 1 else html.Div(),
                ],
            )
        )

    return dbc.Accordion(
        [dbc.AccordionItem(title="サンプル名変更", children=controls)],
        start_collapsed=True,
        flush=True,
        always_open=True,
        style={"marginBottom": "8px"},
    )


# ---------------------------------------------------------------------------
# Spatial マーカーサイズ Auto ボタン
# ---------------------------------------------------------------------------

@callback(
    Output("spatial_marker_size", "value"),
    Input("spatial_marker_auto_btn", "n_clicks"),
    [State("spatial_rotation_store", "data"),
     State("interactive_sample", "value")],
    prevent_initial_call=True,
)
def auto_spatial_marker(n_clicks, rotation_store, sample):
    """通常ビュー: スライダーを自動モード(0)にリセット → 各サンプル個別に自動計算"""
    return 0


@callback(
    Output("fs_spatial_marker_size", "value"),
    Input("fs_spatial_marker_auto_btn", "n_clicks"),
    [State("spatial_rotation_store", "data"),
     State("fs_spatial_sample", "value"),
     State("fs_spatial_height_slider", "value")],
    prevent_initial_call=True,
)
def auto_fs_spatial_marker(n_clicks, rotation_store, sample, height_val):
    """フルスクリーン: スライダーを自動モード(0)にリセット → 各サンプル個別に自動計算"""
    return 0


# ---------------------------------------------------------------------------
# Feature Plot マーカーサイズ Auto ボタン
# ---------------------------------------------------------------------------

@callback(
    Output("feature_marker_size", "value"),
    Input("feature_marker_auto_btn", "n_clicks"),
    [State("spatial_rotation_store", "data"),
     State("feature_sample_select", "value")],
    prevent_initial_call=True,
)
def auto_feature_marker(n_clicks, rotation_store, sample):
    """Feature Plot: スライダーを自動モード(0)にリセット → 各サンプル個別に自動計算"""
    return 0


@callback(
    [Output("spatial_plots_container", "children"),
     Output("last_spatial_figure_store", "data")],
    [Input("interactive_sample", "value"),
     Input("spatial_highlight_cluster", "value"),
     Input("interactive_umap_plot", "selectedData"),
     Input("spatial_rotation_store", "data"),
     Input("spatial_show_labels", "value"),
     Input("spatial_marker_size", "value"),
     Input("spatial_exclude_cluster", "value"),
     Input("spatial_label_size", "value"),
     Input("seurat_rds_path_store", "data"),
     Input("sample_name_map_store", "data"),
     Input("fullscreen_closed_trigger", "data"),
     Input("custom_color_map_store", "data"),
     Input("spatial_columns_per_row", "value")],
    State("accumulated_label_positions", "data"),
)
def update_spatial_plots(sample, highlight_clusters, selected_data,
                         rotation_store, show_labels, marker_size,
                         exclude_clusters, label_size, rds_path, name_map,
                         _fs_trigger, custom_colors, columns_per_row,
                         accumulated_positions):
    df = _interactive_data.get("plot_data")
    if df is None or "SpatialX" not in df.columns:
        return html.Div("空間座標データがありません", className="text-muted p-3"), None

    if not rotation_store:
        rotation_store = {}
    if not name_map:
        name_map = {}

    # UMAP選択セルID
    selected_cell_ids = set()
    if selected_data and selected_data.get("points"):
        for pt in selected_data["points"]:
            if pt.get("text"):
                selected_cell_ids.add(pt["text"])

    color_map = _get_cluster_color_map(df["Cluster"], custom_colors)
    cluster_to_idx, discrete_cscale = _get_cluster_colorscale(df["Cluster"], custom_colors)
    all_pos = _get_merged_label_positions(accumulated_positions)
    spatial_pos = all_pos.get("spatial", {})

    # 表示対象サンプル
    if sample:
        samples_to_show = [sample]
    else:
        samples_to_show = sorted(df["Sample"].unique())

    graphs = []
    representative_fig = None
    for s in samples_to_show:
        df_s = df[df["Sample"] == s]
        # サンプル別の変換設定を取得
        transform = rotation_store.get(s, rotation_store.get("__all__", {"angle": 0, "flip_h": False, "flip_v": False}))
        # 後方互換: 旧形式(int)の場合
        if isinstance(transform, (int, float)):
            transform = {"angle": int(transform), "flip_h": False, "flip_v": False}
        rotation_deg = transform.get("angle", 0)
        flip_h = transform.get("flip_h", False)
        flip_v = transform.get("flip_v", False)
        display_s = _display_name(s, name_map)
        fig = _create_single_spatial_fig(df_s, color_map, highlight_clusters,
                                         selected_cell_ids,
                                         rotation_deg=rotation_deg,
                                         show_labels=show_labels,
                                         flip_h=flip_h, flip_v=flip_v,
                                         title=display_s, embed_legend=True,
                                         cluster_to_idx=cluster_to_idx,
                                         discrete_cscale=discrete_cscale,
                                         marker_size=marker_size or 0,
                                         exclude_clusters=exclude_clusters,
                                         label_size=label_size or 10,
                                         saved_positions=spatial_pos.get(s))
        if representative_fig is None:
            representative_fig = fig
        cfg = dict(_SPATIAL_IMG_CONFIG)
        cfg["toImageButtonOptions"] = dict(cfg["toImageButtonOptions"],
                                           filename=f"Spatial_{display_s}")
        if columns_per_row:
            n_cols = columns_per_row
            gap_total = (n_cols - 1) * 15
            flex_basis = f"calc({100 / n_cols:.2f}% - {gap_total / n_cols:.1f}px)"
            min_w = "0"
        else:
            n_cols = len(samples_to_show)
            flex_basis = f"{max(20, 90 // n_cols)}%"
            min_w = "300px"
        graphs.append(
            html.Div(
                style={"flex": f"1 1 {flex_basis}", "minWidth": min_w,
                        "border": "1px solid #dee2e6", "borderRadius": "6px",
                        "padding": "5px", "backgroundColor": "#fff"},
                children=[
                    dcc.Graph(id={"type": "spatial_graph", "index": s},
                              figure=fig, style={"height": "350px"}, config=cfg),
                ],
            )
        )

    container = html.Div(
        style={"display": "flex", "flexWrap": "wrap", "gap": "15px"},
        children=graphs,
    )
    # 代表figureをStoreに保存（HTMLエクスポート用）
    store_data = representative_fig.to_dict() if representative_fig else None
    return container, store_data


# ---------------------------------------------------------------------------
# Feature プロット（Spatial表示、Parquet高速読み込み優先 → R fallback）
# ---------------------------------------------------------------------------

_FEATURE_IMG_CONFIG = {
    "scrollZoom": True,
    "toImageButtonOptions": {
        "format": "png", "scale": 3,
    },
}


@callback(
    [Output("feature_plot_container", "children"),
     Output("feature_intensity_min", "placeholder"),
     Output("feature_intensity_max", "placeholder")],
    [Input("feature_select", "value"),
     Input("feature_sample_select", "value"),
     Input("feature_marker_size", "value"),
     Input("feature_intensity_min", "value"),
     Input("feature_intensity_max", "value"),
     Input("sample_name_map_store", "data"),
     Input("fullscreen_closed_trigger", "data"),
     Input("feature_columns_per_row", "value")],
    [State("seurat_rds_path_store", "data"),
     State("seurat_cache_dir_store", "data"),
     State("spatial_rotation_store", "data"),
     State("deg_data_store", "data")],
    prevent_initial_call=True,
)
def update_feature_plot(feature_name, sample, marker_size,
                        intensity_min, intensity_max,
                        name_map, _fs_trigger, columns_per_row,
                        rds_path, cache_dir_str, rotation_store,
                        deg_data):
    # 名前変更・フルスクリーン閉鎖トリガーだがFeature未選択 → スキップ
    if ctx.triggered_id in ("sample_name_map_store", "fullscreen_closed_trigger") and not feature_name:
        return no_update, no_update, no_update

    if not feature_name or not rds_path:
        return html.Div("m/z Feature を選択してください", className="text-muted p-3"), no_update, no_update

    df = _interactive_data.get("plot_data")
    if df is None:
        return html.Div("データが読み込まれていません", className="text-muted p-3"), no_update, no_update

    if "SpatialX" not in df.columns:
        return html.Div("空間座標データがありません", className="text-muted p-3"), no_update, no_update

    if not rotation_store:
        rotation_store = {}
    if not name_map:
        name_map = {}

    try:
        # Parquet からの高速読み込みを優先
        expression = None
        if cache_dir_str:
            cache_dir = Path(cache_dir_str)
            expression = _bridge.get_feature_expression_fast(
                cache_dir, feature_name
            )

        # Parquet にない場合は R subprocess で取得
        if expression is None:
            expression = _bridge.get_feature_expression(rds_path, feature_name)

        # expression を df に結合（CellID順で対応）
        df_plot = df.copy()
        df_plot["_expression"] = expression

        # 表示対象サンプル
        if sample:
            samples_to_show = [sample]
        else:
            samples_to_show = sorted(df_plot["Sample"].unique())

        # 全サンプル共通のカラースケール範囲を計算
        expr_vals = df_plot.loc[
            df_plot["Sample"].isin(samples_to_show), "_expression"
        ].values
        global_min = float(np.nanmin(expr_vals))
        global_max = float(np.nanmax(expr_vals))

        # ユーザー指定の Intensity Range（パーセント値）を強度値に変換
        val_range = global_max - global_min
        if intensity_min is not None:
            display_min = global_min + (intensity_min / 100.0) * val_range
        else:
            display_min = global_min
        if intensity_max is not None:
            display_max = global_min + (intensity_max / 100.0) * val_range
        else:
            display_max = global_max

        auto_mode = (marker_size is None or marker_size <= 0)

        graphs = []
        for s in samples_to_show:
            df_s = df_plot[df_plot["Sample"] == s]
            display_s = _display_name(s, name_map)

            # 変換設定を取得
            transform = rotation_store.get(
                s, rotation_store.get("__all__", {"angle": 0, "flip_h": False, "flip_v": False}))
            if isinstance(transform, (int, float)):
                transform = {"angle": int(transform), "flip_h": False, "flip_v": False}

            raw_x = df_s["SpatialX"].values
            raw_y = -df_s["SpatialY"].values  # Y軸反転
            plot_x, plot_y = _transform_coords(
                raw_x, raw_y,
                transform.get("angle", 0),
                flip_h=transform.get("flip_h", False),
                flip_v=transform.get("flip_v", False),
            )

            # マーカーサイズ: 自動モード(0)ならサンプル毎に計算
            if auto_mode:
                m_size = _calc_zero_gap_marker_size(plot_x, plot_y, render_height=280)
            else:
                m_size = marker_size

            # 最後のサンプルのみカラーバーを表示
            is_last = (s == samples_to_show[-1])
            marker_opts = dict(
                size=m_size,
                symbol="square",
                color=df_s["_expression"].values,
                colorscale="Plasma",
                cmin=display_min,
                cmax=display_max,
                showscale=is_last,
            )
            if is_last:
                marker_opts["colorbar"] = dict(
                    title=dict(text="Intensity", side="right"),
                    tickvals=[display_min, display_max],
                    ticktext=[
                        f"{int(intensity_min)}%" if intensity_min is not None else "0%",
                        f"{int(intensity_max)}%" if intensity_max is not None else "100%",
                    ],
                    len=0.8,
                    thickness=15,
                )

            fig = go.Figure()

            # TIC 背景（Greys, opacity=0.5）
            if "TotalCount" in df_s.columns:
                fig.add_trace(go.Scatter(
                    x=plot_x, y=plot_y, mode="markers",
                    marker=dict(size=m_size, symbol="square",
                                color=df_s["TotalCount"].values,
                                colorscale="Greys", opacity=0.5,
                                showscale=False),
                    hoverinfo="skip", showlegend=False,
                ))

            # 発現量オーバーレイ — ポイントごとの opacity で TIC 背景を透過
            # 発現量が低い → 透明（TIC が見える）、高い → 不透明（発現量が見える）
            expr_raw = df_s["_expression"].values
            if display_max > display_min:
                norm = np.clip((expr_raw - display_min) / (display_max - display_min), 0, 1)
            else:
                norm = np.zeros_like(expr_raw)
            marker_opts["opacity"] = np.where(norm > 0.01, 0.3 + 0.7 * norm, 0.0).tolist()
            fig.add_trace(go.Scatter(
                x=plot_x,
                y=plot_y,
                mode="markers",
                marker=marker_opts,
                text=df_s["CellID"],
                hovertemplate=f"{feature_name}: " + "%{marker.color:.4f}<br>%{text}<extra></extra>",
                showlegend=False,
            ))

            r_margin = 80 if is_last else 10
            fig.update_layout(
                title=dict(text=display_s, font=dict(size=14), x=0.5),
                xaxis=dict(showgrid=False, showline=False, zeroline=False,
                           showticklabels=False, title="", visible=False),
                yaxis=dict(scaleanchor="x", showgrid=False, showline=False, zeroline=False,
                           showticklabels=False, title="", visible=False),
                margin=dict(l=10, r=r_margin, t=30, b=10),
                plot_bgcolor="white",
            )

            cfg = dict(_FEATURE_IMG_CONFIG)
            cfg["toImageButtonOptions"] = dict(cfg["toImageButtonOptions"],
                                               filename=f"Feature_{feature_name}_{display_s}")
            if columns_per_row:
                n_cols = columns_per_row
                gap_total = (n_cols - 1) * 15
                flex_basis = f"calc({100 / n_cols:.2f}% - {gap_total / n_cols:.1f}px)"
                min_w = "0"
            else:
                n_cols = len(samples_to_show)
                flex_basis = f"{max(20, 90 // n_cols)}%"
                min_w = "300px"
            graphs.append(
                html.Div(
                    style={"flex": f"1 1 {flex_basis}", "minWidth": min_w,
                            "border": "1px solid #dee2e6", "borderRadius": "6px",
                            "padding": "5px", "backgroundColor": "#fff"},
                    children=[
                        dcc.Graph(figure=fig, style={"height": "350px"}, config=cfg),
                    ],
                )
            )

        container = html.Div(
            style={"display": "flex", "flexWrap": "wrap", "gap": "15px"},
            children=graphs,
        )

        # --- ③ アノテーション見出し ---
        annotation = ""
        if deg_data:
            for r in deg_data:
                if r.get("gene") == feature_name:
                    ann = r.get("annotation", "")
                    if _is_meaningful_annotation(ann, feature_name):
                        annotation = ann
                    break
        title_text = f"{feature_name}  ({annotation})" if annotation else feature_name
        heading = html.H6(
            title_text,
            className="text-center mt-2 mb-1",
            style={"color": "#333", "fontSize": "0.95rem"},
        )

        return html.Div([heading, container]), "0", "100"

    except Exception as e:
        return html.Div(f"エラー: {e}", className="text-danger p-3"), no_update, no_update


# ---------------------------------------------------------------------------
# Feature Plot ブックマークコールバック
# ---------------------------------------------------------------------------

@callback(
    Output("feature_history_store", "data"),
    Input("add_feature_bookmark_btn", "n_clicks"),
    [State("feature_select", "value"),
     State("feature_history_store", "data")],
    prevent_initial_call=True,
)
def add_feature_bookmark(n_clicks, feature_name, current_bookmarks):
    """ブックマーク追加ボタン → 現在の Feature をブックマークストアに保存"""
    if not n_clicks or not feature_name:
        return no_update
    bookmarks = list(current_bookmarks) if current_bookmarks else []
    if feature_name in bookmarks:
        bookmarks.remove(feature_name)
    bookmarks.insert(0, feature_name)
    bookmarks = bookmarks[:50]
    _save_interactive_settings("feature_bookmarks", bookmarks)
    return bookmarks


@callback(
    [Output("feature_history_store", "data", allow_duplicate=True),
     Output("feature_history_select", "value")],
    Input("remove_feature_bookmark_btn", "n_clicks"),
    [State("feature_history_select", "value"),
     State("feature_history_store", "data")],
    prevent_initial_call=True,
)
def remove_feature_bookmark(n_clicks, selected, current_bookmarks):
    """選択中のブックマークを削除"""
    if not n_clicks or not selected or not current_bookmarks:
        return no_update, no_update
    bookmarks = list(current_bookmarks)
    if selected in bookmarks:
        bookmarks.remove(selected)
    _save_interactive_settings("feature_bookmarks", bookmarks)
    return bookmarks, None


@callback(
    Output("feature_history_select", "options"),
    Input("feature_history_store", "data"),
    State("deg_data_store", "data"),
    prevent_initial_call=True,
)
def update_bookmark_options(bookmarks, deg_data):
    """ブックマークストアからドロップダウンオプションを生成"""
    if not bookmarks:
        return []
    ann_map = {}
    if deg_data:
        for r in deg_data:
            gene = r.get("gene", "")
            ann = r.get("annotation", "")
            if gene and _is_meaningful_annotation(ann, gene):
                ann_map[gene] = ann
    return [
        {"label": f"{f} ({ann_map[f]})" if f in ann_map else f, "value": f}
        for f in bookmarks
    ]


@callback(
    Output("feature_select", "value", allow_duplicate=True),
    Input("feature_history_select", "value"),
    prevent_initial_call=True,
)
def bookmark_to_feature(selected):
    """ブックマークドロップダウンから選択 → feature_select に値セット"""
    if not selected:
        return no_update
    return selected


# ---------------------------------------------------------------------------
# PPTX (Google Slides) エクスポート
# ---------------------------------------------------------------------------

def _fig_to_png_bytes(fig_dict, width=1200, height=800, scale=2):
    """Plotly figure dict を PNG バイト列に変換する。
    kaleido が未インストールの場合は None を返す。"""
    try:
        fig = go.Figure(fig_dict)
        fig.update_layout(paper_bgcolor="white", plot_bgcolor="white")
        return pio.to_image(fig, format="png", width=width, height=height, scale=scale)
    except Exception:
        return None


def _get_top_n_features_for_cluster(deg_data, cluster, n=5):
    """指定クラスタの DEG データから Top N up/down regulated feature を取得。

    Returns:
        (up_features: list[str], down_features: list[str])
    """
    if not deg_data:
        return [], []
    cluster_records = [
        r for r in deg_data
        if str(r.get("cluster", "")) == str(cluster)
    ]
    if not cluster_records:
        return [], []

    def _sort_key(r):
        p = r.get("p_val_adj_raw")
        if p is None:
            p = r.get("p_val_adj", 1.0)
        try:
            p = float(p)
            if np.isnan(p):
                p = 1.0
        except (ValueError, TypeError):
            p = 1.0
        fc = r.get("avg_log2FC", 0)
        try:
            fc = float(fc)
            if np.isnan(fc):
                fc = 0.0
        except (ValueError, TypeError):
            fc = 0.0
        return (p, -abs(fc))

    def _extract_top_n(records, n_items):
        sorted_recs = sorted(records, key=_sort_key)
        seen = set()
        result = []
        for r in sorted_recs:
            g = str(r.get("gene", ""))
            if g and g not in seen:
                seen.add(g)
                result.append(g)
                if len(result) >= n_items:
                    break
        return result

    up_records = []
    down_records = []
    for r in cluster_records:
        try:
            fc = float(r.get("avg_log2FC", 0) or 0)
        except (ValueError, TypeError):
            fc = 0.0
        if fc > 0:
            up_records.append(r)
        elif fc < 0:
            down_records.append(r)

    return _extract_top_n(up_records, n), _extract_top_n(down_records, n)


def _build_feature_plot_fig(df, feature_name, cache_dir_path, rds_path,
                            rotation_store=None, name_map=None, marker_size=3,
                            colorbar_tickformat=None, show_colorbar_title=True,
                            auto_marker_size=False):
    """単一 m/z Feature の Spatial Expression Plot figure を生成（PPTX 用）。

    Returns:
        go.Figure or None
    """
    from plotly.subplots import make_subplots

    if rotation_store is None:
        rotation_store = {}
    if name_map is None:
        name_map = {}

    # 発現データ取得
    expression = None
    if cache_dir_path:
        cache_dir = Path(cache_dir_path) if isinstance(cache_dir_path, str) else cache_dir_path
        expression = _bridge.get_feature_expression_fast(cache_dir, feature_name)
    if expression is None and rds_path:
        try:
            expression = _bridge.get_feature_expression(rds_path, feature_name)
        except Exception:
            return None
    if expression is None:
        return None

    df_plot = df.copy()
    df_plot["_expression"] = expression.values if hasattr(expression, "values") else expression

    if "SpatialX" not in df_plot.columns:
        return None

    samples = sorted(df_plot["Sample"].unique())
    n_samples = len(samples)
    subplot_titles = [_display_name(s, name_map) for s in samples]

    fig = make_subplots(
        rows=1, cols=n_samples,
        subplot_titles=subplot_titles,
        horizontal_spacing=0.03,
    )

    expr_vals = df_plot["_expression"].values
    global_min = float(np.nanmin(expr_vals))
    global_max = float(np.nanmax(expr_vals))

    for idx, s in enumerate(samples, 1):
        df_s = df_plot[df_plot["Sample"] == s]
        transform = rotation_store.get(
            s, rotation_store.get("__all__", {"angle": 0, "flip_h": False, "flip_v": False}))
        if isinstance(transform, (int, float)):
            transform = {"angle": int(transform), "flip_h": False, "flip_v": False}

        raw_x = df_s["SpatialX"].values
        raw_y = -df_s["SpatialY"].values
        plot_x, plot_y = _transform_coords(
            raw_x, raw_y,
            transform.get("angle", 0),
            flip_h=transform.get("flip_h", False),
            flip_v=transform.get("flip_v", False),
        )

        # サンプル別マーカーサイズ自動計算
        ms = marker_size
        if auto_marker_size and len(plot_x) > 1:
            sorted_ux = np.sort(np.unique(plot_x))
            if len(sorted_ux) > 1:
                min_sp = float(np.min(np.diff(sorted_ux)))
                xr = float(plot_x.max() - plot_x.min())
                yr = float(plot_y.max() - plot_y.min()) if len(plot_y) > 1 else 1.0
                # Feature Plot は height=280 描画、margin 約40px → 有効240px
                eff_h = 240
                if yr > 0 and xr > 0:
                    eff_w = eff_h * (xr / yr)
                    ms = max(2, round(min_sp * eff_w / xr * 1.5))

        # TIC 背景
        if "TotalCount" in df_s.columns:
            fig.add_trace(go.Scatter(
                x=plot_x, y=plot_y, mode="markers",
                marker=dict(size=ms, symbol="square", color=df_s["TotalCount"].values,
                            colorscale="Greys", opacity=0.5, showscale=False),
                hoverinfo="skip", showlegend=False,
            ), row=1, col=idx)

        # 発現量オーバーレイ
        is_last = (idx == n_samples)
        marker_opts = dict(
            size=ms,
            symbol="square",
            color=df_s["_expression"].values,
            colorscale="Plasma",
            cmin=global_min, cmax=global_max,
            showscale=is_last,
            opacity=0.8,
        )
        if is_last:
            cb_opts = dict(
                len=0.8, thickness=15,
                tickvals=[global_min, global_max],
                ticktext=["0%", "100%"],
            )
            if show_colorbar_title:
                cb_opts["title"] = "Intensity"
            marker_opts["colorbar"] = cb_opts
        fig.add_trace(go.Scatter(
            x=plot_x, y=plot_y, mode="markers",
            marker=marker_opts, showlegend=False,
        ), row=1, col=idx)

    fig.update_layout(
        title=dict(text=feature_name, font=dict(size=20), x=0.5),
        plot_bgcolor="white", paper_bgcolor="white",
        margin=dict(l=10, r=50, t=50, b=10),
    )
    # Subplot titles のフォントサイズ拡大
    fig.update_annotations(font_size=16)
    for i in range(1, n_samples + 1):
        fig.update_xaxes(visible=False, row=1, col=i)
        xanchor = f"x{i}" if i > 1 else "x"
        fig.update_yaxes(visible=False, scaleanchor=xanchor, row=1, col=i)

    return fig


def _build_volcano_fig_for_cluster(deg_data, cluster, fc_thresh=0.5, p_thresh=1.3,
                                   marker_size=8):
    """指定クラスタの Volcano Plot figure を生成（PPTX 用）。

    Returns:
        go.Figure or None
    """
    if not deg_data:
        return None
    df = pd.DataFrame(deg_data)
    if "p_val_adj_raw" in df.columns:
        df["p_num"] = pd.to_numeric(df["p_val_adj_raw"], errors="coerce")
    else:
        df["p_num"] = pd.to_numeric(df["p_val_adj"], errors="coerce")
    df["avg_log2FC"] = pd.to_numeric(df["avg_log2FC"], errors="coerce")
    min_nonzero_p = (
        df.loc[df["p_num"] > 0, "p_num"].min()
        if (df["p_num"] > 0).any()
        else 5e-324
    )
    df["neg_log10_p"] = -np.log10(df["p_num"].clip(lower=min_nonzero_p))

    if "annotation" in df.columns:
        df["display_text"] = df.apply(
            lambda r: f"{r['gene']}\n({r['annotation']})"
            if _is_meaningful_annotation(r.get("annotation", ""), r.get("gene", ""))
            else r["gene"],
            axis=1,
        )
    else:
        df["display_text"] = df["gene"]

    df = df[df["cluster"].astype(str) == str(cluster)]
    if df.empty:
        return None

    fig = go.Figure()
    for reg, color, label in [
        ("Up", "#FF2D2D", "Up-regulated"),
        ("Down", "#1E5BFF", "Down-regulated"),
        ("NS", "#7A7A7A", "Not significant"),
    ]:
        if reg == "Up":
            mask = (df["neg_log10_p"] >= p_thresh) & (df["avg_log2FC"] >= fc_thresh)
        elif reg == "Down":
            mask = (df["neg_log10_p"] >= p_thresh) & (df["avg_log2FC"] <= -fc_thresh)
        else:
            mask = ~(
                (df["neg_log10_p"] >= p_thresh) & (df["avg_log2FC"].abs() >= fc_thresh)
            )
        sub = df[mask]
        if len(sub) > 0:
            fig.add_trace(go.Scattergl(
                x=sub["avg_log2FC"], y=sub["neg_log10_p"],
                mode="markers",
                marker=dict(size=marker_size, color=color, opacity=0.7),
                name=label, text=sub["display_text"],
                hovertemplate=(
                    "<b>%{text}</b><br>"
                    "log2FC: %{x:.3f}<br>"
                    "-log10(p): %{y:.2f}<extra></extra>"
                ),
            ))

    fig.add_hline(y=p_thresh, line_dash="dash", line_color="gray", opacity=0.5)
    fig.add_vline(x=fc_thresh, line_dash="dash", line_color="gray", opacity=0.5)
    fig.add_vline(x=-fc_thresh, line_dash="dash", line_color="gray", opacity=0.5)
    fig.update_layout(
        title=dict(text=f"Volcano Plot - Cluster {cluster}", font=dict(size=20), x=0.5),
        xaxis_title="avg_log2FC",
        yaxis_title="-log10(p_val_adj)",
        template="plotly_white",
        margin=dict(l=50, r=20, t=40, b=40),
    )
    return fig


def _pptx_add_title_bar(slide, title_text):
    """PPTX スライドにタイトルバーを追加する。"""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(8), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def _pptx_add_image(slide, png_bytes, left, top, width, height):
    """PNG バイト列を PPTX スライドの指定位置に配置する。"""
    if png_bytes:
        img_stream = BytesIO(png_bytes)
        slide.shapes.add_picture(img_stream, left, top, width, height)


def _pptx_add_image_preserve_ratio(slide, png_bytes, left, top,
                                    max_width, max_height,
                                    png_w=None, png_h=None):
    """PNG 元比率を保持して max_width × max_height 内に収まるよう配置する。

    png_w, png_h: PNG 生成時のピクセル寸法（比率計算に使用）。
    未指定時は max_width × max_height をそのまま使用（従来互換）。
    """
    if not png_bytes:
        return
    if png_w and png_h:
        aspect = png_w / png_h
        max_aspect = max_width / max_height if max_height else 1.0
        if max_aspect > aspect:
            # 高さ制約: height = max_height, width = height * aspect
            h = max_height
            w = int(h * aspect)
        else:
            # 幅制約: width = max_width, height = width / aspect
            w = max_width
            h = int(w / aspect)
        # 枠内中央寄せ
        left = left + int((max_width - w) / 2)
        top = top + int((max_height - h) / 2)
    else:
        w, h = max_width, max_height
    img_stream = BytesIO(png_bytes)
    slide.shapes.add_picture(img_stream, left, top, w, h)


def _square_tile_dims(tile_w, row_h, margin=0.95):
    """正方形画像のタイル配置寸法を計算する。

    Returns:
        (side, side, offset): 正方形の辺長と、タイル内の左オフセット
    """
    avail_w = tile_w * margin
    side = min(avail_w, row_h)
    offset = (tile_w - side) / 2
    return side, side, offset


def _build_cluster_legend_fig(cluster_list, color_map, font_size=18,
                              marker_size=14):
    """クラスタレジェンド専用のPlotly Figureを生成"""
    fig = go.Figure()
    for cl in sorted(cluster_list, key=_cluster_sort_key):
        fig.add_trace(go.Scatter(
            x=[None], y=[None], mode="markers",
            marker=dict(size=marker_size,
                        color=color_map.get(str(cl), "#999")),
            name=f"Cluster {cl}", showlegend=True,
        ))
    fig.update_layout(
        showlegend=True,
        legend=dict(
            font=dict(size=font_size),
            itemsizing="constant",
            yanchor="middle", y=0.5,
            xanchor="center", x=0.5,
            tracegroupgap=4,
        ),
        margin=dict(l=0, r=0, t=0, b=0),
        xaxis=dict(visible=False), yaxis=dict(visible=False),
        plot_bgcolor="white", paper_bgcolor="white",
        width=200, height=600,
    )
    return fig


def _build_sample_legend_fig(sample_list, sample_color_map, name_map=None,
                              font_size=18, marker_size=14):
    """サンプルレジェンド専用のPlotly Figureを生成"""
    fig = go.Figure()
    for s in sorted(sample_list):
        display = (name_map or {}).get(s, s) if name_map else s
        fig.add_trace(go.Scatter(
            x=[None], y=[None], mode="markers",
            marker=dict(size=marker_size,
                        color=sample_color_map.get(str(s), "#999")),
            name=str(display), showlegend=True,
        ))
    fig.update_layout(
        showlegend=True,
        legend=dict(
            font=dict(size=font_size),
            itemsizing="constant",
            yanchor="middle", y=0.5,
            xanchor="center", x=0.5,
            tracegroupgap=4,
        ),
        margin=dict(l=0, r=0, t=0, b=0),
        xaxis=dict(visible=False), yaxis=dict(visible=False),
        plot_bgcolor="white", paper_bgcolor="white",
        width=200, height=600,
    )
    return fig


def _pptx_add_sections(prs, section_map):
    """PPTX にセクション情報を追加する（PowerPoint のセクションパネルに表示される）。

    Args:
        prs: python-pptx Presentation オブジェクト
        section_map: list of (section_name, start_slide_idx, end_slide_idx)
            slide indices は 0-based、両端を含む。
    """
    import uuid
    from lxml import etree

    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    p14_ns = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    ext_uri = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"

    prs_elem = prs.element  # CT_Presentation lxml element

    # p14 名前空間プレフィックスを登録（PowerPoint互換性のため）
    etree.register_namespace("p14", p14_ns)

    # sldIdLst から全スライド ID を取得
    sldIdLst = prs_elem.find(f"{{{p_ns}}}sldIdLst")
    if sldIdLst is None:
        return

    sld_ids = []
    for sldId in sldIdLst:
        sid = sldId.get("id")
        if sid:
            sld_ids.append(sid)
    if not sld_ids:
        return

    # extLst を検索または作成
    extLst = prs_elem.find(f"{{{p_ns}}}extLst")
    if extLst is None:
        extLst = etree.SubElement(prs_elem, f"{{{p_ns}}}extLst")

    # ext 要素を作成（セクション用 URI）
    ext = etree.SubElement(extLst, f"{{{p_ns}}}ext")
    ext.set("uri", ext_uri)

    # sectionLst を作成
    sectionLst = etree.SubElement(ext, f"{{{p14_ns}}}sectionLst")

    for sec_name, start_idx, end_idx in section_map:
        section = etree.SubElement(sectionLst, f"{{{p14_ns}}}section")
        section.set("name", sec_name)
        section.set("id", "{" + str(uuid.uuid4()).upper() + "}")

        sldIdLst_sec = etree.SubElement(section, f"{{{p14_ns}}}sldIdLst")

        for i in range(start_idx, min(end_idx + 1, len(sld_ids))):
            sldId_ref = etree.SubElement(sldIdLst_sec, f"{{{p14_ns}}}sldId")
            sldId_ref.set("id", sld_ids[i])


def _build_heatmap_for_pptx(heatmap_fig, deg_data, df, cache_dir, top_n,
                             mrm_path=None):
    """PPTX 用にZ-score + アノテーション付きヒートマップ figure を生成/調整する。

    アプリの heatmap_fig がすでに Z-score ならアノテーションだけ補完し、
    Raw なら expression_matrix から再計算して Z-score + アノテーション付きで返す。
    """
    # ---- deg_data / cache_dir がなければアプリ図をそのまま返す ----
    if not deg_data or not cache_dir:
        return heatmap_fig

    cache_dir_path = Path(cache_dir) if isinstance(cache_dir, str) else cache_dir
    expr_path = cache_dir_path / "expression_matrix.parquet" if cache_dir_path else None

    # ---- アノテーション辞書を構築 ----
    gene_to_annotation = {}
    for r in deg_data:
        gene = r.get("gene", "")
        ann = r.get("annotation", "")
        if gene and _is_meaningful_annotation(ann, gene):
            gene_to_annotation[gene] = ann

    # MRM fallback
    if not gene_to_annotation and mrm_path:
        try:
            mz_to_compound = _build_mz_to_compound_map(mrm_path, tolerance=0.1)
        except Exception:
            mz_to_compound = {}
    else:
        mz_to_compound = {}

    # ---- heatmap_fig が Z-score かどうかを判定 ----
    is_zscore = False
    if heatmap_fig:
        traces = heatmap_fig.get("data", [])
        for t in traces:
            if t.get("type") == "heatmap" and t.get("zmid") == 0:
                is_zscore = True
                break

    # ---- Z-score 済みの場合: アノテーション補完のみ ----
    if is_zscore and heatmap_fig:
        fig = go.Figure(heatmap_fig)
        # Y軸ラベルにアノテーションを補完
        for t in fig.data:
            if hasattr(t, "y") and t.y is not None:
                current_labels = list(t.y)
                needs_update = False
                new_labels = []
                for lbl in current_labels:
                    gene = str(lbl).split(" (")[0]  # 既存アノテーション除去
                    if gene in gene_to_annotation:
                        new_labels.append(f"{gene} ({gene_to_annotation[gene]})")
                        needs_update = True
                    elif mz_to_compound:
                        annotated = _annotate_gene_labels(
                            [gene], mz_to_compound, tolerance=0.1)
                        new_labels.append(annotated[0])
                        if annotated[0] != gene:
                            needs_update = True
                    else:
                        new_labels.append(lbl)
                if needs_update:
                    t.y = new_labels
        # 左マージンを再調整
        all_labels = []
        for t in fig.data:
            if hasattr(t, "y") and t.y is not None:
                all_labels.extend(str(l) for l in t.y)
        if all_labels:
            max_len = max(len(l) for l in all_labels)
            left_margin = min(max(max_len * 7, 120), 350)
            fig.update_layout(margin=dict(l=left_margin))
        return fig.to_dict()

    # ---- Raw / heatmap_fig がない場合: expression_matrix から新規生成 ----
    if not expr_path or not expr_path.exists():
        return heatmap_fig  # fallback: アプリ図をそのまま
    if df is None or df.empty:
        return heatmap_fig

    try:
        df_deg = pd.DataFrame(deg_data)
        df_deg["p_num"] = pd.to_numeric(df_deg.get("p_val_adj", ""), errors="coerce")
        top_markers = df_deg.sort_values("p_num").groupby("cluster").head(top_n)
        genes = top_markers["gene"].unique().tolist()
        if not genes:
            return heatmap_fig

        # 利用可能な遺伝子のみ読み込み
        available = []
        for g in genes:
            try:
                pd.read_parquet(expr_path, columns=[g])
                available.append(g)
            except Exception:
                continue
        if not available:
            return heatmap_fig

        expr_df = pd.read_parquet(expr_path, columns=["CellID"] + available)
        merged = expr_df.merge(
            df[["CellID", "Cluster"]], on="CellID", how="inner"
        )
        cluster_means = merged.groupby("Cluster")[available].mean()
        cluster_means = cluster_means.reindex(
            sorted(cluster_means.index, key=_cluster_sort_key)
        )

        # Z-score 変換
        z_data = cluster_means.values.copy()
        col_mean = z_data.mean(axis=0)
        col_std = z_data.std(axis=0)
        col_std[col_std == 0] = 1
        z_data = (z_data - col_mean) / col_std

        # Y軸ラベル（アノテーション付き）
        y_labels = []
        for g in available:
            if g in gene_to_annotation:
                y_labels.append(f"{g} ({gene_to_annotation[g]})")
            elif mz_to_compound:
                annotated = _annotate_gene_labels(
                    [g], mz_to_compound, tolerance=0.1)
                y_labels.append(annotated[0])
            else:
                y_labels.append(g)

        fig = go.Figure(go.Heatmap(
            z=z_data.T,
            x=[f"C{c}" for c in cluster_means.index],
            y=y_labels,
            colorscale="RdBu_r",
            zmid=0,
            hovertemplate=(
                "Cluster: %{x}<br>Gene: %{y}<br>"
                "Value: %{z:.3f}<extra></extra>"
            ),
        ))
        max_label_len = max(len(str(l)) for l in y_labels) if y_labels else 10
        left_margin = min(max(max_label_len * 7, 120), 350)
        fig.update_layout(
            title=dict(
                text=f"Top {top_n} DEG Heatmap (Z-score)",
                font=dict(size=14), x=0.5),
            xaxis_title="Cluster",
            yaxis_title="Gene / m/z",
            template="plotly_white",
            margin=dict(l=left_margin, r=20, t=40, b=40),
            yaxis=dict(autorange="reversed"),
        )
        return fig.to_dict()

    except Exception:
        return heatmap_fig  # エラー時はアプリ図にフォールバック


def _build_pptx(umap_fig, spatial_fig, meta, cluster_stats_data, rds_path,
                 sub_name="", volcano_fig=None, heatmap_fig=None,
                 deg_data=None, top_n=5, df=None, cache_dir=None,
                 custom_colors=None, rotation_store=None, name_map=None,
                 set_progress=None, mrm_path=None,
                 existing_prs=None, progress_offset=0, progress_total=None,
                 saved_positions=None):
    """グローバル概要 + クラスターごとの詳細スライドを含む PPTX を生成し bytes を返す。

    グローバルセクション:
        1. タイトル  2. Heatmap  3. UMAP+Spatial (統合)  4. クラスタ統計
    クラスターセクション (各クラスター × 2 スライド):
        A. Volcano Plot + Top N Feature Plots
        B. UMAP (ハイライト) + Spatial (ハイライト)

    existing_prs: 既存のPresentationオブジェクト。指定時はそこにスライドを追加し、
                  bytes は返さず None を返す。
    progress_offset: 進捗計算のオフセット（複数手法ループ時に使用）
    progress_total: 進捗計算の全体ステップ数（複数手法ループ時に使用）
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    # 進捗計算用
    _clusters_for_progress = []
    if df is not None:
        try:
            _clusters_for_progress = sorted(
                set(str(c) for c in df["Cluster"].unique()),
                key=_cluster_sort_key)
        except Exception:
            pass
    _local_steps = 5 + len(_clusters_for_progress) * 2  # +1 for UMAP by Sample
    _total_steps = progress_total if progress_total else _local_steps
    _current_step = [progress_offset]  # mutable for nested function

    def _progress(label=""):
        _current_step[0] += 1
        if set_progress:
            pct = int(_current_step[0] / _total_steps * 100)
            set_progress((min(pct, 99), 100, label))

    if existing_prs is not None:
        prs = existing_prs
    else:
        prs = Presentation()
        # 16:9 ワイドスクリーン
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # =====================================================================
    # グローバルセクション
    # =====================================================================

    # --- スライド 1: タイトル ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True

    if sub_name:
        p_title = tf.paragraphs[0]
        p_title.text = sub_name
        p_title.font.size = Pt(40)
        p_title.font.bold = True
        p_title.alignment = PP_ALIGN.CENTER

        p_sub = tf.add_paragraph()
        p_sub.text = "MSI Interactive Analysis Report"
        p_sub.font.size = Pt(20)
        p_sub.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p_sub.alignment = PP_ALIGN.CENTER
    else:
        p_title = tf.paragraphs[0]
        p_title.text = "MSI Interactive Analysis Report"
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p2.alignment = PP_ALIGN.CENTER

    if rds_path:
        p3 = tf.add_paragraph()
        p3.text = f"Source: {Path(rds_path).name}"
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p3.alignment = PP_ALIGN.CENTER

    if meta:
        p4 = tf.add_paragraph()
        samples_str = ", ".join(meta.get("samples", []))
        p4.text = (
            f"Cells: {meta.get('n_cells', '?')} | "
            f"Clusters: {meta.get('n_clusters', '?')} | "
            f"Samples: {samples_str}"
        )
        p4.font.size = Pt(12)
        p4.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p4.alignment = PP_ALIGN.CENTER

    _progress("タイトルスライド")

    # --- スライド 2: Heatmap (Z-score + アノテーション) ---
    heatmap_for_pptx = _build_heatmap_for_pptx(
        heatmap_fig, deg_data, df, cache_dir, top_n, mrm_path=mrm_path)
    if heatmap_for_pptx:
        png_bytes = _fig_to_png_bytes(heatmap_for_pptx, width=1200, height=800)
        if png_bytes:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_add_title_bar(slide, f"Heatmap — Top {top_n} DEG (Z-score)")
            _pptx_add_image_preserve_ratio(
                slide, png_bytes,
                int((prs.slide_width - Inches(12)) / 2), Inches(0.9),
                Inches(12), Inches(6.3),
                png_w=1200, png_h=800)

    _progress("Heatmap")

    # --- スライド 3: UMAP + Spatial 統合 (サンプル別) ---
    if df is not None and not df.empty and "SpatialX" in df.columns:
        import math
        color_map_global = _get_cluster_color_map(df["Cluster"], custom_colors)
        all_samples = sorted(df["Sample"].unique())
        if not rotation_store:
            rotation_store = {}
        if not name_map:
            name_map = {}

        if all_samples:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_add_title_bar(slide, "UMAP & Spatial Mapping")

            n_sp = len(all_samples)
            avail_w = 11.0  # 0.3〜11.3" (右端にレジェンド配置)

            # UMAP の全データ軸範囲（サンプル間で統一）
            umap1_min = float(df["UMAP_1"].min())
            umap1_max = float(df["UMAP_1"].max())
            umap2_min = float(df["UMAP_2"].min())
            umap2_max = float(df["UMAP_2"].max())
            umap_pad = max(umap1_max - umap1_min, umap2_max - umap2_min) * 0.05
            umap_xrange = [umap1_min - umap_pad, umap1_max + umap_pad]
            umap_yrange = [umap2_min - umap_pad, umap2_max + umap_pad]

            # 上段: サンプル別UMAP (y=0.9" 〜 y=4.0", 高さ3.1")
            tile_w_umap = avail_w / n_sp
            for idx, s in enumerate(all_samples):
                df_s = df[df["Sample"] == s]
                _umap_pos = (saved_positions or {}).get("umap_integrated", {})
                umap_s = _build_umap_integrated_fig(
                    df_s, color_by="Cluster", highlight_clusters=None,
                    show_legend=False, show_labels=True,
                    title=_display_name(s, name_map),
                    marker_size=3, custom_colors=custom_colors,
                    title_font_size=40, label_size=24,
                    saved_positions=_umap_pos)
                if umap_s is not None:
                    # 全データと同じ軸範囲に固定
                    umap_s.update_xaxes(range=umap_xrange)
                    umap_s.update_yaxes(range=umap_yrange)
                    u_dict = (umap_s.to_dict()
                              if hasattr(umap_s, "to_dict") else umap_s)
                    u_png = _fig_to_png_bytes(
                        u_dict, width=600, height=600, scale=2)
                    _uw, _uh, _uoff = _square_tile_dims(
                        tile_w_umap, 3.0)
                    u_left = Inches(0.3 + idx * tile_w_umap + _uoff)
                    _pptx_add_image(slide, u_png,
                                    int(u_left), Inches(0.9),
                                    Inches(_uw), Inches(_uh))

            # 下段: サンプル別Spatial (y=4.1" 〜 y=7.3", 高さ3.2")
            tile_w_sp = avail_w / n_sp
            for idx, s in enumerate(all_samples):
                df_s = df[df["Sample"] == s]
                transform = rotation_store.get(
                    s, rotation_store.get(
                        "__all__",
                        {"angle": 0, "flip_h": False, "flip_v": False}))
                if isinstance(transform, (int, float)):
                    transform = {"angle": int(transform),
                                 "flip_h": False, "flip_v": False}

                _sp_pos = (saved_positions or {}).get("spatial", {}).get(s, {})
                sp_fig = _create_single_spatial_fig(
                    df_s, color_map_global,
                    highlight_clusters=None,
                    selected_cell_ids=set(),
                    rotation_deg=transform.get("angle", 0),
                    show_labels=True,
                    flip_h=transform.get("flip_h", False),
                    flip_v=transform.get("flip_v", False),
                    title=_display_name(s, name_map),
                    marker_size=0,
                    render_height=560,
                    embed_legend=False,
                    title_font_size=40, label_size=24,
                    saved_positions=_sp_pos)
                if sp_fig is not None:
                    sp_dict = (sp_fig.to_dict()
                               if hasattr(sp_fig, "to_dict") else sp_fig)
                    sp_png = _fig_to_png_bytes(
                        sp_dict, width=600, height=600, scale=2)
                    _sw, _sh, _soff = _square_tile_dims(
                        tile_w_sp, 3.1)
                    sp_left = Inches(0.3 + idx * tile_w_sp + _soff)
                    _pptx_add_image(slide, sp_png,
                                    int(sp_left), Inches(4.1),
                                    Inches(_sw), Inches(_sh))

            # クラスタレジェンド（右端に配置）
            legend_fig = _build_cluster_legend_fig(
                df["Cluster"].unique(), color_map_global)
            legend_png = _fig_to_png_bytes(
                legend_fig.to_dict(), width=200, height=600, scale=2)
            _pptx_add_image_preserve_ratio(slide, legend_png,
                                           Inches(11.5), Inches(0.9),
                                           Inches(1.3), Inches(6.3),
                                           png_w=200, png_h=600)

            # --- UMAP (by Sample) スライド ---
            if len(all_samples) > 1:
                sample_cmap = _get_sample_color_map(all_samples)
                slide_by_sample = prs.slides.add_slide(prs.slide_layouts[6])
                _pptx_add_title_bar(slide_by_sample, "UMAP (by Sample)")

                tile_w_s = avail_w / n_sp
                for idx, s in enumerate(all_samples):
                    df_s = df[df["Sample"] == s]
                    umap_s_sample = _build_umap_integrated_fig(
                        df_s, color_by="Sample",
                        highlight_clusters=None,
                        show_legend=False, show_labels=False,
                        title=_display_name(s, name_map),
                        marker_size=3, custom_colors=sample_cmap,
                        title_font_size=40)
                    if umap_s_sample is not None:
                        umap_s_sample.update_xaxes(range=umap_xrange)
                        umap_s_sample.update_yaxes(range=umap_yrange)
                        u_dict = (umap_s_sample.to_dict()
                                  if hasattr(umap_s_sample, "to_dict")
                                  else umap_s_sample)
                        u_png = _fig_to_png_bytes(
                            u_dict, width=600, height=600, scale=2)
                        _uw, _uh, _uoff = _square_tile_dims(
                            tile_w_s, 5.5)
                        u_left = Inches(0.3 + idx * tile_w_s + _uoff)
                        _pptx_add_image(slide_by_sample, u_png,
                                        int(u_left), Inches(0.9),
                                        Inches(_uw), Inches(_uh))

                # サンプルレジェンド（右端）
                sample_legend = _build_sample_legend_fig(
                    all_samples, sample_cmap, name_map=name_map)
                sample_legend_png = _fig_to_png_bytes(
                    sample_legend.to_dict(), width=200, height=600, scale=2)
                _pptx_add_image_preserve_ratio(
                    slide_by_sample, sample_legend_png,
                    Inches(11.5), Inches(0.9),
                    Inches(1.3), Inches(6.3),
                    png_w=200, png_h=600)
                _progress("UMAP by Sample")

    elif spatial_fig:
        # SpatialX がない場合はアプリの図をフォールバック
        fig_check = go.Figure(spatial_fig)
        if fig_check.data:
            png_bytes = _fig_to_png_bytes(spatial_fig, width=1200, height=800)
            if png_bytes:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _pptx_add_title_bar(slide, "Spatial Mapping")
                _pptx_add_image_preserve_ratio(
                    slide, png_bytes,
                    int((prs.slide_width - Inches(12)) / 2), Inches(0.9),
                    Inches(12), Inches(6.3),
                    png_w=1200, png_h=800)

    _progress("UMAP & Spatial")

    # --- スライド 5: クラスタ統計 ---
    if cluster_stats_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _pptx_add_title_bar(slide, "Cluster Statistics")

        n_rows = len(cluster_stats_data) + 1
        n_cols = 3
        table_w = Inches(6)
        table_h = Inches(min(5.5, 0.4 * n_rows))
        left = int((prs.slide_width - table_w) / 2)
        table_shape = slide.shapes.add_table(
            n_rows, n_cols, left, Inches(1.0), table_w, table_h,
        )
        table = table_shape.table

        for j, h in enumerate(["Cluster", "Pixels", "%"]):
            cell = table.cell(0, j)
            cell.text = h
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.bold = True

        for i, row_data in enumerate(cluster_stats_data):
            table.cell(i + 1, 0).text = str(row_data.get("Cluster", ""))
            table.cell(i + 1, 1).text = str(row_data.get("Pixels", ""))
            table.cell(i + 1, 2).text = str(row_data.get("Percent", ""))
            for j in range(n_cols):
                table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(11)

    _progress("クラスタ統計")

    # =====================================================================
    # クラスターセクション（各クラスター × 2 スライド）
    # =====================================================================
    if df is not None and not df.empty:
        clusters = sorted(df["Cluster"].unique(), key=_cluster_sort_key)
        color_map = _get_cluster_color_map(df["Cluster"], custom_colors)
        has_spatial = "SpatialX" in df.columns
        samples = sorted(df["Sample"].unique()) if has_spatial else []
        if not rotation_store:
            rotation_store = {}
        if not name_map:
            name_map = {}

        # UMAP 軸範囲（Per-cluster Slide B でサンプル間統一に使用）
        cl_umap1_min = float(df["UMAP_1"].min())
        cl_umap1_max = float(df["UMAP_1"].max())
        cl_umap2_min = float(df["UMAP_2"].min())
        cl_umap2_max = float(df["UMAP_2"].max())
        cl_umap_pad = max(cl_umap1_max - cl_umap1_min,
                          cl_umap2_max - cl_umap2_min) * 0.05
        cl_umap_xrange = [cl_umap1_min - cl_umap_pad,
                          cl_umap1_max + cl_umap_pad]
        cl_umap_yrange = [cl_umap2_min - cl_umap_pad,
                          cl_umap2_max + cl_umap_pad]

        rds_path_str = str(rds_path) if rds_path else None
        cache_dir_path = (
            Path(cache_dir) if isinstance(cache_dir, str) and cache_dir
            else cache_dir
        )

        for cl in clusters:
            cl_str = str(cl)

            # === Slide A: Volcano + Feature Plots ===
            slide_a = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_add_title_bar(slide_a, f"Cluster {cl_str} — DEG Analysis")

            # Top N features (up / down)
            up_features, down_features = _get_top_n_features_for_cluster(
                deg_data, cl_str, n=top_n)

            # Volcano Plot を先に生成（レイアウト計算に必要）
            volcano_cl = _build_volcano_fig_for_cluster(deg_data, cl_str)

            # ---- 自動レイアウト計算 ----
            has_up = bool(up_features)
            has_down = bool(down_features)
            has_volcano = volcano_cl is not None

            _avail_top = 0.65      # Feature配置開始Y (タイトルバー下)
            _avail_bottom = 7.35   # スライド下端マージン
            _avail_h = _avail_bottom - _avail_top  # 6.7"
            _label_h = 0.25        # ラベル行の高さ
            _gap = 0.1             # セクション間隙間
            _volcano_h = 2.5       # Volcano固定高さ

            _n_feat_rows = (1 if has_up else 0) + (1 if has_down else 0)

            if _n_feat_rows > 0 and has_volcano:
                _non_feat = _n_feat_rows * (_label_h + _gap) + _gap + _volcano_h
                _feat_h_val = (_avail_h - _non_feat) / _n_feat_rows
            elif _n_feat_rows > 0:
                _non_feat = _n_feat_rows * (_label_h + _gap)
                _feat_h_val = (_avail_h - _non_feat) / _n_feat_rows
            else:
                _feat_h_val = 0

            _feat_h_val = max(1.5, min(_feat_h_val, 3.5))

            feat_w_val = min(2.4, 12.0 / max(top_n, 1))
            feat_w = Inches(feat_w_val)
            feat_h = Inches(_feat_h_val)

            # Y座標を順番に計算
            _cur_y = _avail_top
            _up_label_y = _up_plot_y = 0
            _down_label_y = _down_plot_y = 0
            _volcano_y = 0

            if has_up:
                _up_label_y = _cur_y
                _cur_y += _label_h + _gap
                _up_plot_y = _cur_y
                _cur_y += _feat_h_val + _gap
            if has_down:
                _down_label_y = _cur_y
                _cur_y += _label_h + _gap
                _down_plot_y = _cur_y
                _cur_y += _feat_h_val + _gap
            if has_volcano:
                _volcano_y = _cur_y

            # "▲ Up-regulated" ラベル
            if has_up:
                lbl = slide_a.shapes.add_textbox(
                    Inches(0.3), Inches(_up_label_y), Inches(3), Inches(0.3))
                lp = lbl.text_frame.paragraphs[0]
                lp.text = f"▲ Up-regulated (Top {len(up_features)})"
                lp.font.size = Pt(10)
                lp.font.color.rgb = RGBColor(0xFF, 0x2D, 0x2D)
                lp.font.bold = True

            # Feature Plot 画像配置 — Up
            for i, feat in enumerate(up_features):
                is_last_up = (i == len(up_features) - 1)
                feat_fig = _build_feature_plot_fig(
                    df, feat, cache_dir_path, rds_path_str,
                    rotation_store, name_map, marker_size=3,
                    show_colorbar_title=is_last_up,
                    auto_marker_size=True)
                if feat_fig:
                    png = _fig_to_png_bytes(
                        feat_fig, width=400, height=280, scale=2)
                    left = Inches(0.3 + i * (12.5 / max(top_n, 1)))
                    _pptx_add_image_preserve_ratio(
                        slide_a, png,
                        int(left), Inches(_up_plot_y),
                        feat_w, feat_h,
                        png_w=400, png_h=280)

            # "▼ Down-regulated" ラベル
            if has_down:
                lbl = slide_a.shapes.add_textbox(
                    Inches(0.3), Inches(_down_label_y), Inches(3), Inches(0.3))
                lp = lbl.text_frame.paragraphs[0]
                lp.text = f"▼ Down-regulated (Top {len(down_features)})"
                lp.font.size = Pt(10)
                lp.font.color.rgb = RGBColor(0x1E, 0x5B, 0xFF)
                lp.font.bold = True

            # Feature Plot 画像配置 — Down
            for i, feat in enumerate(down_features):
                is_last_down = (i == len(down_features) - 1)
                feat_fig = _build_feature_plot_fig(
                    df, feat, cache_dir_path, rds_path_str,
                    rotation_store, name_map, marker_size=3,
                    show_colorbar_title=is_last_down,
                    auto_marker_size=True)
                if feat_fig:
                    png = _fig_to_png_bytes(
                        feat_fig, width=400, height=280, scale=2)
                    left = Inches(0.3 + i * (12.5 / max(top_n, 1)))
                    _pptx_add_image_preserve_ratio(
                        slide_a, png,
                        int(left), Inches(_down_plot_y),
                        feat_w, feat_h,
                        png_w=400, png_h=280)

            # Volcano Plot (下段) — XY比を保持して最大サイズで配置
            if has_volcano:
                vpng = _fig_to_png_bytes(volcano_cl, width=800, height=700, scale=2)
                v_aspect = 800 / 700  # ≈ 1.14
                v_height = Inches(_volcano_h)
                v_width = int(v_height * v_aspect)
                v_left = int((prs.slide_width - v_width) / 2)
                _pptx_add_image(slide_a, vpng,
                                v_left, Inches(_volcano_y), v_width, v_height)
            elif not up_features and not down_features:
                # DEG データなし → 注釈
                no_deg = slide_a.shapes.add_textbox(
                    Inches(3), Inches(3), Inches(7), Inches(1))
                np_ = no_deg.text_frame.paragraphs[0]
                np_.text = "No DEG data available for this cluster"
                np_.font.size = Pt(16)
                np_.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
                np_.alignment = PP_ALIGN.CENTER

            _progress(f"Cluster {cl_str} — DEG")

            # === Slide B: UMAP + Spatial (上下2段・サンプル別) ===
            slide_b = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_add_title_bar(slide_b, f"Cluster {cl_str} — UMAP & Spatial")

            n_sp_b = len(samples) if has_spatial and samples else 0
            avail_w_b = 12.7  # 0.3〜13.0"

            if n_sp_b > 0:
                tile_w_b = avail_w_b / n_sp_b

                # 上段: サンプル別UMAP (y=0.9", h=3.0")
                for idx, s in enumerate(samples):
                    df_s = df[df["Sample"] == s]
                    umap_s = _build_umap_integrated_fig(
                        df_s, color_by="Cluster",
                        highlight_clusters=[cl_str],
                        show_legend=False, show_labels=False,
                        title=_display_name(s, name_map),
                        marker_size=3, custom_colors=custom_colors,
                        bg_opacity=1.0, title_font_size=40)
                    if umap_s is not None:
                        umap_s.update_xaxes(range=cl_umap_xrange)
                        umap_s.update_yaxes(range=cl_umap_yrange)
                        u_dict = (umap_s.to_dict()
                                  if hasattr(umap_s, "to_dict") else umap_s)
                        u_png = _fig_to_png_bytes(
                            u_dict, width=600, height=600, scale=2)
                        _buw, _buh, _buoff = _square_tile_dims(
                            tile_w_b, 3.0)
                        u_left = Inches(0.3 + idx * tile_w_b + _buoff)
                        _pptx_add_image(slide_b, u_png,
                                        int(u_left), Inches(0.9),
                                        Inches(_buw), Inches(_buh))

                # 下段: サンプル別Spatial (y=4.1", h=3.1")
                for idx, s in enumerate(samples):
                    df_s = df[df["Sample"] == s]
                    transform = rotation_store.get(
                        s, rotation_store.get(
                            "__all__",
                            {"angle": 0, "flip_h": False, "flip_v": False}))
                    if isinstance(transform, (int, float)):
                        transform = {"angle": int(transform),
                                     "flip_h": False, "flip_v": False}
                    sp_fig = _create_single_spatial_fig(
                        df_s, color_map,
                        highlight_clusters=[cl_str],
                        selected_cell_ids=set(),
                        rotation_deg=transform.get("angle", 0),
                        show_labels=False,
                        flip_h=transform.get("flip_h", False),
                        flip_v=transform.get("flip_v", False),
                        title=_display_name(s, name_map),
                        marker_size=0, render_height=560,
                        title_font_size=40,
                        embed_legend=(idx == 0))
                    if sp_fig is not None:
                        sp_dict = (sp_fig.to_dict()
                                   if hasattr(sp_fig, "to_dict") else sp_fig)
                        sp_png = _fig_to_png_bytes(
                            sp_dict, width=600, height=600, scale=2)
                        _bsw, _bsh, _bsoff = _square_tile_dims(
                            tile_w_b, 3.1)
                        sp_left = Inches(0.3 + idx * tile_w_b + _bsoff)
                        _pptx_add_image(slide_b, sp_png,
                                        int(sp_left), Inches(4.1),
                                        Inches(_bsw), Inches(_bsh))
            else:
                # Spatialデータなし → 単一UMAP（従来互換）
                umap_hl = _build_umap_integrated_fig(
                    df, color_by="Cluster", highlight_clusters=[cl_str],
                    show_legend=True, show_labels=False,
                    marker_size=3, custom_colors=custom_colors,
                    bg_opacity=1.0)
                if umap_hl is not None:
                    umap_dict = (umap_hl.to_dict()
                                 if hasattr(umap_hl, "to_dict") else umap_hl)
                    upng = _fig_to_png_bytes(
                        umap_dict, width=800, height=800, scale=2)
                    _pptx_add_image(slide_b, upng,
                                    Inches(0.3), Inches(0.7),
                                    Inches(4.5), Inches(4.5))

            _progress(f"Cluster {cl_str} — UMAP/Spatial")

    # existing_prs が渡された場合は呼び出し元がまとめて保存するため
    # ここでは保存しない (現在のステップ数を返す)
    if existing_prs is not None:
        return _current_step[0]

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


@callback(
    Output("export_top_n_store", "data"),
    Input("input_export_top_n", "value"),
    prevent_initial_call=True,
)
def sync_export_top_n(value):
    """dbc.Input → dcc.Store ブリッジ (Top N)"""
    return value or 5


@callback(
    [Output("export_method_selector", "options"),
     Output("export_method_selector", "value")],
    Input("interactive_rds_map", "data"),
    prevent_initial_call=True,
)
def update_export_method_options(rds_map):
    """rds_map の変更に応じてエクスポート対象手法セレクタを更新する。"""
    if not rds_map or not isinstance(rds_map, dict):
        return [{"label": "All", "value": "all"}], "all"

    methods = list(rds_map.keys())
    options = [{"label": m, "value": m} for m in methods]
    if len(methods) > 1:
        options.append({"label": "Both", "value": "all"})

    default_val = "all" if len(methods) > 1 else methods[0]
    return options, default_val


@callback(
    [Output("dl_report_pptx", "data"),
     Output("div_export_status", "children")],
    Input("btn_export_report", "n_clicks"),
    [State("interactive_umap_plot", "figure"),
     State("last_spatial_figure_store", "data"),
     State("seurat_rds_path_store", "data"),
     State("cluster_stats_table", "data"),
     State("interactive_sub_project_select", "options"),
     State("interactive_sub_project_select", "value"),
     State("volcano_plot", "figure"),
     State("heatmap_plot", "figure"),
     State("deg_data_store", "data"),
     State("custom_color_map_store", "data"),
     State("spatial_rotation_store", "data"),
     State("sample_name_map_store", "data"),
     State("export_top_n_store", "data"),
     State("seurat_cache_dir_store", "data"),
     State("mrm_path", "value"),
     State("interactive_rds_map", "data"),
     State("interactive_result_folder", "value"),
     State("interactive_integration_method", "value"),
     State("export_method_selector", "value")],
    background=True,
    running=[
        (Output("btn_export_report", "disabled"), True, False),
        (Output("export_progress_container", "style"),
         {"display": "block"}, {"display": "none"}),
    ],
    progress=[
        Output("export_progress_bar", "value"),
        Output("export_progress_bar", "max"),
        Output("export_progress_label", "children"),
    ],
    prevent_initial_call=True,
)
def cb_export_report(set_progress, n_clicks, umap_fig, spatial_fig, rds_path,
                     cluster_stats_data, sub_options, sub_value,
                     volcano_fig, heatmap_fig, deg_data, custom_colors,
                     rotation_store, name_map, top_n, cache_dir_str,
                     mrm_path_str, rds_map, result_folder, current_method,
                     export_method_selection):
    """PPTX レポートをバックグラウンド生成してダウンロード。

    export_method_selection:
        "all" → 全手法（比較スライド付き）
        特定手法名 → その手法のみ
    """
    if not n_clicks:
        raise PreventUpdate

    set_progress((0, 100, "準備中..."))

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        return no_update, (
            "python-pptx がインストールされていません。"
            "pip install python-pptx を実行してください。"
        )

    try:
        import kaleido  # noqa: F401
    except ImportError:
        return no_update, (
            "kaleido がインストールされていません。"
            "pip install kaleido を実行してください。"
        )

    if not umap_fig:
        return no_update, "UMAPプロットが見つかりません。データを読み込んでください。"

    try:
        # サブプロジェクト名を取得
        sub_name = ""
        if sub_options and sub_value:
            for opt in sub_options:
                if opt.get("value") == sub_value:
                    sub_name = opt.get("label", "")
                    break

        # ファイル名を決定
        if sub_name:
            safe_name = re.sub(r'[\\/*?:"<>|]', '_', sub_name)
            filename = f"{safe_name}.pptx"
        else:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"MSI_Report_{timestamp}.pptx"

        top_n = top_n or 5
        saved_positions = _load_label_positions()

        # ------------------------------------------------------------------
        # 出力対象手法リストの決定（export_method_selector に基づく）
        # ------------------------------------------------------------------
        methods_to_export = []
        if rds_map and isinstance(rds_map, dict):
            if export_method_selection and export_method_selection != "all":
                # 特定手法のみ選択
                if export_method_selection in rds_map:
                    methods_to_export = [export_method_selection]
                else:
                    methods_to_export = []
            elif len(rds_map) > 1:
                # "all" → 全手法
                if current_method and current_method in rds_map:
                    methods_to_export = [current_method] + [
                        m for m in rds_map if m != current_method
                    ]
                else:
                    methods_to_export = list(rds_map.keys())

        # ------------------------------------------------------------------
        # 単一手法の場合（従来の動作と完全互換）
        # ------------------------------------------------------------------
        if not methods_to_export:
            cache_dir_path = Path(cache_dir_str) if cache_dir_str else None
            df = None
            meta = {}
            if cache_dir_path:
                plot_parquet = cache_dir_path / "plot_data.parquet"
                plot_csv = cache_dir_path / "plot_data.csv"
                if plot_parquet.exists():
                    df = pd.read_parquet(plot_parquet)
                elif plot_csv.exists():
                    df = pd.read_csv(plot_csv)

                meta_file = cache_dir_path / "extraction_meta.json"
                if meta_file.exists():
                    with open(meta_file, "r", encoding="utf-8") as f:
                        meta = json.load(f)

            pptx_bytes = _build_pptx(
                umap_fig, spatial_fig, meta, cluster_stats_data, rds_path,
                sub_name=sub_name, volcano_fig=volcano_fig,
                heatmap_fig=heatmap_fig,
                deg_data=deg_data, top_n=top_n, df=df,
                cache_dir=str(cache_dir_path) if cache_dir_path else None,
                custom_colors=custom_colors, rotation_store=rotation_store,
                name_map=name_map, set_progress=set_progress,
                mrm_path=mrm_path_str,
                saved_positions=saved_positions,
            )

            return (
                dcc.send_bytes(pptx_bytes, filename=filename),
                f"✓ PPTXファイルを出力しました: {filename}",
            )

        # ------------------------------------------------------------------
        # 複数手法 or セレクタで指定された手法 → 1つの PPTX に結合
        # ------------------------------------------------------------------
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        is_multi = len(methods_to_export) > 1

        # 全手法のステップ数（Phase 1 後に再計算する）
        total_steps = len(methods_to_export) * 2  # 暫定: データ読込用
        progress_offset = 0
        exported_methods = []
        section_map = []  # (name, start_idx, end_idx) — ①-5

        # ==================================================================
        # Phase 1: 全手法のデータを事前抽出（キャッシュして二重抽出を防止）
        # ==================================================================
        extracted_data = {}  # method_name → dict

        for method_name in methods_to_export:
            method_rds = rds_map.get(method_name)
            if not method_rds or not Path(method_rds).exists():
                print(f"[Export] {method_name}: "
                      f"RDSファイルが見つかりません → スキップ")
                continue

            set_progress((
                min(int(progress_offset / total_steps * 100), 99), 100,
                f"{method_name} のデータを読み込み中..."
            ))

            try:
                result = _bridge.extract_data(method_rds)
            except Exception as e:
                print(f"[Export] {method_name}: データ抽出エラー: {e}")
                continue

            method_df = result["plot_data"]
            method_meta = result["meta"]
            method_cache_dir = result.get("cache_dir")

            # DEG 結果読み込み
            method_deg_data = None
            if result_folder:
                result_base = Path(result_folder)
                method_deg_data = _load_deg_results(
                    result_base, method_name)
            else:
                rds_dir = Path(method_rds).parent
                result_base = (rds_dir.parent
                               if rds_dir.name == "RDS_Files"
                               else rds_dir)
                method_deg_data = _load_deg_results(
                    result_base, method_name)

            extracted_data[method_name] = {
                "df": method_df,
                "meta": method_meta,
                "cache_dir": method_cache_dir,
                "deg_data": method_deg_data,
                "rds_path": method_rds,
            }
            progress_offset += 1

        if not extracted_data:
            return no_update, "エクスポート可能な手法がありません。"

        # Phase 1 完了: 実際のクラスタ数に基づいて total_steps を再計算
        total_steps = progress_offset                    # Phase 1 消費済み
        if is_multi:
            total_steps += len(extracted_data) * 2        # Phase 2: 比較 + sample UMAP
        for _ed in extracted_data.values():
            total_steps += 1                              # セパレータスライド
            _method_df = _ed["df"]
            if _method_df is not None:
                _n_clusters = len(_method_df["Cluster"].unique())
                total_steps += 5 + _n_clusters * 2       # _build_pptx 実ステップ (+1 sample UMAP)
            else:
                total_steps += 4

        # ==================================================================
        # Phase 2: 比較セクション — 全手法の UMAP & Spatial を先頭に配置
        #          (①-3: 複数手法の場合のみ)
        # ==================================================================
        if is_multi and len(extracted_data) > 1:
            comparison_start = len(prs.slides)

            for method_name, ed in extracted_data.items():
                m_df = ed["df"]
                if m_df is None or m_df.empty:
                    continue
                has_spatial_cmp = "SpatialX" in m_df.columns

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _pptx_add_title_bar(
                    slide,
                    f"UMAP & Spatial Mapping — {method_name}")

                color_map_cmp = _get_cluster_color_map(
                    m_df["Cluster"], custom_colors)
                all_samples_cmp = sorted(m_df["Sample"].unique())
                n_sp_cmp = len(all_samples_cmp)
                avail_w_cmp = 11.0  # 右端にレジェンド配置

                # UMAP axis ranges (全サンプル統一)
                u1_min = float(m_df["UMAP_1"].min())
                u1_max = float(m_df["UMAP_1"].max())
                u2_min = float(m_df["UMAP_2"].min())
                u2_max = float(m_df["UMAP_2"].max())
                u_pad = max(u1_max - u1_min,
                            u2_max - u2_min) * 0.05
                u_xrange = [u1_min - u_pad, u1_max + u_pad]
                u_yrange = [u2_min - u_pad, u2_max + u_pad]

                # 上段: サンプル別 UMAP
                tile_w_cmp = avail_w_cmp / max(n_sp_cmp, 1)
                _umap_pos_cmp = (saved_positions or {}).get("umap_integrated", {})
                for idx_s, s in enumerate(all_samples_cmp):
                    df_s = m_df[m_df["Sample"] == s]
                    umap_s = _build_umap_integrated_fig(
                        df_s, color_by="Cluster",
                        highlight_clusters=None,
                        show_legend=False, show_labels=True,
                        title=_display_name(s, name_map),
                        marker_size=3, custom_colors=custom_colors,
                        title_font_size=40, label_size=24,
                        saved_positions=_umap_pos_cmp)
                    if umap_s is not None:
                        umap_s.update_xaxes(range=u_xrange)
                        umap_s.update_yaxes(range=u_yrange)
                        u_dict = (umap_s.to_dict()
                                  if hasattr(umap_s, "to_dict")
                                  else umap_s)
                        u_png = _fig_to_png_bytes(
                            u_dict, width=600, height=600, scale=2)
                        _cw, _ch, _coff = _square_tile_dims(
                            tile_w_cmp, 3.0)
                        u_left = Inches(
                            0.3 + idx_s * tile_w_cmp + _coff)
                        _pptx_add_image(slide, u_png,
                                        int(u_left), Inches(0.9),
                                        Inches(_cw), Inches(_ch))

                # 下段: サンプル別 Spatial
                if has_spatial_cmp:
                    _rot_store = rotation_store or {}
                    for idx_s, s in enumerate(all_samples_cmp):
                        df_s = m_df[m_df["Sample"] == s]
                        transform = _rot_store.get(
                            s, _rot_store.get(
                                "__all__",
                                {"angle": 0, "flip_h": False,
                                 "flip_v": False}))
                        if isinstance(transform, (int, float)):
                            transform = {
                                "angle": int(transform),
                                "flip_h": False, "flip_v": False}
                        _sp_pos_cmp = (saved_positions or {}).get("spatial", {}).get(s, {})
                        sp_fig_cmp = _create_single_spatial_fig(
                            df_s, color_map_cmp,
                            highlight_clusters=None,
                            selected_cell_ids=set(),
                            rotation_deg=transform.get("angle", 0),
                            show_labels=True,
                            flip_h=transform.get("flip_h", False),
                            flip_v=transform.get("flip_v", False),
                            title=_display_name(s, name_map),
                            marker_size=0, render_height=560,
                            embed_legend=False,
                            title_font_size=40, label_size=24,
                            saved_positions=_sp_pos_cmp)
                        if sp_fig_cmp is not None:
                            sp_dict = (sp_fig_cmp.to_dict()
                                       if hasattr(sp_fig_cmp, "to_dict")
                                       else sp_fig_cmp)
                            sp_png = _fig_to_png_bytes(
                                sp_dict, width=600, height=600,
                                scale=2)
                            _csw, _csh, _csoff = _square_tile_dims(
                                tile_w_cmp, 3.1)
                            sp_left = Inches(
                                0.3 + idx_s * tile_w_cmp + _csoff)
                            _pptx_add_image(
                                slide, sp_png,
                                int(sp_left), Inches(4.1),
                                Inches(_csw), Inches(_csh))

                # クラスタレジェンド（右端に配置）
                legend_fig_cmp = _build_cluster_legend_fig(
                    m_df["Cluster"].unique(), color_map_cmp)
                legend_png_cmp = _fig_to_png_bytes(
                    legend_fig_cmp.to_dict(),
                    width=200, height=600, scale=2)
                _pptx_add_image_preserve_ratio(slide, legend_png_cmp,
                                               Inches(11.5), Inches(0.9),
                                               Inches(1.3), Inches(6.3),
                                               png_w=200, png_h=600)

                # --- UMAP (by Sample) comparison slide ---
                if n_sp_cmp > 1:
                    sample_cmap_cmp = _get_sample_color_map(all_samples_cmp)
                    slide_s_cmp = prs.slides.add_slide(prs.slide_layouts[6])
                    _pptx_add_title_bar(
                        slide_s_cmp,
                        f"UMAP (by Sample) \u2014 {method_name}")

                    for idx_s, s in enumerate(all_samples_cmp):
                        df_s = m_df[m_df["Sample"] == s]
                        umap_s_cmp = _build_umap_integrated_fig(
                            df_s, color_by="Sample",
                            highlight_clusters=None,
                            show_legend=False, show_labels=False,
                            title=_display_name(s, name_map),
                            marker_size=3, custom_colors=sample_cmap_cmp,
                            title_font_size=40)
                        if umap_s_cmp is not None:
                            umap_s_cmp.update_xaxes(range=u_xrange)
                            umap_s_cmp.update_yaxes(range=u_yrange)
                            u_dict = (umap_s_cmp.to_dict()
                                      if hasattr(umap_s_cmp, "to_dict")
                                      else umap_s_cmp)
                            u_png = _fig_to_png_bytes(
                                u_dict, width=600, height=600, scale=2)
                            _cw, _ch, _coff = _square_tile_dims(
                                tile_w_cmp, 5.5)
                            u_left = Inches(
                                0.3 + idx_s * tile_w_cmp + _coff)
                            _pptx_add_image(slide_s_cmp, u_png,
                                            int(u_left), Inches(0.9),
                                            Inches(_cw), Inches(_ch))

                    # サンプルレジェンド（右端）
                    sample_legend_cmp = _build_sample_legend_fig(
                        all_samples_cmp, sample_cmap_cmp,
                        name_map=name_map)
                    sample_legend_png_cmp = _fig_to_png_bytes(
                        sample_legend_cmp.to_dict(),
                        width=200, height=600, scale=2)
                    _pptx_add_image_preserve_ratio(
                        slide_s_cmp, sample_legend_png_cmp,
                        Inches(11.5), Inches(0.9),
                        Inches(1.3), Inches(6.3),
                        png_w=200, png_h=600)

                progress_offset += 1

            comparison_end = len(prs.slides) - 1
            if comparison_end >= comparison_start:
                section_map.append(
                    ("Comparison", comparison_start, comparison_end))

        # ==================================================================
        # Phase 3: 各手法のフルセクション
        # ==================================================================
        for method_name in methods_to_export:
            if method_name not in extracted_data:
                continue

            ed = extracted_data[method_name]
            method_df = ed["df"]
            method_meta = ed["meta"]
            method_cache_dir = ed["cache_dir"]
            method_deg_data = ed["deg_data"]
            method_rds = ed["rds_path"]

            method_start_idx = len(prs.slides)

            set_progress((
                min(int(progress_offset / total_steps * 100), 99),
                100,
                f"{method_name} のスライドを生成中..."
            ))

            # --- セパレータスライド ---
            sep_slide = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_add_title_bar(sep_slide, f"═══ {method_name} ═══")
            txBox = sep_slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(11), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"Integration Method: {method_name}"
            p.font.size = Pt(28)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            p2 = tf.add_paragraph()
            p2.text = (
                f"Cells: {method_meta.get('n_cells', '?')} | "
                f"Clusters: {method_meta.get('n_clusters', '?')}"
            )
            p2.font.size = Pt(18)
            p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            p2.alignment = PP_ALIGN.CENTER
            progress_offset += 1

            # --- UMAP 図を生成 ---
            _umap_pos_m = (saved_positions or {}).get("umap_integrated", {})
            method_umap_fig = _build_umap_integrated_fig(
                method_df, color_by="Cluster",
                highlight_clusters=None,
                show_legend=True, show_labels=True,
                custom_colors=custom_colors,
                saved_positions=_umap_pos_m,
            )

            # --- クラスタ統計の生成 ---
            method_cluster_stats = []
            try:
                clusters_sorted = sorted(
                    method_df["Cluster"].unique(),
                    key=_cluster_sort_key)
                n_total = len(method_df)
                for c in clusters_sorted:
                    n_c = int((method_df["Cluster"] == c).sum())
                    pct = (f"{n_c / n_total * 100:.1f}"
                           if n_total else "0")
                    method_cluster_stats.append({
                        "Cluster": str(c),
                        "Pixels": n_c,
                        "Percent": pct,
                    })
            except Exception:
                method_cluster_stats = []

            # --- _build_pptx でフルセットを追加 ---
            method_sub_name = (
                f"{sub_name} [{method_name}]"
                if sub_name else method_name
            )
            returned = _build_pptx(
                method_umap_fig, None, method_meta,
                method_cluster_stats, method_rds,
                sub_name=method_sub_name,
                volcano_fig=None, heatmap_fig=None,
                deg_data=method_deg_data, top_n=top_n,
                df=method_df,
                cache_dir=(str(method_cache_dir)
                           if method_cache_dir else None),
                custom_colors=custom_colors,
                rotation_store=rotation_store,
                name_map=name_map,
                set_progress=set_progress,
                mrm_path=mrm_path_str,
                existing_prs=prs,
                progress_offset=progress_offset,
                progress_total=total_steps,
                saved_positions=saved_positions,
            )
            if isinstance(returned, int):
                progress_offset = returned

            method_end_idx = len(prs.slides) - 1
            section_map.append(
                (method_name, method_start_idx, method_end_idx))
            exported_methods.append(method_name)

        if not exported_methods:
            return no_update, "エクスポート可能な手法がありません。"

        # ==================================================================
        # Phase 4: セクション情報を PPTX XML に追加（①-5）
        # ==================================================================
        if section_map:
            try:
                _pptx_add_sections(prs, section_map)
            except Exception as e:
                print(f"[Export] セクション追加エラー（無視）: {e}")

        # --- 結合 PPTX をバイト列に変換 ---
        output = BytesIO()
        prs.save(output)
        output.seek(0)

        methods_str = " + ".join(exported_methods)
        return (
            dcc.send_bytes(output.getvalue(), filename=filename),
            f"✓ PPTXファイルを出力しました ({methods_str}): {filename}",
        )

    except Exception as e:
        return no_update, f"エクスポートエラー: {e}"


# ---------------------------------------------------------------------------
# プロジェクト / サブプロジェクト選択コールバック
# ---------------------------------------------------------------------------

@callback(
    Output("interactive_project_row", "style"),
    Input("interactive_entry_mode", "data"),
    prevent_initial_call=True,
)
def toggle_project_dropdown_visibility(entry_mode):
    """entry_mode に応じてプロジェクトドロップダウンの表示/非表示を切り替え"""
    if entry_mode == "sub_project":
        return {"display": "none"}
    return {}


@callback(
    Output("interactive_project_select", "options"),
    [Input("main_tabs", "active_tab"),
     Input("current_page", "data")],
    prevent_initial_call=True,
)
def populate_interactive_projects(active_tab, current_page):
    """interactiveタブがアクティブになった時にプロジェクト一覧を取得"""
    if current_page != "analysis" or active_tab != "interactive":
        return no_update
    from app.services.project_manager import list_projects
    projects = list_projects()
    return [{"label": p["name"], "value": p["id"]} for p in projects]


@callback(
    [Output("interactive_sub_project_select", "options", allow_duplicate=True),
     Output("interactive_sub_project_select", "value", allow_duplicate=True)],
    Input("interactive_project_select", "value"),
    State("interactive_entry_mode", "data"),
    prevent_initial_call=True,
)
def populate_interactive_sub_projects(project_id, entry_mode):
    """プロジェクト選択時にサブプロジェクト一覧を取得。
    sub_project モード: options のみ更新（value は sub_action_interactive が設定済み）
    standalone モード: options を更新、value は None（ユーザーが選択）"""
    if not project_id:
        return [], None
    from app.services.project_manager import list_sub_projects
    subs = list_sub_projects(project_id)
    options = [{"label": s["name"], "value": s["id"]} for s in subs]
    if entry_mode == "sub_project":
        return options, no_update
    return options, None


@callback(
    [Output("interactive_viz_container", "style", allow_duplicate=True),
     Output("interactive_data_info", "children", allow_duplicate=True),
     Output("sap_skip_reset", "data", allow_duplicate=True),
     Output("sap_btn_wrapper", "style", allow_duplicate=True)],
    Input("interactive_project_select", "value"),
    State("sap_skip_reset", "data"),
    prevent_initial_call=True,
)
def reset_interactive_on_project_change(project_id, skip_reset):
    """プロジェクト変更時にインタラクティブデータをリセット。
    sub_project_select の value が同一でも確実にクリアされる。
    sap_skip_reset=True の場合はリセットをスキップ（保存後の自動切替時）。"""
    if skip_reset:
        return no_update, no_update, False, no_update

    _interactive_data["plot_data"] = None
    _interactive_data["cluster_stats"] = None
    _interactive_data["features_list"] = None
    _interactive_data["meta"] = None
    _interactive_data["rds_path"] = None
    _interactive_data["cache_dir"] = None

    if not project_id:
        return {"display": "none"}, "", False, {"display": "none"}
    return {"display": "none"}, "データを読み込んでください", False, {"display": "none"}


@callback(
    [Output("interactive_result_folder", "value", allow_duplicate=True),
     Output("interactive_msi_folder", "value", allow_duplicate=True),
     Output("interactive_data_info", "children", allow_duplicate=True),
     Output("int_cal_ms_instrument", "data"),
     Output("interactive_viz_container", "style", allow_duplicate=True),
     Output("sap_skip_reset", "data", allow_duplicate=True),
     Output("sap_btn_wrapper", "style", allow_duplicate=True)],
    Input("interactive_sub_project_select", "value"),
    [State("interactive_project_select", "value"),
     State("sap_skip_reset", "data")],
    prevent_initial_call=True,
)
def set_interactive_folders_from_sub_project(sub_id, project_id, skip_reset):
    """サブプロジェクト選択時にフォルダパスを自動設定 + データリセット
    sap_skip_reset=True の場合はリセットをスキップ（保存後の自動切替時）。"""
    if skip_reset:
        return (no_update,) * 6 + (False,)

    # 前のプロジェクトのデータをクリア
    _interactive_data["plot_data"] = None
    _interactive_data["cluster_stats"] = None
    _interactive_data["features_list"] = None
    _interactive_data["meta"] = None
    _interactive_data["rds_path"] = None
    _interactive_data["cache_dir"] = None

    if not sub_id or not project_id:
        return no_update, no_update, no_update, no_update, {"display": "none"}, False, {"display": "none"}
    from app.services.project_manager import get_sub_project
    sub = get_sub_project(project_id, sub_id)
    if not sub:
        return no_update, no_update, no_update, no_update, {"display": "none"}, False, {"display": "none"}
    result_dir = sub.get("last_result_dir") or sub.get("output_dir", "")
    data_folder = sub.get("data_folder", "")
    ms_instrument = sub.get("ms_instrument", "TIMS")
    # 未設定フォルダの警告メッセージ
    warnings = []
    if not result_dir:
        warnings.append("結果フォルダが未設定です")
    if not data_folder:
        warnings.append("MSIデータフォルダが未設定です")
    msg = "⚠ " + "、".join(warnings) if warnings else "データを読み込んでください"
    return (result_dir, data_folder, msg, ms_instrument, {"display": "none"}, False, {"display": "none"})


# ---------------------------------------------------------------------------
# フルスクリーン拡大モーダル
# ---------------------------------------------------------------------------

@callback(
    [Output("fullscreen_plot_modal", "is_open"),
     Output("fullscreen_modal_title", "children"),
     Output("fullscreen_modal_body", "children")],
    [Input("expand_umap_btn", "n_clicks"),
     Input("expand_feature_btn", "n_clicks"),
     Input("expand_spatial_btn", "n_clicks"),
     Input("expand_deg_btn", "n_clicks")],
    [State("interactive_umap_plot", "figure"),
     State("feature_plot_container", "children"),
     State("last_spatial_figure_store", "data"),
     State("deg_data_store", "data"),
     State("spatial_rotation_store", "data"),
     State("custom_color_map_store", "data"),
     State("spatial_columns_per_row", "value")],
    prevent_initial_call=True,
)
def toggle_fullscreen(umap_n, feat_n, spatial_n, deg_n,
                      umap_fig, feat_container_children, spatial_fig_data, deg_data,
                      rotation_store, custom_colors, spatial_columns_per_row):
    trigger = ctx.triggered_id
    if not trigger:
        return False, "", ""

    fs_graph_style = {"height": "80vh"}
    fs_config = {
        "scrollZoom": True,
        "edits": {"annotationPosition": True},
        "toImageButtonOptions": {"format": "png", "scale": 3},
    }

    # ===== UMAP (インタラクティブ) =====
    if trigger == "expand_umap_btn":
        df = _interactive_data.get("plot_data")
        if df is None:
            return False, "", ""

        clusters = sorted(df["Cluster"].unique(), key=_cluster_sort_key)
        cluster_opts = [{"label": f"Cluster {c}", "value": str(c)} for c in clusters]
        color_map = _get_cluster_color_map(df["Cluster"], custom_colors)

        # タイトル（RDSファイル名から生成）
        rds_path = _interactive_data.get("rds_path", "")
        umap_title = Path(rds_path).stem if rds_path else "UMAP"

        # 初期グラフ（統合モード）
        init_fig = _build_umap_integrated_fig(df, "Cluster", None, True, False,
                                               title=umap_title,
                                               custom_colors=custom_colors)
        init_fs_config = dict(fs_config)
        init_fs_config["toImageButtonOptions"] = dict(init_fs_config["toImageButtonOptions"],
                                                       filename=f"UMAP_{umap_title}")
        init_graph = dcc.Graph(figure=init_fig, style={"height": "78vh"}, config=init_fs_config)

        body = html.Div([
            dbc.Row(className="mb-2 align-items-center", children=[
                dbc.Col(width=2, children=[
                    dbc.RadioItems(id="fs_umap_display_mode",
                                   options=[{"label": "統合", "value": "integrated"},
                                            {"label": "サンプル別", "value": "per_sample"}],
                                   value="integrated", inline=True),
                ]),
                dbc.Col(width=2, children=[
                    dbc.RadioItems(id="fs_umap_color_by",
                                   options=[{"label": "Cluster", "value": "Cluster"},
                                            {"label": "Sample", "value": "Sample"}],
                                   value="Cluster", inline=True),
                ]),
                dbc.Col(width=2, children=[
                    dcc.Dropdown(id="fs_umap_highlight_cluster",
                                 options=cluster_opts, multi=True,
                                 placeholder="ハイライト"),
                ]),
                dbc.Col(width=3, children=[
                    dcc.Dropdown(id="fs_umap_exclude_cluster",
                                 options=cluster_opts, multi=True,
                                 placeholder="除去するクラスタ"),
                ]),
                dbc.Col(width=1, children=[
                    dbc.Checkbox(id="fs_umap_show_labels", label="ラベル", value=False),
                ]),
                dbc.Col(width=1, children=[
                    dbc.Checkbox(id="fs_umap_show_legend", label="凡例", value=True),
                ]),
            ]),
            dbc.Row(className="mt-1 align-items-center", children=[
                dbc.Col(width=2, children=[
                    dbc.Label("点サイズ", className="small mb-0"),
                    dcc.Slider(
                        id="fs_umap_marker_size",
                        min=1, max=10, step=1, value=2,
                        marks={1: "1", 5: "5", 10: "10"},
                        tooltip={"placement": "bottom", "always_visible": False},
                    ),
                ]),
                dbc.Col(width=2, children=[
                    dbc.Label("ラベルサイズ", className="small mb-0"),
                    dcc.Slider(
                        id="fs_umap_label_size",
                        min=6, max=24, step=1, value=14,
                        marks={6: "6", 14: "14", 24: "24"},
                        tooltip={"placement": "bottom", "always_visible": False},
                    ),
                ]),
                dbc.Col(width=2, children=[
                    dbc.Label("高さ", className="small mb-0"),
                    dcc.Slider(
                        id="fs_umap_height_slider",
                        min=40, max=95, step=5, value=78,
                        marks={40: "40", 78: "78", 95: "95"},
                        tooltip={"placement": "bottom", "always_visible": False},
                    ),
                ]),
                dbc.Col(width=2, children=[
                    dbc.Label("横幅", className="small mb-0"),
                    dcc.Slider(
                        id="fs_umap_width_slider",
                        min=40, max=100, step=5, value=95,
                        marks={40: "40", 70: "70", 95: "95"},
                        tooltip={"placement": "bottom", "always_visible": False},
                    ),
                ]),
                dbc.Col(width=2, className="d-flex align-items-end", children=[
                    dbc.Button("ラベル位置保存",
                               id={"type": "fs_save_label_pos", "index": "umap"},
                               size="sm", color="secondary", className="mb-1"),
                ]),
            ]),
            html.Div(id="fs_umap_graph_container", children=[init_graph]),
        ])
        return True, "UMAP", body

    # ===== Feature Plot (コンテナごと拡大) =====
    if trigger == "expand_feature_btn" and feat_container_children:
        return (
            True, "Feature Plot",
            html.Div(feat_container_children),
        )

    # ===== Spatial Mapping (インタラクティブ) =====
    if trigger == "expand_spatial_btn":
        df = _interactive_data.get("plot_data")
        if df is None or "SpatialX" not in df.columns:
            return False, "", ""

        samples = sorted(df["Sample"].unique())
        name_map = _interactive_data.get("_name_map") or {}
        sample_opts = [{"label": _display_name(s, name_map), "value": s} for s in samples]
        clusters = sorted(df["Cluster"].unique(), key=_cluster_sort_key)
        cluster_opts = [{"label": f"Cluster {c}", "value": str(c)} for c in clusters]
        color_map = _get_cluster_color_map(df["Cluster"], custom_colors)
        cluster_to_idx, discrete_cscale = _get_cluster_colorscale(df["Cluster"], custom_colors)

        # 初期グラフ（全サンプル、rotation_store適用）
        if not rotation_store:
            rotation_store = {}
        init_graphs = []
        for s in samples:
            display_s = _display_name(s, name_map)
            df_s = df[df["Sample"] == s]
            transform = rotation_store.get(
                s, rotation_store.get("__all__", {"angle": 0, "flip_h": False, "flip_v": False}))
            if isinstance(transform, (int, float)):
                transform = {"angle": int(transform), "flip_h": False, "flip_v": False}
            fig = _create_single_spatial_fig(df_s, color_map, None, set(),
                                             rotation_deg=transform.get("angle", 0),
                                             flip_h=transform.get("flip_h", False),
                                             flip_v=transform.get("flip_v", False),
                                             title=display_s, embed_legend=True,
                                             cluster_to_idx=cluster_to_idx,
                                             discrete_cscale=discrete_cscale)
            if spatial_columns_per_row:
                n_cols = spatial_columns_per_row
                gap_total = (n_cols - 1) * 15
                flex_basis = f"calc({100 / n_cols:.2f}% - {gap_total / n_cols:.1f}px)"
                min_w = "0"
            else:
                n_cols = len(samples)
                flex_basis = f"{max(20, 90 // n_cols)}%"
                min_w = "350px"
            init_cfg = dict(fs_config)
            init_cfg["toImageButtonOptions"] = dict(init_cfg["toImageButtonOptions"],
                                                     filename=f"Spatial_{display_s}")
            init_graphs.append(
                html.Div(
                    style={"flex": f"1 1 {flex_basis}", "minWidth": min_w,
                            "border": "1px solid #dee2e6", "borderRadius": "6px",
                            "padding": "5px", "backgroundColor": "#fff"},
                    children=[
                        dcc.Graph(figure=fig, style={"height": "60vh"}, config=init_cfg),
                    ],
                )
            )
        init_container = html.Div(
            style={"display": "flex", "flexWrap": "wrap", "gap": "15px"},
            children=init_graphs,
        )

        # フルスクリーン用サンプル別コントロール（1つのAccordionItemに統合）
        name_map = _interactive_data.get("_name_map") or {}
        fs_all_controls = []
        for i, s in enumerate(samples):
            t = rotation_store.get(
                s, rotation_store.get("__all__", {"angle": 0, "flip_h": False, "flip_v": False}))
            if isinstance(t, (int, float)):
                t = {"angle": int(t), "flip_h": False, "flip_v": False}
            display_s = _display_name(s, name_map)
            fs_all_controls.append(
                html.Div(
                    style={"padding": "4px 8px"},
                    children=[
                        html.Label(display_s or s, className="fw-bold small mb-1"),
                        dcc.Slider(
                            id={"type": "per_sample_rotation", "index": s},
                            min=0, max=270, step=90,
                            value=t.get("angle", 0),
                            marks={0: "0°", 90: "90°", 180: "180°", 270: "270°"},
                        ),
                        html.Div(className="d-flex gap-2 justify-content-center", children=[
                            dbc.Checkbox(
                                id={"type": "per_sample_flip_h", "index": s},
                                label="↔ 左右", value=t.get("flip_h", False),
                            ),
                            dbc.Checkbox(
                                id={"type": "per_sample_flip_v", "index": s},
                                label="↕ 上下", value=t.get("flip_v", False),
                            ),
                        ]),
                        html.Hr(className="my-1") if i < len(samples) - 1 else html.Div(),
                    ],
                )
            )
        fs_accordion_items = [dbc.AccordionItem(title="回転/反転", children=fs_all_controls)]

        body = html.Div([
            dbc.Row(className="mb-2 align-items-center", children=[
                dbc.Col(width=2, children=[
                    dcc.Dropdown(id="fs_spatial_sample", options=sample_opts,
                                 placeholder="サンプル(空=全表示)", clearable=True),
                ]),
                dbc.Col(width=2, children=[
                    dcc.Dropdown(id="fs_spatial_highlight_cluster",
                                 options=cluster_opts, multi=True,
                                 placeholder="ハイライト"),
                ]),
                dbc.Col(width=2, children=[
                    dcc.Dropdown(id="fs_spatial_exclude_cluster",
                                 options=cluster_opts, multi=True,
                                 placeholder="除去"),
                ]),
                dbc.Col(width=1, children=[
                    dbc.Checkbox(id="fs_spatial_show_labels", label="番号", value=False),
                ]),
            ]),
            dbc.Row(className="mt-1 align-items-center", children=[
                dbc.Col(width=2, children=[
                    html.Div(style={"display": "flex", "alignItems": "center", "gap": "4px"}, children=[
                        dbc.Label("マーカー", className="small mb-0"),
                        dbc.Button("Auto", id="fs_spatial_marker_auto_btn",
                                   size="sm", outline=True, color="info",
                                   style={"padding": "0 5px", "fontSize": "10px",
                                          "lineHeight": "1.2"}),
                    ]),
                    dcc.Slider(
                        id="fs_spatial_marker_size",
                        min=0, max=30, step=1, value=0,
                        marks={0: "自動", 10: "10", 20: "20", 30: "30"},
                        tooltip={"placement": "bottom", "always_visible": False},
                    ),
                ]),
                dbc.Col(width=2, children=[
                    dbc.Label("ラベル", className="small mb-0"),
                    dcc.Slider(
                        id="fs_spatial_label_size",
                        min=6, max=24, step=1, value=10,
                        marks={6: "6", 14: "14", 24: "24"},
                        tooltip={"placement": "bottom", "always_visible": False},
                    ),
                ]),
                dbc.Col(width=2, children=[
                    dbc.Label("高さ", className="small mb-0"),
                    dcc.Slider(
                        id="fs_spatial_height_slider",
                        min=30, max=85, step=5, value=60,
                        marks={30: "30", 60: "60", 85: "85"},
                        tooltip={"placement": "bottom", "always_visible": False},
                    ),
                ]),
                dbc.Col(width=2, children=[
                    dbc.Label("横幅", className="small mb-0"),
                    dcc.Slider(
                        id="fs_spatial_width_slider",
                        min=40, max=100, step=5, value=95,
                        marks={40: "40", 70: "70", 95: "95"},
                        tooltip={"placement": "bottom", "always_visible": False},
                    ),
                ]),
                dbc.Col(width=2, className="d-flex align-items-end", children=[
                    dbc.Button("ラベル位置保存",
                               id={"type": "fs_save_label_pos", "index": "spatial"},
                               size="sm", color="secondary", className="mb-1"),
                ]),
            ]),
            dbc.Accordion(
                fs_accordion_items, start_collapsed=True,
                flush=True, always_open=True,
                style={"marginBottom": "8px"},
            ),
            html.Div(id="fs_spatial_graph_container", children=[init_container]),
        ])
        return True, "Spatial Mapping", body

    # ===== DEG テーブル + Volcano + Heatmap =====
    if trigger == "expand_deg_btn" and deg_data:
        # クラスタ選択肢
        deg_clusters = sorted(
            set(str(r.get("cluster", "")) for r in deg_data),
            key=_cluster_sort_key,
        )
        deg_cluster_opts = [{"label": f"Cluster {c}", "value": c} for c in deg_clusters]

        fs_deg_body = dbc.Tabs(active_tab="fs_deg_volcano_tab", children=[
            dbc.Tab(label="Volcano Plot", tab_id="fs_deg_volcano_tab", children=[
                html.P("メイン画面のVolcano Plotタブで操作してください。",
                       className="text-muted small mt-2"),
            ]),
            dbc.Tab(label="Heatmap", tab_id="fs_deg_heatmap_tab", children=[
                html.P("メイン画面のHeatmapタブで操作してください。",
                       className="text-muted small mt-2"),
            ]),
        ])
        return True, "DEG マーカー", fs_deg_body

    return False, "", ""


# ---------------------------------------------------------------------------
# フルスクリーン閉鎖 → メインプロット再描画トリガー
# ---------------------------------------------------------------------------

@callback(
    Output("fullscreen_closed_trigger", "data"),
    Input("fullscreen_plot_modal", "is_open"),
    State("fullscreen_closed_trigger", "data"),
    prevent_initial_call=True,
)
def on_fullscreen_close(is_open, current_val):
    """フルスクリーンモーダルが閉じた時にトリガー値をインクリメントし、
    メインプロットの再描画をトリガーする"""
    if not is_open:
        return (current_val or 0) + 1
    return no_update


# ---------------------------------------------------------------------------
# フルスクリーン UMAP インタラクティブ更新
# ---------------------------------------------------------------------------

@callback(
    Output("fs_umap_graph_container", "children"),
    [Input("fs_umap_display_mode", "value"),
     Input("fs_umap_color_by", "value"),
     Input("fs_umap_highlight_cluster", "value"),
     Input("fs_umap_show_labels", "value"),
     Input("fs_umap_show_legend", "value"),
     Input("fs_umap_height_slider", "value"),
     Input("fs_umap_width_slider", "value"),
     Input("fs_umap_marker_size", "value"),
     Input("fs_umap_exclude_cluster", "value"),
     Input("fs_umap_label_size", "value")],
    [State("custom_color_map_store", "data"),
     State("umap_columns_per_row", "value"),
     State("accumulated_label_positions", "data")],
    prevent_initial_call=True,
)
def update_fs_umap(display_mode, color_by, highlight, show_labels, show_legend,
                   height_val, width_val, marker_size, exclude_clusters, label_size,
                   custom_color_map, columns_per_row, accumulated_positions):
    height_val = height_val or 78
    width_val = width_val or 95
    df = _interactive_data.get("plot_data")
    if df is None:
        return ""
    custom_colors = custom_color_map if custom_color_map else None
    color_map = _get_cluster_color_map(df["Cluster"], custom_colors)
    fs_config = {"scrollZoom": True, "edits": {"annotationPosition": True}, "toImageButtonOptions": {"format": "png", "scale": 3}}
    all_pos = _get_merged_label_positions(accumulated_positions)

    # タイトル（RDSファイル名から生成）
    rds_path = _interactive_data.get("rds_path", "")
    umap_title = Path(rds_path).stem if rds_path else "UMAP"

    if display_mode == "integrated":
        fig = _build_umap_integrated_fig(df, color_by, highlight, show_legend, show_labels,
                                          title=umap_title,
                                          marker_size=marker_size or 2,
                                          exclude_clusters=exclude_clusters,
                                          label_size=label_size or 14,
                                          saved_positions=all_pos.get("umap_integrated"),
                                          custom_colors=custom_colors)
        fs_cfg = dict(fs_config)
        fs_cfg["toImageButtonOptions"] = dict(fs_cfg["toImageButtonOptions"],
                                               filename=f"UMAP_{umap_title}")
        return html.Div(
            style={"width": f"{width_val}vw", "margin": "0 auto"},
            children=[dcc.Graph(id="fs_umap_integrated_graph", figure=fig, style={"height": f"{height_val}vh"}, config=fs_cfg)],
        )
    else:
        per_h = max(height_val // 2, 25)
        name_map = _interactive_data.get("_name_map") or {}
        graphs = _build_umap_per_sample_graphs(df, color_map, highlight,
                                                show_labels, graph_height=f"{per_h}vh",
                                                marker_size=marker_size or 2,
                                                exclude_clusters=exclude_clusters,
                                                label_size=label_size or 11,
                                                saved_positions=all_pos.get("umap_per_sample"),
                                                show_legend=bool(show_legend),
                                                name_map=name_map,
                                                columns_per_row=columns_per_row or 0)
        return html.Div(
            style={"display": "flex", "flexWrap": "wrap", "gap": "15px",
                   "width": f"{width_val}vw", "margin": "0 auto"},
            children=graphs,
        )


# ---------------------------------------------------------------------------
# フルスクリーン Spatial インタラクティブ更新
# ---------------------------------------------------------------------------

@callback(
    Output("fs_spatial_graph_container", "children"),
    [Input("fs_spatial_sample", "value"),
     Input("spatial_rotation_store", "data"),
     Input("fs_spatial_show_labels", "value"),
     Input("fs_spatial_highlight_cluster", "value"),
     Input("fs_spatial_exclude_cluster", "value"),
     Input("fs_spatial_marker_size", "value"),
     Input("fs_spatial_height_slider", "value"),
     Input("fs_spatial_width_slider", "value"),
     Input("fs_spatial_label_size", "value")],
    [State("custom_color_map_store", "data"),
     State("spatial_columns_per_row", "value"),
     State("accumulated_label_positions", "data")],
    prevent_initial_call=True,
)
def update_fs_spatial(sample, rotation_store, show_labels, highlight,
                      exclude_clusters, marker_size, height_val, width_val,
                      label_size, custom_colors, columns_per_row,
                      accumulated_positions):
    height_val = height_val or 60
    width_val = width_val or 95
    df = _interactive_data.get("plot_data")
    if df is None or "SpatialX" not in df.columns:
        return ""
    color_map = _get_cluster_color_map(df["Cluster"], custom_colors)
    cluster_to_idx, discrete_cscale = _get_cluster_colorscale(df["Cluster"], custom_colors)
    all_pos = _get_merged_label_positions(accumulated_positions)
    spatial_pos = all_pos.get("spatial", {})
    if not rotation_store:
        rotation_store = {}

    if sample:
        samples_to_show = [sample]
    else:
        samples_to_show = sorted(df["Sample"].unique())

    fs_config = {"scrollZoom": True, "edits": {"annotationPosition": True}, "toImageButtonOptions": {"format": "png", "scale": 3}}
    name_map = _interactive_data.get("_name_map") or {}
    graphs = []
    for s in samples_to_show:
        display_s = _display_name(s, name_map)
        df_s = df[df["Sample"] == s]
        transform = rotation_store.get(
            s, rotation_store.get("__all__", {"angle": 0, "flip_h": False, "flip_v": False}))
        if isinstance(transform, (int, float)):
            transform = {"angle": int(transform), "flip_h": False, "flip_v": False}
        render_h = round(height_val * 10.8)  # vh → px概算
        fig = _create_single_spatial_fig(df_s, color_map, highlight, set(),
                                         rotation_deg=transform.get("angle", 0),
                                         show_labels=show_labels,
                                         flip_h=transform.get("flip_h", False),
                                         flip_v=transform.get("flip_v", False),
                                         title=display_s, embed_legend=True,
                                         cluster_to_idx=cluster_to_idx,
                                         discrete_cscale=discrete_cscale,
                                         marker_size=marker_size or 0,
                                         exclude_clusters=exclude_clusters,
                                         label_size=label_size or 10,
                                         saved_positions=spatial_pos.get(s),
                                         render_height=render_h)
        if columns_per_row:
            n_cols = columns_per_row
            gap_total = (n_cols - 1) * 15
            flex_basis = f"calc({100 / n_cols:.2f}% - {gap_total / n_cols:.1f}px)"
            min_w = "0"
        else:
            n_cols = len(samples_to_show)
            flex_basis = f"{max(20, 90 // n_cols)}%"
            min_w = "350px"
        fs_cfg = dict(fs_config)
        fs_cfg["toImageButtonOptions"] = dict(fs_cfg["toImageButtonOptions"],
                                               filename=f"Spatial_{display_s}")
        graphs.append(
            html.Div(
                style={"flex": f"1 1 {flex_basis}", "minWidth": min_w,
                        "border": "1px solid #dee2e6", "borderRadius": "6px",
                        "padding": "5px", "backgroundColor": "#fff"},
                children=[
                    dcc.Graph(id={"type": "fs_spatial_graph", "index": s},
                              figure=fig, style={"height": f"{height_val}vh"}, config=fs_cfg),
                ],
            )
        )
    return html.Div(
        style={"display": "flex", "flexWrap": "wrap", "gap": "15px",
               "width": f"{width_val}vw", "margin": "0 auto"},
        children=graphs,
    )



# ---------------------------------------------------------------------------
# ラベル位置の永続保存（v2: relayoutData 蓄積 + DOM スナップショット二重方式）
# ---------------------------------------------------------------------------
# Primary: relayoutData からアノテーション位置をリアルタイム蓄積（サーバーサイド）
# Backup: 保存ボタン押下時に Plotly.js DOM を直接読み取り（クライアントサイド）


# --- メカニズム1: relayoutData 蓄積 ---
# v4: 通常モード / FS モードに分割（動的 Input がコールバック全体をブロックする問題の対策）


def _accumulate_core(triggered_id, existing, excl_fn):
    """蓄積コールバックの共通ロジック。

    Args:
        triggered_id: ctx.triggered_id
        existing: accumulated_label_positions Store の現在値
        excl_fn: triggered_id に応じた除外セットを返す関数
    Returns:
        更新された existing dict、または PreventUpdate
    """
    # triggered[0]["value"] から relayoutData を取得
    rd = None
    for t in ctx.triggered:
        if t.get("value") and isinstance(t["value"], dict):
            rd = t["value"]
            break
    if not rd:
        raise PreventUpdate

    # annotation 位置変更のみ処理（zoom/pan はスキップ）
    if not any(k.startswith("annotations[") for k in rd):
        raise PreventUpdate

    df = _interactive_data.get("plot_data")
    if df is None:
        raise PreventUpdate

    existing = dict(existing) if existing else {}
    excl = excl_fn(triggered_id)

    # 文字列 ID → UMAP 統合
    if isinstance(triggered_id, str):
        clusters = [c for c in sorted(df["Cluster"].unique(), key=_cluster_sort_key)
                    if str(c) not in excl]
        pos = _extract_annotation_positions_by_name(rd, clusters)
        if pos:
            umap_saved = dict(existing.get("umap_integrated", {}))
            _merge_label_positions(umap_saved, pos)
            existing["umap_integrated"] = umap_saved

    # dict ID → パターンマッチ（per_sample / spatial）
    elif isinstance(triggered_id, dict):
        graph_type = triggered_id.get("type")
        sample_name = str(triggered_id.get("index", ""))

        if graph_type in ("umap_per_sample_graph",):
            section = "umap_per_sample"
        elif graph_type in ("spatial_graph", "fs_spatial_graph"):
            section = "spatial"
        else:
            raise PreventUpdate

        sample_df = df[df["Sample"] == sample_name]
        if sample_df.empty:
            raise PreventUpdate
        clusters = [c for c in sorted(sample_df["Cluster"].unique(), key=_cluster_sort_key)
                    if str(c) not in excl]
        pos = _extract_annotation_positions_by_name(rd, clusters)
        if pos:
            section_saved = dict(existing.get(section, {}))
            sample_saved = dict(section_saved.get(sample_name, {}))
            _merge_label_positions(sample_saved, pos)
            section_saved[sample_name] = sample_saved
            existing[section] = section_saved
    else:
        raise PreventUpdate

    return existing


def _excl_set(val):
    """exclude dropdown の値から除外セットを返す"""
    if not val:
        return set()
    return set(str(c) for c in val)


# 1a: 通常モード蓄積（静的 Input のみ → 常に発火可能）
@callback(
    Output("accumulated_label_positions", "data"),
    [Input("interactive_umap_plot", "relayoutData"),
     Input({"type": "umap_per_sample_graph", "index": ALL}, "relayoutData"),
     Input({"type": "spatial_graph", "index": ALL}, "relayoutData")],
    [State("accumulated_label_positions", "data"),
     State("umap_exclude_cluster", "value"),
     State("spatial_exclude_cluster", "value")],
    prevent_initial_call=True,
)
def accumulate_annotation_positions_normal(umap_rd, umap_ps_rds,
                                            spatial_rds, existing,
                                            umap_exclude, spatial_exclude):
    """通常モード: relayoutData のアノテーション位置変更をリアルタイムで蓄積。"""
    triggered_id = ctx.triggered_id
    if not triggered_id:
        raise PreventUpdate

    def _get_excl(tid):
        if isinstance(tid, dict):
            gtype = tid.get("type")
            if gtype == "spatial_graph":
                return _excl_set(spatial_exclude)
            else:
                return _excl_set(umap_exclude)
        return _excl_set(umap_exclude)

    return _accumulate_core(triggered_id, existing, _get_excl)


# 1b: FS UMAP 蓄積（Input 1個: 文字列 ID → UMAP FS 時のみ存在）
@callback(
    Output("accumulated_label_positions", "data", allow_duplicate=True),
    Input("fs_umap_integrated_graph", "relayoutData"),
    [State("accumulated_label_positions", "data"),
     State("fs_umap_exclude_cluster", "value")],
    prevent_initial_call=True,
)
def accumulate_annotation_positions_fs_umap(fs_umap_rd, existing, fs_umap_exclude):
    """FS UMAP: relayoutData のアノテーション位置変更を蓄積。"""
    triggered_id = ctx.triggered_id
    if not triggered_id:
        raise PreventUpdate
    return _accumulate_core(triggered_id, existing, lambda _: _excl_set(fs_umap_exclude))


# 1c: FS Spatial 蓄積（Input 1個: パターンマッチ → Spatial FS 時のみ存在）
@callback(
    Output("accumulated_label_positions", "data", allow_duplicate=True),
    Input({"type": "fs_spatial_graph", "index": ALL}, "relayoutData"),
    [State("accumulated_label_positions", "data"),
     State("fs_spatial_exclude_cluster", "value")],
    prevent_initial_call=True,
)
def accumulate_annotation_positions_fs_spatial(fs_spatial_rds, existing, fs_spatial_exclude):
    """FS Spatial: relayoutData のアノテーション位置変更を蓄積。"""
    triggered_id = ctx.triggered_id
    if not triggered_id:
        raise PreventUpdate
    return _accumulate_core(triggered_id, existing, lambda _: _excl_set(fs_spatial_exclude))


# --- メカニズム2: DOM スナップショット（バックアップ） ---
# v4: 通常 / FS 分割

# 2a: 通常モード DOM スナップショット
clientside_callback(
    ClientsideFunction(namespace="annotation_ns",
                       function_name="capture_annotations_normal"),
    Output("annotation_snapshot_store", "data"),
    [Input("save_label_pos_btn", "n_clicks"),
     Input("save_spatial_label_pos_btn", "n_clicks")],
    prevent_initial_call=True,
)

# 2b: FS DOM スナップショット（パターンマッチング: UMAP/Spatial 共通、ALL で 0個以上にマッチ）
clientside_callback(
    ClientsideFunction(namespace="annotation_ns",
                       function_name="capture_annotations_fs"),
    Output("annotation_snapshot_store", "data", allow_duplicate=True),
    Input({"type": "fs_save_label_pos", "index": ALL}, "n_clicks"),
    prevent_initial_call=True,
)


# --- 保存コールバック: 蓄積データ + スナップショット → JSON ---
# v4: 通常 / FS 分割、共通ロジックはヘルパー関数に抽出


def _do_save_label_positions(accumulated, snapshot):
    """ラベル位置保存の共通ロジック。蓄積データ + DOM スナップショット → JSON。"""
    try:
        path = _get_label_positions_path()
        if not path:
            return no_update, "ラベル位置の保存に失敗しました（データ未読込）", True

        # JSON ファイルの既存データ
        existing = _load_label_positions()

        # --- Primary: 蓄積データからマージ ---
        acc = accumulated or {}
        for section in ("umap_integrated", "umap_per_sample", "spatial"):
            acc_section = acc.get(section)
            if not acc_section:
                continue
            saved_section = existing.get(section, {})
            if section == "umap_integrated":
                _merge_label_positions(saved_section, acc_section)
            else:
                for sample_name, pos_dict in acc_section.items():
                    sample_saved = saved_section.get(sample_name, {})
                    _merge_label_positions(sample_saved, pos_dict)
                    saved_section[sample_name] = sample_saved
            existing[section] = saved_section

        # --- Backup: DOM スナップショットからマージ ---
        if snapshot and snapshot.get("timestamp"):
            def _anns_to_dict(anns_list):
                d = {}
                for a in (anns_list or []):
                    txt = (a.get("text") or "").strip()
                    if txt and a.get("x") is not None and a.get("y") is not None:
                        d[txt] = {"x": a["x"], "y": a["y"]}
                return d

            # UMAP 統合（FS 優先）
            umap_anns = snapshot.get("fs_umap_integrated") or []
            if not umap_anns:
                umap_anns = snapshot.get("umap_integrated") or []
            umap_dict = _anns_to_dict(umap_anns)
            if umap_dict:
                umap_saved = existing.get("umap_integrated", {})
                _merge_label_positions(umap_saved, umap_dict)
                existing["umap_integrated"] = umap_saved

            # サンプル別 UMAP
            for sample_name, anns in (snapshot.get("umap_per_sample") or {}).items():
                sd = _anns_to_dict(anns)
                if sd:
                    ps_saved = existing.get("umap_per_sample", {})
                    ss = ps_saved.get(sample_name, {})
                    _merge_label_positions(ss, sd)
                    ps_saved[sample_name] = ss
                    existing["umap_per_sample"] = ps_saved

            # Spatial（FS 優先）
            for src_key in ("spatial", "fs_spatial"):
                for sample_name, anns in (snapshot.get(src_key) or {}).items():
                    sd = _anns_to_dict(anns)
                    if sd:
                        sp_saved = existing.get("spatial", {})
                        ss = sp_saved.get(sample_name, {})
                        _merge_label_positions(ss, sd)
                        sp_saved[sample_name] = ss
                        existing["spatial"] = sp_saved

        # JSON 書き込み
        path.write_text(json.dumps(existing, indent=2, ensure_ascii=False),
                        encoding="utf-8")

        print(f"[LABEL] ラベル位置を保存しました: {path}")
        return datetime.now().isoformat(), "ラベル位置を保存しました", True
    except Exception as e:
        print(f"[LABEL] ラベル位置保存エラー: {e}")
        import traceback; traceback.print_exc()
        return no_update, f"ラベル位置の保存エラー: {e}", True


# 3a: 通常モード保存
@callback(
    [Output("label_pos_save_status", "data"),
     Output("notification_toast", "children", allow_duplicate=True),
     Output("notification_toast", "is_open", allow_duplicate=True)],
    [Input("save_label_pos_btn", "n_clicks"),
     Input("save_spatial_label_pos_btn", "n_clicks")],
    [State("accumulated_label_positions", "data"),
     State("annotation_snapshot_store", "data")],
    prevent_initial_call=True,
)
def save_label_positions_normal(n1, n2, accumulated, snapshot):
    """通常モード: ラベル位置保存"""
    if not any([n1, n2]):
        raise PreventUpdate
    return _do_save_label_positions(accumulated, snapshot)


# 3b: FS 保存（パターンマッチング: UMAP/Spatial 共通、ALL で 0個以上にマッチ）
@callback(
    [Output("label_pos_save_status", "data", allow_duplicate=True),
     Output("notification_toast", "children", allow_duplicate=True),
     Output("notification_toast", "is_open", allow_duplicate=True)],
    Input({"type": "fs_save_label_pos", "index": ALL}, "n_clicks"),
    [State("accumulated_label_positions", "data"),
     State("annotation_snapshot_store", "data")],
    prevent_initial_call=True,
)
def save_label_positions_fs(n_clicks_list, accumulated, snapshot):
    """FS (UMAP/Spatial 共通): ラベル位置保存"""
    if not any(n_clicks_list):
        raise PreventUpdate
    return _do_save_label_positions(accumulated, snapshot)


# ---------------------------------------------------------------------------
# Volcano Plot（DEG インタラクティブ可視化）
# ---------------------------------------------------------------------------

@callback(
    Output("volcano_cluster_select", "options"),
    Input("deg_data_store", "data"),
    prevent_initial_call=True,
)
def update_volcano_cluster_options(deg_data):
    """DEGデータからVolcano Plotのクラスタ選択肢を生成"""
    if not deg_data:
        return []
    clusters = sorted(
        set(str(r.get("cluster", "")) for r in deg_data),
        key=_cluster_sort_key,
    )
    return [{"label": f"Cluster {c}", "value": c} for c in clusters]


@callback(
    Output("volcano_plot", "figure"),
    [Input("volcano_cluster_select", "value"),
     Input("volcano_fc_threshold", "value"),
     Input("volcano_p_threshold", "value"),
     Input("volcano_y_max", "value"),
     Input("volcano_marker_size", "value")],
    State("deg_data_store", "data"),
    prevent_initial_call=True,
)
def update_volcano_plot(cluster, fc_thresh, p_thresh, y_max, marker_size, deg_data):
    """DEGデータからVolcano Plotを生成"""
    if not deg_data:
        return go.Figure()

    df = pd.DataFrame(deg_data)
    # p_val_adj_raw があればそちらを使用（文字列変換前の精度を保持）
    if "p_val_adj_raw" in df.columns:
        df["p_num"] = pd.to_numeric(df["p_val_adj_raw"], errors="coerce")
    else:
        df["p_num"] = pd.to_numeric(df["p_val_adj"], errors="coerce")
    df["avg_log2FC"] = pd.to_numeric(df["avg_log2FC"], errors="coerce")
    # p=0 は非ゼロ最小p値にclip（真の -log10 を使用、p=0は最上部に集約）
    min_nonzero_p = df.loc[df["p_num"] > 0, "p_num"].min() if (df["p_num"] > 0).any() else 5e-324
    df["neg_log10_p"] = -np.log10(df["p_num"].clip(lower=min_nonzero_p))

    # annotation列があれば、表示テキストに化合物名を含める
    if "annotation" in df.columns:
        df["display_text"] = df.apply(
            lambda r: f"{r['gene']}\n({r['annotation']})"
            if _is_meaningful_annotation(r.get('annotation', ''), r.get('gene', ''))
            else r['gene'],
            axis=1,
        )
    else:
        df["display_text"] = df["gene"]

    if cluster:
        df = df[df["cluster"].astype(str) == str(cluster)]

    fc_thresh = fc_thresh or 0.5
    p_thresh = p_thresh or 1.3
    marker_size = marker_size or 8

    fig = go.Figure()
    for reg, color, label in [
        ("Up", "#FF2D2D", "Up-regulated"),
        ("Down", "#1E5BFF", "Down-regulated"),
        ("NS", "#7A7A7A", "Not significant"),
    ]:
        if reg == "Up":
            mask = (df["neg_log10_p"] >= p_thresh) & (df["avg_log2FC"] >= fc_thresh)
        elif reg == "Down":
            mask = (df["neg_log10_p"] >= p_thresh) & (df["avg_log2FC"] <= -fc_thresh)
        else:
            mask = ~(
                (df["neg_log10_p"] >= p_thresh)
                & (df["avg_log2FC"].abs() >= fc_thresh)
            )
        sub = df[mask]
        if len(sub) > 0:
            fig.add_trace(go.Scattergl(
                x=sub["avg_log2FC"],
                y=sub["neg_log10_p"],
                mode="markers",
                marker=dict(size=marker_size, color=color, opacity=0.7),
                name=label,
                text=sub["display_text"],
                hovertemplate=(
                    "<b>%{text}</b><br>"
                    "log2FC: %{x:.3f}<br>"
                    "-log10(p): %{y:.2f}<extra></extra>"
                ),
            ))

    # 閾値ライン
    fig.add_hline(y=p_thresh, line_dash="dash", line_color="gray", opacity=0.5)
    fig.add_vline(x=fc_thresh, line_dash="dash", line_color="gray", opacity=0.5)
    fig.add_vline(x=-fc_thresh, line_dash="dash", line_color="gray", opacity=0.5)

    title = (f"Volcano Plot - Cluster {cluster}" if cluster
             else "Volcano Plot (全クラスタ)")
    yaxis_opts = {}
    if y_max is not None and y_max > 0:
        yaxis_opts["range"] = [0, y_max]
    fig.update_layout(
        title=dict(text=title, font=dict(size=14), x=0.5),
        xaxis_title="avg_log2FC",
        yaxis_title="-log10(p_val_adj)",
        yaxis=yaxis_opts,
        template="plotly_white",
        margin=dict(l=50, r=20, t=40, b=40),
    )
    return fig


# ---------------------------------------------------------------------------
# Heatmap（DEG Top N マーカーのクラスタ別平均発現量）
# ---------------------------------------------------------------------------

def _build_annotation_csv_map(csv_path, ion_mode="Positive", adduct_patterns=None, tolerance=0.01):
    """アノテーションCSV（TraceFinder/HMDB）→ {m/z: compound_name} マッピング。

    1行目でフォーマットを自動判定し、イオンモード・付加イオンでフィルタして返す。
    Returns:
        dict[float, str] — m/z数値 → 化合物名
    """
    if not csv_path or not Path(csv_path).exists():
        return {}

    try:
        with open(csv_path, encoding="utf-8", errors="replace") as f:
            first_line = f.readline()
    except Exception:
        return {}

    pol_need = "+" if ion_mode == "Positive" else "-"

    # --- TraceFinder format ---
    if "TraceFinder" in first_line:
        try:
            df = pd.read_csv(csv_path, skiprows=2, encoding="utf-8", on_bad_lines="skip")
        except Exception:
            return {}
        if "CompoundName" not in df.columns:
            return {}
        # ExtractedMass / Adduct / Polarity の繰返しブロックを検出
        # pandas は重複カラム名に .1, .2, .3 サフィックスを付与するため対応
        import re as _re2
        mass_cols = [i for i, c in enumerate(df.columns)
                     if c == "ExtractedMass" or _re2.match(r"ExtractedMass\.\d+$", c)]
        add_cols  = [i for i, c in enumerate(df.columns)
                     if c == "Adduct" or _re2.match(r"Adduct\.\d+$", c)]
        pol_cols  = [i for i, c in enumerate(df.columns)
                     if c == "Polarity" or _re2.match(r"Polarity\.\d+$", c)]
        if not mass_cols or not add_cols or not pol_cols:
            return {}
        n_blocks = min(len(mass_cols), len(add_cols), len(pol_cols))
        rows = []
        for k in range(n_blocks):
            for _, row in df.iterrows():
                name = str(row.iloc[0]).strip()  # CompoundName は常に0列目
                try:
                    mz = float(row.iloc[mass_cols[k]])
                except (ValueError, TypeError):
                    continue
                adduct = str(row.iloc[add_cols[k]]).strip()
                pol_raw = str(row.iloc[pol_cols[k]]).strip().upper()
                pol = "+" if pol_raw in ("+", "POS", "P", "POSITIVE") else (
                      "-" if pol_raw in ("-", "NEG", "N", "NEGATIVE") else "")
                if pol == pol_need and name and mz > 0:
                    rows.append((mz, name, adduct))
        # フィルタ: adduct_patterns
        if adduct_patterns:
            filtered = [(mz, nm, ad) for mz, nm, ad in rows
                        if any(p in ad for p in adduct_patterns)]
        else:
            filtered = rows
        return {mz: f"{nm} {ad}" for mz, nm, ad in filtered}

    # --- HMDB format ---
    adduct_col_map = {
        "[M+H]+":  ("M+H",  "+"),
        "[M+Na]+": ("M+Na", "+"),
        "[M+K]+":  ("M+K",  "+"),
        "[M+NH4]+":("M+NH4","+"),
        "[M]+":    ("M+",   "+"),
        "[M-H]-":  ("M-H",  "-"),
        "[M]-":    ("M-",   "-"),
    }
    try:
        df = pd.read_csv(csv_path, encoding="utf-8", on_bad_lines="skip")
    except Exception:
        return {}
    # 化合物名カラム検出
    name_col = None
    for c in ("name (accession)", "name", "Name", "CompoundName"):
        if c in df.columns:
            name_col = c
            break
    if name_col is None:
        return {}

    import re as _re
    result = {}
    for col_name, (adduct_str, pol) in adduct_col_map.items():
        if col_name not in df.columns or pol != pol_need:
            continue
        # adduct_patterns フィルタ
        if adduct_patterns:
            if not any(p in adduct_str for p in adduct_patterns):
                continue
        for _, row in df.iterrows():
            try:
                mz = float(row[col_name])
            except (ValueError, TypeError):
                continue
            if mz <= 0 or pd.isna(mz):
                continue
            name = str(row[name_col]).strip()
            # HMDB ID 除去
            name = _re.sub(r"\s*\(HMDB\d+\)\s*$", "", name)
            if name:
                result[mz] = f"{name} {adduct_str}"
    return result


def _build_mz_to_compound_map(mrm_path_str, tolerance=0.1):
    """MRMファイルから m/z値 → 化合物名 のマッピング辞書を構築する。

    Parent m/z と Daughter m/z の両方を対象にマッチングを行う。
    Returns:
        dict[float, str] — m/z数値 → 化合物名
    """
    if not mrm_path_str:
        return {}
    from app.services.data_manager import load_mrm_file
    mrm_df = load_mrm_file(mrm_path_str)
    if mrm_df is None or mrm_df.empty:
        return {}

    # カラム名の正規化（R側と同様のロジック）
    col_map = {}
    for col in mrm_df.columns:
        cl = col.lower().replace(" ", ".").replace("_", ".")
        if cl in ("compound", "name", "metabolite", "metabolite.name",
                  "analyte", "analyte.name"):
            col_map[col] = "Compound"
        elif cl in ("parent.m.z", "parent.mz", "parent", "precursor",
                    "q1", "q1.m.z", "precursor.m.z", "precursor.mz"):
            col_map[col] = "Parent_mz"
        elif cl in ("daughter.m.z", "daughter.mz", "daughter", "product",
                    "q3", "q3.m.z", "product.m.z", "product.mz"):
            col_map[col] = "Daughter_mz"
    mrm_df = mrm_df.rename(columns=col_map)

    if "Compound" not in mrm_df.columns:
        return {}

    # m/z → 化合物名 マッピング (Parent と Daughter 両方)
    mz_map = {}
    for _, row in mrm_df.iterrows():
        name = str(row.get("Compound", "")).strip()
        if not name:
            continue
        for mz_col in ("Parent_mz", "Daughter_mz"):
            if mz_col in mrm_df.columns:
                try:
                    mz_val = float(row[mz_col])
                    mz_map[mz_val] = name
                except (ValueError, TypeError):
                    continue
    return mz_map


def _annotate_gene_labels(gene_list, mz_to_compound, tolerance=0.1):
    """遺伝子/m/zラベルリストに化合物名を付与して返す。

    Returns:
        list[str] — アノテーション済みラベル（例: "mz_123.456 (Compound A)"）
    """
    if not mz_to_compound:
        return gene_list

    mrm_mz_values = np.array(sorted(mz_to_compound.keys()))
    annotated = []
    for g in gene_list:
        label = g
        # m/z数値を抽出（"mz_123.456" や "123.456" 形式に対応）
        import re
        match = re.search(r"(\d+\.\d+)", str(g))
        if match and len(mrm_mz_values) > 0:
            mz_val = float(match.group(1))
            # 最近傍マッチ
            idx = np.argmin(np.abs(mrm_mz_values - mz_val))
            if abs(mrm_mz_values[idx] - mz_val) <= tolerance:
                compound = mz_to_compound[mrm_mz_values[idx]]
                label = f"{g} ({compound})"
        annotated.append(label)
    return annotated


@callback(
    Output("heatmap_plot", "figure"),
    [Input("heatmap_top_n", "value"),
     Input("heatmap_scale", "value"),
     Input("heatmap_annotation_switch", "value")],
    [State("deg_data_store", "data"),
     State("seurat_cache_dir_store", "data"),
     State("mrm_path", "value")],
    prevent_initial_call=True,
)
def update_heatmap(top_n, scale, annotation_on, deg_data, cache_dir_str, mrm_path_str):
    """DEG Top N マーカーのクラスタ別平均発現量ヒートマップを生成"""
    if not deg_data or not cache_dir_str:
        return go.Figure()

    top_n = top_n or 5
    df_deg = pd.DataFrame(deg_data)
    df_deg["p_num"] = pd.to_numeric(df_deg["p_val_adj"], errors="coerce")

    # 各クラスタの Top N マーカーを抽出
    top_markers = df_deg.sort_values("p_num").groupby("cluster").head(top_n)
    genes = top_markers["gene"].unique().tolist()

    if not genes:
        fig = go.Figure()
        fig.add_annotation(
            text="マーカーが見つかりません", showarrow=False,
            xref="paper", yref="paper", x=0.5, y=0.5,
        )
        return fig

    # expression_matrix.parquet から発現量取得
    cache_dir = Path(cache_dir_str)
    expr_path = cache_dir / "expression_matrix.parquet"
    if not expr_path.exists():
        fig = go.Figure()
        fig.add_annotation(
            text="発現量データがありません", showarrow=False,
            xref="paper", yref="paper", x=0.5, y=0.5,
        )
        return fig

    # 利用可能な遺伝子のみ読み込み
    available = []
    for g in genes:
        try:
            pd.read_parquet(expr_path, columns=[g])
            available.append(g)
        except Exception:
            continue

    if not available:
        fig = go.Figure()
        fig.add_annotation(
            text="発現量カラムが見つかりません", showarrow=False,
            xref="paper", yref="paper", x=0.5, y=0.5,
        )
        return fig

    expr_df = pd.read_parquet(expr_path, columns=["CellID"] + available)

    # クラスタ情報を結合
    plot_data = _interactive_data.get("plot_data")
    if plot_data is None:
        return go.Figure()
    merged = expr_df.merge(
        plot_data[["CellID", "Cluster"]], on="CellID", how="inner"
    )

    # クラスタ別平均発現量
    cluster_means = merged.groupby("Cluster")[available].mean()
    cluster_means = cluster_means.reindex(
        sorted(cluster_means.index, key=_cluster_sort_key)
    )

    # Z-score 変換（scipy なしで手動計算）
    z_data = cluster_means.values.copy()
    if scale == "zscore":
        col_mean = z_data.mean(axis=0)
        col_std = z_data.std(axis=0)
        col_std[col_std == 0] = 1
        z_data = (z_data - col_mean) / col_std

    # Y軸ラベル: アノテーション表示
    y_labels = available
    if annotation_on:
        # 1次ソース: CSV annotation列（deg_dataから取得）
        gene_to_annotation = {}
        if deg_data:
            for r in deg_data:
                gene = r.get("gene", "")
                ann = r.get("annotation", "")
                if gene and _is_meaningful_annotation(ann, gene):
                    gene_to_annotation[gene] = ann
        if gene_to_annotation:
            y_labels = [
                f"{g} ({gene_to_annotation[g]})" if g in gene_to_annotation else g
                for g in available
            ]
        elif mrm_path_str:
            # フォールバック: MRMファイルマッチング
            mz_to_compound = _build_mz_to_compound_map(mrm_path_str, tolerance=0.1)
            y_labels = _annotate_gene_labels(available, mz_to_compound, tolerance=0.1)

    fig = go.Figure(go.Heatmap(
        z=z_data.T,
        x=[f"C{c}" for c in cluster_means.index],
        y=y_labels,
        colorscale="RdBu_r",
        zmid=0 if scale == "zscore" else None,
        hovertemplate=(
            "Cluster: %{x}<br>Gene: %{y}<br>"
            "Value: %{z:.3f}<extra></extra>"
        ),
    ))
    # Y軸ラベルの余白を動的に調整
    max_label_len = max(len(str(l)) for l in y_labels) if y_labels else 10
    left_margin = min(max(max_label_len * 7, 120), 350)
    fig.update_layout(
        title=dict(text=f"Top {top_n} DEG Heatmap", font=dict(size=14), x=0.5),
        xaxis_title="Cluster",
        yaxis_title="Gene / m/z",
        template="plotly_white",
        margin=dict(l=left_margin, r=20, t=40, b=40),
        yaxis=dict(autorange="reversed"),
    )
    return fig


# ===========================================================================
# インタラクティブ m/z キャリブレーション UI コールバック (INT-CB1〜CB11)
# ===========================================================================

# INT-CB1: キャリブレーション詳細パネルの表示切替
@callback(
    Output("int_cal_detail_panel", "style"),
    Input("int_cal_enable", "value"),
    prevent_initial_call=True,
)
def toggle_int_cal_panel(enabled):
    if enabled:
        return {"display": "block", "marginTop": "10px",
                "padding": "10px", "background": "#f8f9fa",
                "borderRadius": "5px"}
    return {"display": "none"}


# INT-CB2: マトリックス/イオンモード変更 → 参照m/zテーブル初期化
@callback(
    [Output("int_cal_table_data", "data", allow_duplicate=True),
     Output("int_cal_restore_pending", "data", allow_duplicate=True)],
    [Input("int_cal_matrix", "value"),
     Input("int_cal_ion_mode", "value")],
    State("int_cal_restore_pending", "data"),
    prevent_initial_call=True,
)
def update_int_cal_table_on_matrix(matrix_type, ion_mode, is_restoring):
    # データ読み込み直後の復元フェーズではテーブル上書きをスキップ
    if is_restoring:
        return no_update, False
    from app.config import MATRIX_REFERENCE_MZ
    if not matrix_type or matrix_type == "custom":
        return no_update, False
    polarity = ion_mode or "Positive"
    ref_list = MATRIX_REFERENCE_MZ.get(matrix_type, {}).get(polarity, [])
    if not ref_list:
        return [], False
    return [{"ref_mz": round(r, 4), "formula": "", "obs_mz": "",
             "ppm_drift": "--", "use": "Yes"} for r in ref_list], False


# INT-CB3: Store → DataTable 同期
@callback(
    [Output("int_cal_table", "data"),
     Output("int_cal_table", "selected_rows")],
    Input("int_cal_table_data", "data"),
    prevent_initial_call=True,
)
def sync_int_cal_store_to_table(store_data):
    data = store_data or []
    selected = [i for i, r in enumerate(data) if r.get("use") == "Yes"]
    return data, selected


# INT-CB4: チェックボックス → use フィールド同期
@callback(
    Output("int_cal_table_data", "data", allow_duplicate=True),
    Input("int_cal_table", "selected_rows"),
    State("int_cal_table", "data"),
    prevent_initial_call=True,
)
def sync_int_cal_selection_to_use(selected_rows, table_data):
    if table_data is None:
        return no_update
    selected_set = set(selected_rows or [])
    changed = False
    new_data = []
    for i, row in enumerate(table_data):
        new_use = "Yes" if i in selected_set else "No"
        if row.get("use") != new_use:
            changed = True
        new_data.append({**row, "use": new_use})
    if not changed:
        return no_update
    return new_data


# INT-CB5: 行追加
@callback(
    Output("int_cal_table_data", "data", allow_duplicate=True),
    Input("int_cal_add_row", "n_clicks"),
    State("int_cal_table", "data"),
    prevent_initial_call=True,
)
def add_int_cal_row(n, data):
    if not n:
        return no_update
    data = list(data or [])
    data.append({"ref_mz": "", "formula": "", "obs_mz": "", "ppm_drift": "--", "use": "Yes"})
    return data


# INT-CB6: 選択行削除
@callback(
    Output("int_cal_table_data", "data", allow_duplicate=True),
    Input("int_cal_delete_rows", "n_clicks"),
    [State("int_cal_table", "selected_rows"),
     State("int_cal_table", "data")],
    prevent_initial_call=True,
)
def delete_int_cal_rows(n, selected, data):
    if not n or not selected or not data:
        return no_update
    return [r for i, r in enumerate(data) if i not in set(selected)]


# INT-CB7: ピーク自動検出
@callback(
    [Output("int_cal_table_data", "data", allow_duplicate=True),
     Output("int_cal_status_text", "children")],
    Input("int_cal_auto_detect", "n_clicks"),
    [State("int_cal_table", "data"),
     State("int_cal_search_window", "value"),
     State("seurat_cache_dir_store", "data"),
     State("data_folder", "value"),
     State("analysis_method", "value"),
     State("analysis_method_tims", "value")],
    prevent_initial_call=True,
)
def auto_detect_int_cal_peaks(n, table_data, search_window, cache_dir,
                              data_folder, analysis_method, analysis_method_tims):
    if not n or not table_data:
        return no_update, no_update

    # --- データ読み込み: 優先順位 1) cache_dir  2) data_folder 生データ ---
    from app.services.data_manager import read_raw_mz_spectrum

    expr_df = None
    if cache_dir:
        expr_path = Path(cache_dir) / "expression_matrix.parquet"
        if expr_path.exists():
            try:
                expr_df = pd.read_parquet(expr_path)
            except Exception:
                expr_df = None

    if expr_df is None:
        is_tims = bool(analysis_method_tims)
        expr_df = read_raw_mz_spectrum(data_folder, is_tims=is_tims)

    if expr_df is None:
        return no_update, "データが見つかりません。データフォルダを確認してください。"

    mz_values = {}
    for col in expr_df.columns:
        match = re.search(r"(\d+\.?\d*)", col)
        if match:
            mz_values[col] = float(match.group(1))
    if not mz_values:
        return no_update, "m/z値を含むフィーチャーが見つかりません。"

    mz_array = np.array(list(mz_values.values()))
    feature_names = list(mz_values.keys())
    avg_spectrum = {f: float(expr_df[f].mean()) for f in feature_names}

    sw = float(search_window or 0.5)
    matched_count = 0
    updated_data = []
    for row in table_data:
        row = dict(row)
        ref = row.get("ref_mz")
        if not ref or ref == "":
            updated_data.append(row)
            continue
        try:
            ref_f = float(ref)
        except (ValueError, TypeError):
            updated_data.append(row)
            continue
        within = np.where(np.abs(mz_array - ref_f) <= sw)[0]
        if len(within) == 0:
            row["obs_mz"] = ""
            row["ppm_drift"] = "--"
            updated_data.append(row)
            continue
        best_idx = max(within, key=lambda i: avg_spectrum.get(feature_names[i], 0))
        obs = float(mz_array[best_idx])
        ppm = (obs - ref_f) / ref_f * 1e6
        row["obs_mz"] = round(obs, 5)
        row["ppm_drift"] = f"{ppm:+.1f}"
        matched_count += 1
        updated_data.append(row)

    status = f"検出完了: {matched_count}/{len(table_data)} ピークがマッチしました"
    return updated_data, status


# INT-CB8: テーブル編集時にΔppm再計算
@callback(
    Output("int_cal_table_data", "data", allow_duplicate=True),
    Input("int_cal_table", "data_timestamp"),
    State("int_cal_table", "data"),
    prevent_initial_call=True,
)
def recalculate_int_cal_ppm(ts, table_data):
    if not table_data:
        return no_update
    changed = False
    updated = []
    for row in table_data:
        row = dict(row)
        ref = row.get("ref_mz")
        obs = row.get("obs_mz")
        if ref and obs and str(ref).strip() and str(obs).strip():
            try:
                ref_f = float(ref)
                obs_f = float(obs)
                ppm = (obs_f - ref_f) / ref_f * 1e6
                new_drift = f"{ppm:+.1f}"
                if row.get("ppm_drift") != new_drift:
                    changed = True
                row["ppm_drift"] = new_drift
            except (ValueError, TypeError):
                pass
        else:
            if row.get("ppm_drift") != "--":
                changed = True
            row["ppm_drift"] = "--"
        updated.append(row)
    if not changed:
        return no_update
    return updated


# INT-CB9: 自動保存
@callback(
    Output("int_cal_save_trigger", "data"),
    [Input("int_cal_enable", "value"),
     Input("int_cal_ion_mode", "value"),
     Input("int_cal_adduct_filter", "value"),
     Input("int_cal_matrix", "value"),
     Input("int_cal_table_data", "data"),
     Input("int_cal_search_window", "value"),
     Input("int_cal_min_peaks", "value"),
     Input("int_cal_regression_mode", "value"),
     Input("int_cal_mrm_path", "value")],
    prevent_initial_call=True,
)
def auto_save_int_cal(enable, ion_mode, adduct_filter, matrix, table_data,
                      search_window, min_peaks, regression_mode,
                      mrm_path):
    _save_interactive_settings("int_calibration", {
        "enable": enable or False,
        "ion_mode": ion_mode or "Positive",
        "adduct_filter": adduct_filter or DEFAULT_ADDUCT_POSITIVE,
        "matrix": matrix or "DHB",
        "table_data": table_data or [],
        "search_window": search_window or 0.5,
        "min_peaks": min_peaks or 2,
        "regression_mode": regression_mode or "poly3",
        "mrm_path": mrm_path or "",
    })
    return None


# INT-CB10: List保存ボタン
@callback(
    Output("int_cal_status_text", "children", allow_duplicate=True),
    Input("int_cal_save_list", "n_clicks"),
    [State("int_cal_enable", "value"),
     State("int_cal_ion_mode", "value"),
     State("int_cal_adduct_filter", "value"),
     State("int_cal_matrix", "value"),
     State("int_cal_table_data", "data"),
     State("int_cal_search_window", "value"),
     State("int_cal_min_peaks", "value"),
     State("int_cal_regression_mode", "value"),
     State("int_cal_mrm_path", "value")],
    prevent_initial_call=True,
)
def save_int_cal_list(n, enable, ion_mode, adduct_filter, matrix, table_data,
                      search_window, min_peaks, regression_mode,
                      mrm_path):
    if not n:
        return no_update
    _save_interactive_settings("int_calibration", {
        "enable": enable or False,
        "ion_mode": ion_mode or "Positive",
        "adduct_filter": adduct_filter or DEFAULT_ADDUCT_POSITIVE,
        "matrix": matrix or "DHB",
        "table_data": table_data or [],
        "search_window": search_window or 0.5,
        "min_peaks": min_peaks or 2,
        "regression_mode": regression_mode or "poly3",
        "mrm_path": mrm_path or "",
    })
    return "設定を保存しました"


# INT-CB-ADDUCT: イオンモード変更時に付加イオン自動切替
@callback(
    Output("int_cal_adduct_filter", "value", allow_duplicate=True),
    Input("int_cal_ion_mode", "value"),
    prevent_initial_call=True,
)
def auto_switch_int_cal_adduct(ion_mode):
    from app.config import DEFAULT_ADDUCT_POSITIVE, DEFAULT_ADDUCT_NEGATIVE
    if ion_mode == "Positive":
        return DEFAULT_ADDUCT_POSITIVE
    return DEFAULT_ADDUCT_NEGATIVE


# INT-CB-MRM: MRMセクション表示制御（DESIのみ表示）
@callback(
    Output("int_cal_mrm_section", "style"),
    Input("int_cal_ms_instrument", "data"),
    prevent_initial_call=True,
)
def toggle_int_cal_mrm_section(ms_instrument):
    if ms_instrument and str(ms_instrument).upper() == "DESI":
        return {}
    return {"display": "none"}


# INT-CB11: キャリブレーション適用
@callback(
    [Output("deg_data_store", "data", allow_duplicate=True),
     Output("int_cal_apply_status", "children")],
    Input("int_cal_apply", "n_clicks"),
    [State("int_cal_enable", "value"),
     State("int_cal_table_data", "data"),
     State("int_cal_search_window", "value"),
     State("int_cal_min_peaks", "value"),
     State("int_cal_regression_mode", "value"),
     State("int_cal_mrm_path", "value"),
     State("tolerance_mz", "value"),
     State("interactive_result_folder", "value"),
     State("interactive_integration_method", "value"),
     State("interactive_rds_map", "data"),
     State("default_annotation_csv", "value"),
     State("int_cal_ion_mode", "value"),
     State("int_cal_adduct_filter", "value")],
    prevent_initial_call=True,
)
def apply_int_calibration(n_clicks, cal_enable, cal_table_data,
                          cal_search_window, cal_min_peaks,
                          cal_regression_mode, mrm_path_str,
                          tolerance_mz, result_folder,
                          integration_method, rds_map,
                          annotation_csv, ion_mode, adduct_filter):
    if not n_clicks:
        raise PreventUpdate

    # DEGデータをディスクから再読み込み
    deg_data = None
    if result_folder:
        result_base = Path(result_folder)
        deg_data = _load_deg_results(result_base, integration_method or "")
    elif _interactive_data.get("rds_path"):
        rds_dir = Path(_interactive_data["rds_path"]).parent
        result_base = rds_dir.parent if rds_dir.name == "RDS_Files" else rds_dir
        deg_data = _load_deg_results(result_base, integration_method or "")

    if not deg_data:
        return no_update, "DEGデータが見つかりません。データを先に読み込んでください。"

    if not cal_enable:
        # キャリブレーション無効 → 元のDEGデータをそのまま返す
        return deg_data, "キャリブレーションは無効です。有効にしてから適用してください。"

    if not cal_table_data:
        return no_update, "キャリブレーションテーブルにデータがありません。"

    # テーブルから use="Yes" のペアを抽出
    matched_pairs = []
    ref_only_mz = []
    for row in cal_table_data:
        if row.get("use") != "Yes":
            continue
        ref = row.get("ref_mz")
        obs = row.get("obs_mz")
        if ref and obs and str(ref).strip() and str(obs).strip():
            try:
                ref_f = float(ref)
                obs_f = float(obs)
                ppm = (obs_f - ref_f) / ref_f * 1e6
                matched_pairs.append({
                    "ref_mz": ref_f, "obs_mz": obs_f, "ppm_drift": ppm,
                })
            except (ValueError, TypeError):
                pass
        elif ref and str(ref).strip():
            try:
                ref_only_mz.append(float(ref))
            except (ValueError, TypeError):
                pass

    features_list = _interactive_data.get("features_list", [])
    if not features_list:
        return no_update, "フィーチャーリストが読み込まれていません。"

    mp = int(cal_min_peaks or 2)
    reg_mode = cal_regression_mode or "poly3"
    cal_result = None

    try:
        if len(matched_pairs) >= mp:
            cal_result = _calibrate_mz_from_pairs(
                features_list, matched_pairs, regression_mode=reg_mode,
            )
        elif ref_only_mz:
            cache_dir = _interactive_data.get("cache_dir")
            expr_path = (Path(cache_dir) / "expression_matrix.parquet"
                         if cache_dir else None)
            if expr_path and expr_path.exists():
                expr_df = pd.read_parquet(expr_path)
                sw = float(cal_search_window or 0.5)
                cal_result = _calibrate_mz(
                    features_list, expr_df, ref_only_mz,
                    search_window=sw, min_peaks=mp,
                    regression_mode=reg_mode,
                )
            else:
                return no_update, "expression_matrix.parquet が見つかりません。ピーク自動検出を先に実行してください。"
        else:
            return no_update, "有効なリファレンス/実測値ペアがありません。テーブルを確認してください。"

        if cal_result and cal_result.get("calibrated"):
            tol = float(tolerance_mz or 0.1)
            deg_data = _reannotate_with_calibration(
                deg_data, cal_result["corrected_mz_map"],
                mrm_path_str, tolerance=tol,
                annotation_csv_path=annotation_csv,
                ion_mode=ion_mode or "Positive",
                adduct_patterns=adduct_filter,
            )
            _interactive_data["_calibration_result"] = cal_result
            r2 = cal_result.get("r_squared", 0)
            mode = cal_result.get("regression_mode", "")
            status = f"キャリブレーション適用完了 (R²={r2:.4f}, {mode})"
            # 設定を永続化
            _save_interactive_settings("int_calibration", {
                "enable": cal_enable,
                "table_data": cal_table_data,
                "search_window": cal_search_window,
                "min_peaks": cal_min_peaks,
                "regression_mode": cal_regression_mode,
                "mrm_path": mrm_path_str or "",
            })
            return deg_data, status
        else:
            return no_update, "キャリブレーションに失敗しました。ピーク数が不足している可能性があります。"

    except Exception as e:
        return no_update, f"キャリブレーションエラー: {e}"


# =========================================================================
# プロジェクトとして保存
# =========================================================================


@callback(
    Output("save_as_project_modal", "is_open"),
    [Input("open_save_as_project_modal", "n_clicks"),
     Input("close_save_as_project_modal", "n_clicks")],
    State("save_as_project_modal", "is_open"),
    prevent_initial_call=True,
)
def toggle_save_as_project_modal(open_clicks, close_clicks, is_open):
    """モーダルの開閉"""
    if ctx.triggered_id in ("open_save_as_project_modal",
                            "close_save_as_project_modal"):
        return not is_open
    return no_update


@callback(
    [Output("sap_new_project_section", "style"),
     Output("sap_existing_project_section", "style"),
     Output("sap_new_sub_section", "style"),
     Output("sap_existing_sub_section", "style")],
    Input("sap_action_type", "value"),
    prevent_initial_call=True,
)
def switch_sap_action_type(action_type):
    """アクション切替でフィールドの表示/非表示を制御"""
    show, hide = {}, {"display": "none"}
    if action_type == "new_all":
        return show, hide, show, hide
    elif action_type == "add_sub":
        return hide, show, show, hide
    elif action_type == "link_existing":
        return hide, show, hide, show
    return no_update, no_update, no_update, no_update


@callback(
    Output("sap_project_select", "options"),
    Input("save_as_project_modal", "is_open"),
    prevent_initial_call=True,
)
def populate_sap_projects(is_open):
    """モーダル表示時にプロジェクト一覧を取得"""
    if not is_open:
        return no_update
    from app.services.project_manager import list_projects
    projects = list_projects()
    return [{"label": p["name"], "value": p["id"]} for p in projects]


@callback(
    Output("sap_sub_select", "options"),
    Input("sap_project_select", "value"),
    prevent_initial_call=True,
)
def populate_sap_sub_projects(project_id):
    """プロジェクト選択時にサブプロジェクト一覧を取得"""
    if not project_id:
        return []
    from app.services.project_manager import list_sub_projects
    subs = list_sub_projects(project_id)
    return [{"label": s["name"], "value": s["id"]} for s in subs]


@callback(
    [Output("sap_result_folder_display", "children"),
     Output("sap_msi_folder_display", "children")],
    Input("save_as_project_modal", "is_open"),
    [State("interactive_result_folder", "value"),
     State("interactive_msi_folder", "value")],
    prevent_initial_call=True,
)
def display_sap_paths(is_open, result_folder, msi_folder):
    """モーダル表示時にパス情報を表示"""
    if not is_open:
        return no_update, no_update
    return (
        f"結果フォルダ: {result_folder or '(未指定)'}",
        f"MSIデータフォルダ: {msi_folder or '(未指定)'}",
    )


@callback(
    [Output("sap_status", "children"),
     Output("save_as_project_modal", "is_open", allow_duplicate=True),
     Output("interactive_project_select", "options", allow_duplicate=True),
     Output("interactive_project_select", "value", allow_duplicate=True),
     Output("interactive_sub_project_select", "value", allow_duplicate=True),
     Output("interactive_entry_mode", "data", allow_duplicate=True),
     Output("sap_skip_reset", "data", allow_duplicate=True)],
    Input("execute_save_as_project", "n_clicks"),
    [State("sap_action_type", "value"),
     State("sap_project_name", "value"),
     State("sap_project_date", "value"),
     State("sap_project_select", "value"),
     State("sap_sub_name", "value"),
     State("sap_sub_date", "value"),
     State("sap_target_compound", "value"),
     State("sap_ms_instrument", "value"),
     State("sap_polarity", "value"),
     State("sap_sub_select", "value"),
     State("interactive_result_folder", "value"),
     State("interactive_msi_folder", "value")],
    prevent_initial_call=True,
)
def execute_save_as_project(
    n_clicks, action_type,
    proj_name, proj_date, existing_proj_id,
    sub_name, sub_date, target_compound, ms_instrument, polarity,
    existing_sub_id,
    result_folder, msi_folder,
):
    """プロジェクトとして保存を実行"""
    if not n_clicks:
        return (no_update,) * 7

    from app.services.project_manager import (
        create_project, create_sub_project, update_sub_project,
        list_projects, get_sub_project,
    )

    _no = (no_update,) * 7

    try:
        if action_type == "new_all":
            # --- 新規プロジェクト + 新規サブプロジェクト ---
            if not proj_name:
                return (
                    dbc.Alert("プロジェクト名を入力してください。",
                              color="warning"),
                ) + (no_update,) * 6
            if not sub_name:
                return (
                    dbc.Alert("サブプロジェクト名を入力してください。",
                              color="warning"),
                ) + (no_update,) * 6
            proj = create_project(
                name=proj_name,
                experiment_date=proj_date or "",
            )
            sub = create_sub_project(
                project_id=proj["id"],
                name=sub_name,
                experiment_date=sub_date or "",
                target_compound=target_compound or "",
                ms_instrument=ms_instrument or "",
                polarity=polarity or [],
                data_folder=msi_folder or "",
                output_dir=result_folder or "",
                extra_fields={"last_result_dir": result_folder or ""},
            )
            pid, sid = proj["id"], sub["id"]

        elif action_type == "add_sub":
            # --- 既存プロジェクトにサブプロジェクト追加 ---
            if not existing_proj_id:
                return (
                    dbc.Alert("プロジェクトを選択してください。",
                              color="warning"),
                ) + (no_update,) * 6
            if not sub_name:
                return (
                    dbc.Alert("サブプロジェクト名を入力してください。",
                              color="warning"),
                ) + (no_update,) * 6
            sub = create_sub_project(
                project_id=existing_proj_id,
                name=sub_name,
                experiment_date=sub_date or "",
                target_compound=target_compound or "",
                ms_instrument=ms_instrument or "",
                polarity=polarity or [],
                data_folder=msi_folder or "",
                output_dir=result_folder or "",
                extra_fields={"last_result_dir": result_folder or ""},
            )
            pid, sid = existing_proj_id, sub["id"]

        elif action_type == "link_existing":
            # --- 既存サブプロジェクトに紐付け ---
            if not existing_proj_id or not existing_sub_id:
                return (
                    dbc.Alert("プロジェクトとサブプロジェクトを選択してください。",
                              color="warning"),
                ) + (no_update,) * 6
            updates = {"last_result_dir": result_folder or ""}
            if msi_folder:
                updates["data_folder"] = msi_folder
            update_sub_project(existing_proj_id, existing_sub_id, updates)
            pid, sid = existing_proj_id, existing_sub_id
        else:
            return _no

        # プロジェクト一覧を更新
        projects = list_projects()
        proj_options = [
            {"label": p["name"], "value": p["id"]} for p in projects
        ]
        return (
            "",             # sap_status (クリア)
            False,          # モーダルを閉じる
            proj_options,   # プロジェクトドロップダウン更新
            pid,            # プロジェクト選択
            sid,            # サブプロジェクト選択
            "sub_project",  # entry_mode (populate が value を保持)
            True,           # sap_skip_reset (リセット抑止)
        )
    except Exception as e:
        return (
            dbc.Alert(f"保存エラー: {e}", color="danger"),
        ) + (no_update,) * 6
